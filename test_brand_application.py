#!/usr/bin/env python3
"""
Test script to verify brand application is working correctly
"""

from pptx.dml.color import RGBColor
from slide_templates import get_brand_styling, render_management_team_slide
from pptx import Presentation

def test_brand_application():
    """Test if brand configuration is properly applied to slides"""
    
    print("🧪 Testing brand application...")
    
    # Create a test brand configuration (manual style)
    test_brand_config = {
        'color_scheme': {
            'primary': RGBColor(255, 0, 0),      # Red primary
            'secondary': RGBColor(0, 255, 0),    # Green secondary  
            'text': RGBColor(0, 0, 0),           # Black text
            'background': RGBColor(255, 255, 255), # White background
            'light_grey': RGBColor(240, 240, 240),
            'footer_grey': RGBColor(128, 128, 128)
        },
        'typography': {
            'primary_font': 'Arial',
            'title_size': 24,
            'content_size': 12
        },
        'source': 'test'
    }
    
    print(f"🧪 Created test brand config with {len(test_brand_config['color_scheme'])} colors")
    
    # Test 1: Test get_brand_styling function
    print("\n🧪 Test 1: Testing get_brand_styling function")
    colors, fonts = get_brand_styling(brand_config=test_brand_config, template_name="modern")
    
    print(f"🧪 Returned colors: {list(colors.keys())}")
    print(f"🧪 Returned fonts: {list(fonts.keys())}")
    
    # Check if primary color is correct
    primary_color = colors.get('primary')
    if isinstance(primary_color, RGBColor):
        hex_str = str(primary_color)
        r = int(hex_str[0:2], 16)
        g = int(hex_str[2:4], 16) 
        b = int(hex_str[4:6], 16)
        print(f"🧪 Primary color: RGB({r},{g},{b}) - Should be RGB(255,0,0)")
        
        if r == 255 and g == 0 and b == 0:
            print("✅ Primary color correctly applied!")
        else:
            print(f"❌ Primary color not applied correctly. Expected RGB(255,0,0), got RGB({r},{g},{b})")
    else:
        print(f"❌ Primary color wrong type: {type(primary_color)}")
    
    # Test 2: Test slide renderer with brand config
    print("\n🧪 Test 2: Testing slide renderer with brand config")
    
    test_data = {
        'title': 'Test Management Team',
        'left_column_profiles': [
            {'name': 'Test CEO', 'role': 'Chief Executive Officer', 'experience': ['10+ years experience']}
        ],
        'right_column_profiles': []
    }
    
    prs = Presentation()
    
    # Render slide with brand config
    result_prs = render_management_team_slide(
        data=test_data,
        prs=prs,
        brand_config=test_brand_config,
        template_name="modern",
        company_name="Test Company"
    )
    
    if result_prs and hasattr(result_prs, 'slides') and len(result_prs.slides) > 0:
        print(f"✅ Slide rendered successfully! Total slides: {len(result_prs.slides)}")
        
        # Save test presentation to check visually
        test_output = "/home/user/webapp/test_brand_output.pptx"
        result_prs.save(test_output)
        print(f"✅ Test presentation saved to: {test_output}")
        
    else:
        print("❌ Slide rendering failed!")
    
    print("\n🧪 Brand application test completed!")

if __name__ == "__main__":
    test_brand_application()
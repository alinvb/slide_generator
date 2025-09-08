#!/usr/bin/env python3

"""
Debug the remaining issues in Historical Financial Performance slide:
1. Chart bars showing identical heights despite different values
2. Key metrics box appearing blank/empty
"""

from pptx import Presentation
from slide_templates import render_historical_financial_performance_slide

def test_saudi_exact_data():
    """Test with the exact Saudi Aramco data format from user's render plan"""
    print("=== Debugging Historical Financial Performance Issues ===")
    
    # Exact data from user's render plan
    saudi_data = {
        "title": "Historical Financial Performance",
        "chart": {
            "title": "Revenue & EBITDA Growth (USD Millions)",
            "categories": ["2020", "2021", "2022", "2023", "2024E"],
            "revenue": [229000, 400000, 495100, 480570, 461560],
            "ebitda": [100000, 180000, 239000, 223000, 215000]
        },
        "key_metrics": {
            "metrics": [
                "Record annual revenue of $495B in 2022",
                "Industry-leading EBITDA margin above 46%",
                "Consistent production above 12 mmboe/d"
            ]
        },
        "revenue_growth": {
            "title": "Revenue Growth Drivers",
            "points": [
                "Growth in gas and chemicals output",
                "Expansion in high-growth international markets",
                "Efficiency gains from digital transformation"
            ]
        },
        "banker_view": {
            "title": "Banker View",
            "text": "Saudi Aramco demonstrates resilient performance through commodity cycles, backed by scale, cost leadership, and diversified downstream assets."
        }
    }
    
    print("1. Data verification:")
    revenue_data = saudi_data["chart"]["revenue"]
    ebitda_data = saudi_data["chart"]["ebitda"]
    print(f"Revenue data: {revenue_data}")
    print(f"EBITDA data: {ebitda_data}")
    print(f"Revenue range: {min(revenue_data):,} to {max(revenue_data):,}")
    print(f"EBITDA range: {min(ebitda_data):,} to {max(ebitda_data):,}")
    
    # Check if chart scaling should work
    all_values = revenue_data + ebitda_data
    data_max = max(all_values)
    expected_axis_max = int(data_max * 1.2)
    print(f"Expected chart axis max: {expected_axis_max:,}")
    
    # Check metrics
    metrics = saudi_data["key_metrics"]["metrics"]
    print(f"Metrics count: {len(metrics)}")
    print(f"Metrics: {metrics}")
    
    print("\n2. Rendering slide...")
    prs = Presentation()
    prs = render_historical_financial_performance_slide(
        data=saudi_data,
        company_name="Saudi Aramco",
        prs=prs
    )
    
    test_file = "debug_historical_saudi_exact.pptx"
    prs.save(test_file)
    print(f"✓ Slide saved as: {test_file}")
    
    return True

def test_chart_data_conversion():
    """Test if there's an issue with chart data conversion"""
    print("\n=== Testing Chart Data Conversion ===")
    
    # Test with simplified data to isolate chart issue
    simple_data = {
        "title": "Historical Financial Performance",
        "chart": {
            "title": "Revenue Test",
            "categories": ["2020", "2021", "2022"],
            "revenue": [100, 200, 300],  # Simple ascending values
            "ebitda": [50, 100, 150]     # Simple ascending values
        },
        "key_metrics": {
            "metrics": [
                "Test metric 1",
                "Test metric 2", 
                "Test metric 3"
            ]
        }
    }
    
    print("Simple test data:")
    print(f"Revenue: {simple_data['chart']['revenue']}")
    print(f"EBITDA: {simple_data['chart']['ebitda']}")
    
    prs = Presentation()
    prs = render_historical_financial_performance_slide(
        data=simple_data,
        company_name="Test Company",
        prs=prs
    )
    
    test_file = "debug_historical_simple.pptx"
    prs.save(test_file)
    print(f"✓ Simple test saved as: {test_file}")
    
    return True

def test_object_format_metrics():
    """Test with object format metrics to see if that fixes the display"""
    print("\n=== Testing Object Format Metrics ===")
    
    # Convert string metrics to object format
    object_data = {
        "title": "Historical Financial Performance",
        "chart": {
            "title": "Revenue & EBITDA Growth (USD Millions)",
            "categories": ["2020", "2021", "2022", "2023", "2024E"],
            "revenue": [229000, 400000, 495100, 480570, 461560],
            "ebitda": [100000, 180000, 239000, 223000, 215000]
        },
        "key_metrics": {
            "metrics": [
                {
                    "title": "Annual Revenue",
                    "value": "$495B",
                    "period": "(2022)",
                    "note": "✓ Record performance"
                },
                {
                    "title": "EBITDA Margin", 
                    "value": "46%+",
                    "period": "(2024E)",
                    "note": "✓ Industry leading"
                },
                {
                    "title": "Production",
                    "value": "12M boe/d",
                    "period": "(Consistent)",
                    "note": "✓ Stable output"
                }
            ]
        },
        "revenue_growth": {
            "title": "Revenue Growth Drivers",
            "points": [
                "Growth in gas and chemicals output",
                "Expansion in high-growth international markets",
                "Efficiency gains from digital transformation"
            ]
        },
        "banker_view": {
            "title": "Banker View",
            "text": "Saudi Aramco demonstrates resilient performance through commodity cycles, backed by scale, cost leadership, and diversified downstream assets."
        }
    }
    
    print("Object format metrics:")
    for i, metric in enumerate(object_data["key_metrics"]["metrics"]):
        print(f"Metric {i+1}: {metric}")
    
    prs = Presentation()
    prs = render_historical_financial_performance_slide(
        data=object_data,
        company_name="Saudi Aramco",
        prs=prs
    )
    
    test_file = "debug_historical_object_metrics.pptx"
    prs.save(test_file)
    print(f"✓ Object metrics test saved as: {test_file}")
    
    return True

def main():
    """Run all debugging tests"""
    print("Debugging Historical Financial Performance slide remaining issues...")
    
    try:
        success1 = test_saudi_exact_data()
        success2 = test_chart_data_conversion()
        success3 = test_object_format_metrics()
        
        if success1 and success2 and success3:
            print(f"\n🔍 DEBUGGING TESTS COMPLETED!")
            print("Generated debug files:")
            print("- debug_historical_saudi_exact.pptx (your exact data)")
            print("- debug_historical_simple.pptx (simple test data)")
            print("- debug_historical_object_metrics.pptx (object format metrics)")
            
            print(f"\n📊 Check these files to identify:")
            print("1. Whether chart bars show proportional heights")
            print("2. Whether metrics box displays content properly")
            print("3. Which format works best for your use case")
            
        else:
            print("\n❌ Some debugging tests failed.")
    except Exception as e:
        print(f"\n💥 Error: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()
import json
import io
from pathlib import Path
import requests
import streamlit as st
import pandas as pd
import zipfile
from datetime import datetime
import re

# Local libs
from executor import execute_plan
from catalog_loader import TemplateCatalog
from brand_extractor import BrandExtractor

# ADD THESE IMPORTS FOR BRAND FUNCTIONALITY
try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    from pptx.enum.dml import MSO_COLOR_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    st.error("python-pptx not installed. Please run: pip install python-pptx")

# validators are optional
try:
    from validators import validate_render_plan_against_catalog, summarize_issues
    HAS_VALIDATORS = True
except Exception:
    HAS_VALIDATORS = False

st.set_page_config(page_title="AI Deck Builder", page_icon="🤖", layout="wide")
st.title("🤖 AI Deck Builder – LLM-Powered Pitch Deck Generator")

# JSON CLEANING FUNCTIONS - ALL SAFE STRING OPERATIONS
def clean_json_string(json_str):
    """Clean and fix common JSON formatting issues using safe string operations"""
    if not json_str:
        return "{}"
    
    # Remove common prefixes and suffixes that break JSON
    json_str = json_str.strip()
    
    # Remove markdown code block indicators using safe string operations
    if json_str.startswith('```json'):
        json_str = json_str[7:].strip()  # Remove ```json
    elif json_str.startswith('```'):
        json_str = json_str[3:].strip()   # Remove ```
    
    if json_str.endswith('```'):
        json_str = json_str[:-3].strip()  # Remove trailing ```
    
    # Remove any leading/trailing non-JSON content
    # Find the first '{' and last '}'
    start_idx = json_str.find('{')
    end_idx = json_str.rfind('}')
    
    if start_idx == -1 or end_idx == -1 or start_idx >= end_idx:
        return "{}"
    
    json_str = json_str[start_idx:end_idx+1]
    
    # Fix common JSON issues using safe replacements
    # Remove trailing commas before }
    lines = json_str.split('\n')
    cleaned_lines = []
    
    for line in lines:
        # Remove trailing comma before } or ]
        line = line.rstrip()
        if line.endswith(',}'):
            line = line[:-2] + '}'
        elif line.endswith(',]'):
            line = line[:-2] + ']'
        cleaned_lines.append(line)
    
    json_str = '\n'.join(cleaned_lines)
    
    return json_str

def extract_jsons_from_response(response_text):
    """Extract both Content IR and Render Plan JSONs from AI response using safe string operations"""
    content_ir = None
    render_plan = None
    
    # Method 1: Look for JSON code blocks using safe string operations
    json_blocks = []
    
    # Find all ```json blocks
    start_pos = 0
    while True:
        start_marker = '```json'
        end_marker = '```'
        
        start_idx = response_text.find(start_marker, start_pos)
        if start_idx == -1:
            break
        
        content_start = start_idx + len(start_marker)
        end_idx = response_text.find(end_marker, content_start)
        if end_idx == -1:
            break
        
        json_content = response_text[content_start:end_idx].strip()
        if json_content:
            json_blocks.append(json_content)
        
        start_pos = end_idx + len(end_marker)
    
    # Method 2: Look for any code blocks if no ```json blocks found
    if not json_blocks:
        start_pos = 0
        while True:
            start_marker = '```'
            end_marker = '```'
            
            start_idx = response_text.find(start_marker, start_pos)
            if start_idx == -1:
                break
            
            # Skip if this is part of ```json (already handled)
            if start_pos == 0 or response_text[start_idx:start_idx+7] != '```json':
                content_start = start_idx + len(start_marker)
                # Skip to next line to avoid immediate end marker
                newline_pos = response_text.find('\n', content_start)
                if newline_pos != -1:
                    content_start = newline_pos + 1
                
                end_idx = response_text.find(end_marker, content_start)
                if end_idx == -1:
                    break
                
                json_content = response_text[content_start:end_idx].strip()
                if json_content and (json_content.startswith('{') or json_content.startswith('[')):
                    json_blocks.append(json_content)
            
            start_pos = start_idx + len(start_marker)
    
    # Try to parse each potential JSON block
    for i, block in enumerate(json_blocks):
        try:
            cleaned_block = clean_json_string(block)
            
            if not cleaned_block or cleaned_block == "{}":
                continue
                
            parsed = json.loads(cleaned_block)
            
            if not isinstance(parsed, dict):
                continue
            
            # Identify which JSON is which based on structure
            if ("entities" in parsed or "management_team" in parsed or 
                "historical_financials" in parsed or "strategic_buyers" in parsed):
                content_ir = parsed
                print(f"✅ Successfully extracted Content IR from block {i+1}")
                
            elif "slides" in parsed and isinstance(parsed.get("slides"), list):
                render_plan = parsed
                print(f"✅ Successfully extracted Render Plan from block {i+1}")
                
        except json.JSONDecodeError as e:
            print(f"❌ Failed to parse JSON block {i+1}: {e}")
            continue
        except Exception as e:
            print(f"❌ Unexpected error parsing block {i+1}: {e}")
            continue
    
    # Method 3: Aggressive line-by-line extraction if needed
    if not content_ir or not render_plan:
        print("🔍 Attempting line-by-line JSON extraction...")
        
        content_ir_markers = ["CONTENT IR JSON:", "Content IR:", "content_ir"]
        render_plan_markers = ["RENDER PLAN JSON:", "Render Plan:", "render_plan"]
        
        lines = response_text.split('\n')
        current_json = []
        json_mode = None
        brace_count = 0
        
        for line in lines:
            line_lower = line.lower()
            
            # Check if we're starting a new JSON section
            if any(marker.lower() in line_lower for marker in content_ir_markers):
                json_mode = "content_ir"
                current_json = []
                brace_count = 0
                continue
            elif any(marker.lower() in line_lower for marker in render_plan_markers):
                json_mode = "render_plan"
                current_json = []
                brace_count = 0
                continue
            
            # If we're in JSON mode, collect lines
            if json_mode and line.strip():
                current_json.append(line)
                brace_count += line.count('{') - line.count('}')
                
                # If braces are balanced and we have content, try to parse
                if brace_count == 0 and len(current_json) > 1:
                    try:
                        json_str = '\n'.join(current_json)
                        cleaned_json = clean_json_string(json_str)
                        parsed = json.loads(cleaned_json)
                        
                        if json_mode == "content_ir" and not content_ir:
                            content_ir = parsed
                            print(f"✅ Extracted Content IR via line-by-line method")
                        elif json_mode == "render_plan" and not render_plan:
                            render_plan = parsed
                            print(f"✅ Extracted Render Plan via line-by-line method")
                            
                    except:
                        pass
                    
                    json_mode = None
                    current_json = []
    
    return content_ir, render_plan

# COMPREHENSIVE SLIDE VALIDATION SYSTEM
def validate_individual_slides(content_ir, render_plan):
    """Validate each slide individually to ensure no empty boxes or missing content"""
    
    validation_results = {
        'overall_valid': True,
        'slide_validations': [],
        'critical_issues': [],
        'warnings': [],
        'summary': {
            'total_slides': 0,
            'valid_slides': 0,
            'invalid_slides': 0,
            'slides_with_warnings': 0
        }
    }
    
    if not render_plan or 'slides' not in render_plan:
        validation_results['critical_issues'].append("No render plan or slides found")
        validation_results['overall_valid'] = False
        return validation_results
    
    slides = render_plan['slides']
    validation_results['summary']['total_slides'] = len(slides)
    
    # Define validation rules for each template
    template_validators = {
        'business_overview': validate_business_overview_slide,
        'investor_considerations': validate_investor_considerations_slide,
        'product_service_footprint': validate_product_service_footprint_slide,
        'product_service_overview': validate_product_service_overview_slide,
        'buyer_profiles': validate_buyer_profiles_slide,
        'historical_financial_performance': validate_historical_financial_performance_slide,
        'management_team': validate_management_team_slide,
        'growth_strategy_projections': validate_growth_strategy_slide,
        'competitive_positioning': validate_competitive_positioning_slide,
        'valuation_overview': validate_valuation_overview_slide,
        'trading_comparables': validate_trading_comparables_slide,
        'precedent_transactions': validate_precedent_transactions_slide,
        'margin_cost_resilience': validate_margin_cost_resilience_slide,
        'financial_summary': validate_financial_summary_slide,
        'transaction_overview': validate_transaction_overview_slide,
        'appendix': validate_appendix_slide,
        'sea_conglomerates': validate_sea_conglomerates_slide
    }
    
    # Validate each slide
    for i, slide in enumerate(slides):
        slide_num = i + 1
        template = slide.get('template', 'unknown')
        
        slide_validation = {
            'slide_number': slide_num,
            'template': template,
            'valid': True,
            'issues': [],
            'warnings': [],
            'missing_fields': [],
            'empty_fields': []
        }
        
        # Basic slide structure validation
        if not slide.get('data'):
            slide_validation['issues'].append("Missing 'data' section")
            slide_validation['valid'] = False
        
        # Template-specific validation
        if template in template_validators:
            template_validator = template_validators[template]
            template_validation = template_validator(slide, content_ir)
            
            slide_validation['issues'].extend(template_validation.get('issues', []))
            slide_validation['warnings'].extend(template_validation.get('warnings', []))
            slide_validation['missing_fields'].extend(template_validation.get('missing_fields', []))
            slide_validation['empty_fields'].extend(template_validation.get('empty_fields', []))
            
            if template_validation.get('issues') or template_validation.get('missing_fields') or template_validation.get('empty_fields'):
                slide_validation['valid'] = False
        else:
            slide_validation['warnings'].append(f"Unknown template type: {template}")
        
        # Update summary counts
        if slide_validation['valid']:
            validation_results['summary']['valid_slides'] += 1
        else:
            validation_results['summary']['invalid_slides'] += 1
            validation_results['overall_valid'] = False
        
        if slide_validation['warnings']:
            validation_results['summary']['slides_with_warnings'] += 1
        
        validation_results['slide_validations'].append(slide_validation)
    
    return validation_results

# FIXED SLIDE-SPECIFIC VALIDATORS
def validate_business_overview_slide(slide, content_ir):
    """Validate business overview slide for completeness"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', {})
    
    # Required fields for business overview
    required_fields = {
        'title': 'Slide title',
        'description': 'Business description',
        'highlights': 'Key highlights',
        'services': 'Services/products list',
        'positioning_desc': 'Market positioning description'
    }
    
    for field, description in required_fields.items():
        if field not in data:
            validation['missing_fields'].append(f"Missing {description} ({field})")
        elif not data[field] or (isinstance(data[field], list) and len(data[field]) == 0):
            validation['empty_fields'].append(f"Empty {description} ({field})")
        elif isinstance(data[field], str) and (data[field].strip() == '' or '[' in data[field]):
            validation['empty_fields'].append(f"Placeholder or empty {description} ({field})")
    
    # Validate highlights array
    if 'highlights' in data and isinstance(data['highlights'], list):
        if len(data['highlights']) < 3:
            validation['warnings'].append("Less than 3 highlights provided - consider adding more")
        for i, highlight in enumerate(data['highlights']):
            if not highlight or highlight.strip() == '' or '[' in highlight:
                validation['empty_fields'].append(f"Empty highlight #{i+1}")
    
    # Validate services array - ENHANCED FOR DETAILED DESCRIPTIONS
    if 'services' in data and isinstance(data['services'], list):
        if len(data['services']) < 6:
            validation['warnings'].append("Less than 6 services provided - business overview should have 6-8 detailed services for comprehensive coverage")
        if len(data['services']) > 8:
            validation['warnings'].append("More than 8 services - consider consolidating for optimal layout")
        
        # Check for category: description format
        non_detailed_services = [s for s in data['services'] if ':' not in s or len(s) < 40]
        if non_detailed_services:
            validation['warnings'].append(f"Found {len(non_detailed_services)} services without detailed descriptions - use format 'Category: Detailed description of scope and capabilities'")
            
        for i, service in enumerate(data['services']):
            if not service or service.strip() == '' or '[' in service:
                validation['empty_fields'].append(f"Empty service #{i+1}")
            elif len(service) < 30:
                validation['warnings'].append(f"Service #{i+1} is very brief - consider adding more descriptive detail")
    
    return validation

def validate_product_service_footprint_slide(slide, content_ir):
    """Validate product service footprint slide - the one with empty boxes"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', {})
    
    # Required fields
    if 'title' not in data or not data['title']:
        validation['missing_fields'].append("Missing slide title")
    
    if 'services' not in data:
        validation['missing_fields'].append("Missing services array")
    elif not isinstance(data['services'], list) or len(data['services']) == 0:
        validation['empty_fields'].append("Empty services array")
    else:
        # Validate each service entry
        for i, service in enumerate(data['services']):
            service_num = i + 1
            if not isinstance(service, dict):
                validation['issues'].append(f"Service #{service_num} is not a proper object")
                continue
                
            if 'title' not in service or not service['title'] or service['title'].strip() == '':
                validation['empty_fields'].append(f"Service #{service_num} missing title")
            elif '[' in service['title']:
                validation['empty_fields'].append(f"Service #{service_num} has placeholder title")
                
            if 'desc' not in service or not service['desc'] or service['desc'].strip() == '':
                validation['empty_fields'].append(f"Service #{service_num} missing description")
            elif '[' in service['desc']:
                validation['empty_fields'].append(f"Service #{service_num} has placeholder description")
    
    # ENHANCED: Check for market coverage data (right side of slide) - REQUIRES 3-4 COLUMNS
    if 'coverage_table' in data:
        coverage_data = data['coverage_table']
        if not coverage_data or (isinstance(coverage_data, list) and len(coverage_data) == 0):
            validation['empty_fields'].append("Empty coverage table section")
        elif isinstance(coverage_data, list) and len(coverage_data) > 0:
            # Check column structure
            if isinstance(coverage_data[0], list):
                num_cols = len(coverage_data[0])
                if num_cols < 3:
                    validation['issues'].append(f"Coverage table has only {num_cols} columns - market comparison tables should have 3-4 columns")
                elif num_cols > 4:
                    validation['warnings'].append(f"Coverage table has {num_cols} columns - optimal range is 3-4 columns for readability")
                
                # Validate header row content
                if len(coverage_data) > 0 and isinstance(coverage_data[0], list):
                    header_row = coverage_data[0]
                    required_concepts = ['region', 'market', 'segment', 'business', 'asset', 'coverage', 'product', 'service']
                    header_text = ' '.join(str(h).lower() for h in header_row)
                    has_market_concepts = any(concept in header_text for concept in required_concepts)
                    if not has_market_concepts:
                        validation['warnings'].append("Table headers don't seem to contain market/business concepts - should include terms like Region, Market Segment, Assets, Coverage, etc.")
            else:
                validation['issues'].append("Coverage table format invalid - should be array of arrays (rows and columns)")
    else:
        validation['issues'].append("Missing coverage_table - market coverage comparison is required for this slide")
    
    if 'metrics' in data:
        metrics = data['metrics']
        if not metrics or (isinstance(metrics, dict) and len(metrics) == 0):
            validation['empty_fields'].append("Empty metrics section")
    else:
        validation['warnings'].append("No operational metrics - may result in empty boxes")
    
    return validation

def validate_buyer_profiles_slide(slide, content_ir):
    """Validate buyer profiles slide - FIXED to handle both approaches correctly"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', {})
    
    # Check for content_ir_key (preferred) or table_rows (fallback)
    has_content_ir_key = 'content_ir_key' in slide
    has_table_rows = 'table_rows' in data
    
    if not has_content_ir_key and not has_table_rows:
        validation['issues'].append("Missing content_ir_key - table will be empty")
    elif has_content_ir_key:
        content_key = slide['content_ir_key']
        
        # Verify the key exists in content_ir and has data
        if content_key not in content_ir:
            validation['issues'].append(f"content_ir_key '{content_key}' not found in Content IR")
        elif not content_ir[content_key] or len(content_ir[content_key]) == 0:
            validation['empty_fields'].append(f"Empty {content_key} array in Content IR")
        else:
            # Validate buyer data completeness
            buyers = content_ir[content_key]
            if not isinstance(buyers, list):
                validation['issues'].append(f"content_ir_key '{content_key}' should be an array")
            else:
                for i, buyer in enumerate(buyers):
                    buyer_num = i + 1
                    if not isinstance(buyer, dict):
                        validation['issues'].append(f"Buyer #{buyer_num} should be an object")
                        continue
                    
                    # Check for required buyer fields - FIXED for your data structure
                    required_buyer_fields = ['buyer_name', 'strategic_rationale', 'fit']
                    for field in required_buyer_fields:
                        if field not in buyer:
                            validation['empty_fields'].append(f"Buyer #{buyer_num} missing {field}")
                        elif not buyer[field] or str(buyer[field]).strip() == '':
                            validation['empty_fields'].append(f"Buyer #{buyer_num} has empty {field}")
                        elif '[' in str(buyer[field]):
                            validation['empty_fields'].append(f"Buyer #{buyer_num} has placeholder {field}")
    
    elif has_table_rows and not has_content_ir_key:
        # Validate table_rows content - FIXED to handle your data structure
        validation['warnings'].append("Using hardcoded table_rows - content_ir_key preferred for dynamic data")
        
        table_rows = data.get('table_rows', [])
        if not table_rows or len(table_rows) == 0:
            validation['empty_fields'].append("Empty table_rows array")
        else:
            for i, row in enumerate(table_rows):
                row_num = i + 1
                # Your table_rows contain dictionaries, not lists
                if isinstance(row, dict):
                    # Check if it has required fields
                    required_fields = ['buyer_name', 'strategic_rationale']
                    for field in required_fields:
                        if field not in row or not row[field] or str(row[field]).strip() == '':
                            validation['empty_fields'].append(f"Table row #{row_num} missing or empty {field}")
                elif isinstance(row, list):
                    if len(row) == 0:
                        validation['empty_fields'].append(f"Table row #{row_num} is empty")
                    else:
                        for j, cell in enumerate(row):
                            cell_num = j + 1
                            if not cell or str(cell).strip() == '' or '[' in str(cell):
                                validation['empty_fields'].append(f"Table row #{row_num}, cell #{cell_num} is empty or placeholder")
                else:
                    validation['empty_fields'].append(f"Table row #{row_num} has invalid structure")
    
    # Validate required fields
    required_fields = ['title', 'table_headers']
    for field in required_fields:
        if field not in data:
            validation['missing_fields'].append(f"Missing {field}")
        elif not data[field]:
            validation['empty_fields'].append(f"Empty {field}")
        elif field == 'table_headers' and isinstance(data[field], list):
            if len(data[field]) < 2:
                validation['warnings'].append("Less than 2 table headers - table may look incomplete")
            for i, header in enumerate(data[field]):
                if not header or str(header).strip() == '':
                    validation['empty_fields'].append(f"Table header #{i+1} is empty")
    
    return validation

def validate_management_team_slide(slide, content_ir):
    """Validate management team slide - ENHANCED with name and layout requirements"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    # Check if using content_ir_key approach
    if 'content_ir_key' in slide:
        content_key = slide['content_ir_key']
        if content_key not in content_ir:
            validation['issues'].append(f"content_ir_key '{content_key}' not found in Content IR")
            return validation
        mgmt_data = content_ir[content_key]
    else:
        # Check data section
        data = slide.get('data', {})
        if 'management_team' not in content_ir:
            validation['issues'].append("No management_team data in Content IR")
            return validation
        mgmt_data = content_ir['management_team']
    
    # Check for required profile arrays with ENHANCED requirements
    for column in ['left_column_profiles', 'right_column_profiles']:
        if column not in mgmt_data:
            validation['missing_fields'].append(f"Missing {column}")
        elif not isinstance(mgmt_data[column], list) or len(mgmt_data[column]) == 0:
            validation['empty_fields'].append(f"Empty {column}")
        elif len(mgmt_data[column]) < 2:
            validation['warnings'].append(f"{column} should have at least 2 profiles for comprehensive management overview")
        else:
            # ENHANCED: Validate individual profiles with NAME and LENGTH requirements
            if len(mgmt_data[column]) > 3:
                validation['warnings'].append(f"{column} has {len(mgmt_data[column])} profiles - recommend max 3 per column for optimal layout")
                
            for i, profile in enumerate(mgmt_data[column]):
                profile_num = i + 1
                
                # Check for required fields
                if 'role_title' not in profile or not profile['role_title']:
                    validation['empty_fields'].append(f"{column} profile #{profile_num} missing role_title")
                elif '[' in str(profile['role_title']):
                    validation['empty_fields'].append(f"{column} profile #{profile_num} has placeholder role_title")
                else:
                    role_title = profile['role_title']
                    # ENHANCED: Check for name requirement (look for common name patterns)
                    if not any(indicator in role_title.lower() for indicator in [' - ', 'ceo', 'cfo', 'coo', 'president']) or len(role_title) < 10:
                        validation['warnings'].append(f"{column} profile #{profile_num} role_title should include executive name (e.g., 'CEO - John Smith')")
                    
                    # ENHANCED: Check length limit
                    if len(role_title) > 50:
                        validation['warnings'].append(f"{column} profile #{profile_num} role_title too long ({len(role_title)} chars) - max 50 chars to prevent overlap")
                
                # Check experience bullets with enhanced validation
                if 'experience_bullets' not in profile or not profile['experience_bullets']:
                    validation['empty_fields'].append(f"{column} profile #{profile_num} missing experience_bullets")
                elif not isinstance(profile['experience_bullets'], list) or len(profile['experience_bullets']) == 0:
                    validation['empty_fields'].append(f"{column} profile #{profile_num} experience_bullets must be non-empty array")
                else:
                    bullets = profile['experience_bullets']
                    if len(bullets) > 4:
                        validation['warnings'].append(f"{column} profile #{profile_num} has {len(bullets)} bullets - recommend max 4 for optimal layout")
                    
                    # Check individual bullet length
                    for j, bullet in enumerate(bullets):
                        if len(bullet) > 80:
                            validation['warnings'].append(f"{column} profile #{profile_num} bullet #{j+1} too long ({len(bullet)} chars) - max 80 chars to prevent overlap")
                        elif len(bullet) < 20:
                            validation['warnings'].append(f"{column} profile #{profile_num} bullet #{j+1} very brief - consider adding more detail")
    
    return validation

def validate_historical_financial_performance_slide(slide, content_ir):
    """Validate historical financial performance slide"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', {})
    
    # Required fields for historical financial performance
    required_fields = {
        'title': 'Slide title',
        'chart': 'Financial performance chart data',
        'key_metrics': 'Key financial metrics'
    }
    
    for field, description in required_fields.items():
        if field not in data:
            validation['missing_fields'].append(f"Missing {description}")
        elif not data[field]:
            validation['empty_fields'].append(f"Empty {description}")
    
    # Validate chart data
    if 'chart' in data and isinstance(data['chart'], dict):
        chart = data['chart']
        chart_required = ['categories', 'revenue', 'ebitda']
        for field in chart_required:
            if field not in chart or not chart[field]:
                validation['empty_fields'].append(f"Missing chart {field} data")
    
    # Validate key metrics
    if 'key_metrics' in data and isinstance(data['key_metrics'], dict):
        metrics = data['key_metrics']
        if 'metrics' in metrics and isinstance(metrics['metrics'], list):
            if len(metrics['metrics']) < 3:
                validation['warnings'].append("Less than 3 key metrics - consider adding more")
        else:
            validation['empty_fields'].append("Missing metrics array in key_metrics")
    
    return validation

def validate_growth_strategy_slide(slide, content_ir):
    """Validate growth strategy slide"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', {})
    
    # Get actual data structure - check for slide_data wrapper
    if 'slide_data' in data:
        actual_data = data['slide_data']
    else:
        actual_data = data
    
    # Required fields for growth strategy
    required_fields = {
        'title': 'Slide title',
        'growth_strategy': 'Growth strategy section',
        'financial_projections': 'Financial projections'
    }
    
    for field, description in required_fields.items():
        if field not in actual_data:
            validation['missing_fields'].append(f"Missing {description}")
        elif not actual_data[field]:
            validation['empty_fields'].append(f"Empty {description}")
    
    # Validate growth strategy
    if 'growth_strategy' in actual_data and isinstance(actual_data['growth_strategy'], dict):
        growth_strat = actual_data['growth_strategy']
        if 'strategies' in growth_strat and isinstance(growth_strat['strategies'], list):
            if len(growth_strat['strategies']) < 3:
                validation['warnings'].append("Less than 3 growth strategies - consider adding more")
        else:
            validation['empty_fields'].append("Missing strategies array in growth_strategy")
    
    return validation

def validate_competitive_positioning_slide(slide, content_ir):
    """ENHANCED: Validate competitive positioning slide - iCar Asia format requirements"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', {})
    
    # ENHANCED: Check for iCar Asia format requirements
    required_fields = {
        'title': 'Slide title',
        'competitors': 'Competitors list for revenue chart',
        'assessment': 'Competitive assessment table (REQUIRED for iCar Asia format)',
        'advantages': 'Competitive advantages',
        'barriers': 'Market barriers to entry'
    }
    
    for field, description in required_fields.items():
        if field not in data:
            if field == 'assessment':
                validation['issues'].append(f"Missing {description} - critical for professional competitive analysis")
            else:
                validation['missing_fields'].append(f"Missing {description}")
        elif not data[field]:
            validation['empty_fields'].append(f"Empty {description}")
    
    # ENHANCED: Validate 5-column assessment table structure (iCar Asia format)
    if 'assessment' in data:
        assessment = data['assessment']
        if not assessment or not isinstance(assessment, list) or len(assessment) == 0:
            validation['issues'].append("Empty competitive assessment table - iCar Asia format requires 5-column structure")
        elif len(assessment) > 0:
            # Check column structure
            if isinstance(assessment[0], list):
                num_cols = len(assessment[0])
                if num_cols < 5:
                    validation['issues'].append(f"Assessment table has only {num_cols} columns - iCar Asia format requires 5 columns: Company, Market Share, Tech Platform, Coverage, Revenue")
                elif num_cols > 5:
                    validation['warnings'].append(f"Assessment table has {num_cols} columns - optimal is 5 for iCar Asia format")
                
                # Validate header structure
                if len(assessment) > 0:
                    header_row = assessment[0]
                    expected_concepts = ['company', 'market', 'tech', 'platform', 'coverage', 'revenue']
                    header_text = ' '.join(str(h).lower() for h in header_row)
                    has_expected_headers = any(concept in header_text for concept in expected_concepts)
                    if not has_expected_headers:
                        validation['warnings'].append("Table headers don't match iCar Asia format - should include Company, Market Share, Tech Platform, Coverage, Revenue")
                
                # Check for star ratings in data rows
                if len(assessment) > 1:
                    data_rows = assessment[1:]
                    has_star_ratings = False
                    for row in data_rows[:3]:  # Check first few data rows
                        if isinstance(row, list) and len(row) > 1:
                            for cell in row[1:-1]:  # Skip company name and revenue columns
                                if '⭐' in str(cell) or '★' in str(cell):
                                    has_star_ratings = True
                                    break
                    if not has_star_ratings:
                        validation['warnings'].append("Assessment table should use star ratings (⭐⭐⭐⭐) for visual comparison like iCar Asia format")
            else:
                validation['issues'].append("Assessment table format invalid - should be array of arrays (rows and columns)")
    
    # ENHANCED: Validate competitors for revenue chart
    if 'competitors' in data and isinstance(data['competitors'], list):
        if len(data['competitors']) < 3:
            validation['warnings'].append("Less than 3 competitors for revenue chart - consider adding more for comprehensive analysis")
        for i, competitor in enumerate(data['competitors']):
            comp_num = i + 1
            if isinstance(competitor, dict):
                if 'name' not in competitor or not competitor['name']:
                    validation['empty_fields'].append(f"Competitor #{comp_num} missing name")
                if 'revenue' not in competitor or not competitor['revenue']:
                    validation['empty_fields'].append(f"Competitor #{comp_num} missing revenue data")
            elif not competitor or '[' in str(competitor):
                validation['empty_fields'].append(f"Competitor #{comp_num} is empty or placeholder")
    
    # ENHANCED: Check advantages and barriers content quality
    for section in ['barriers', 'advantages']:
        if section in data:
            items = data[section]
            if isinstance(items, list):
                if len(items) > 6:
                    validation['warnings'].append(f"Too many {section} ({len(items)}) - consider limiting to 4-6 for clean layout")
                for i, item in enumerate(items):
                    if isinstance(item, dict):
                        if not item.get('title') or not item.get('desc'):
                            validation['empty_fields'].append(f"{section.title()} #{i+1} missing title or description")
                    elif isinstance(item, str):
                        if len(item) > 80:
                            validation['warnings'].append(f"{section.title()} #{i+1} text too long ({len(item)} chars) - consider shortening for clean layout")
    
    return validation

def validate_valuation_overview_slide(slide, content_ir):
    """Validate valuation overview slide - FIXED for correct field names"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', {})
    
    # FIXED: Use the correct field names from your data structure
    required_fields = {
        'title': 'Slide title',
        'valuation_data': 'Valuation methodologies data'  # FIXED
    }
    
    for field, description in required_fields.items():
        if field not in data:
            validation['missing_fields'].append(f"Missing {description}")
        elif not data[field]:
            validation['empty_fields'].append(f"Empty {description}")
    
    # Validate valuation_data array
    if 'valuation_data' in data and isinstance(data['valuation_data'], list):
        if len(data['valuation_data']) < 2:
            validation['warnings'].append("Less than 2 valuation methodologies - consider adding more")
        for i, method in enumerate(data['valuation_data']):
            method_num = i + 1
            if isinstance(method, dict):
                required_method_fields = ['methodology', 'enterprise_value', 'commentary']
                for field in required_method_fields:
                    if field not in method or not method[field]:
                        validation['empty_fields'].append(f"Methodology #{method_num} missing {field}")
            elif not method or '[' in str(method):
                validation['empty_fields'].append(f"Methodology #{method_num} is empty or placeholder")
    
    return validation

def validate_trading_comparables_slide(slide, content_ir):
    """Validate trading comparables slide"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', {})
    
    # Required fields
    required_fields = {
        'title': 'Slide title',
        'comparable_companies': 'Comparable companies list',
        'metrics': 'Financial metrics comparison'
    }
    
    for field, description in required_fields.items():
        if field not in data:
            validation['missing_fields'].append(f"Missing {description}")
        elif not data[field]:
            validation['empty_fields'].append(f"Empty {description}")
    
    # Validate comparable companies
    if 'comparable_companies' in data and isinstance(data['comparable_companies'], list):
        if len(data['comparable_companies']) < 4:
            validation['warnings'].append("Less than 4 comparable companies - consider adding more")
        for i, company in enumerate(data['comparable_companies']):
            comp_num = i + 1
            if isinstance(company, dict):
                required_comp_fields = ['name', 'market_cap', 'revenue', 'ebitda_multiple']
                for field in required_comp_fields:
                    if field not in company or not company[field]:
                        validation['empty_fields'].append(f"Comparable #{comp_num} missing {field}")
                    elif '[' in str(company[field]):
                        validation['empty_fields'].append(f"Comparable #{comp_num} has placeholder {field}")
    
    return validation

def validate_precedent_transactions_slide(slide, content_ir):
    """Validate precedent transactions slide"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', {})
    
    # Required fields
    required_fields = {
        'title': 'Slide title',
        'transactions': 'Precedent transactions list'
    }
    
    for field, description in required_fields.items():
        if field not in data:
            validation['missing_fields'].append(f"Missing {description}")
        elif not data[field]:
            validation['empty_fields'].append(f"Empty {description}")
    
    # Validate transactions
    if 'transactions' in data and isinstance(data['transactions'], list):
        if len(data['transactions']) < 3:
            validation['warnings'].append("Less than 3 precedent transactions - consider adding more")
        for i, transaction in enumerate(data['transactions']):
            trans_num = i + 1
            if isinstance(transaction, dict):
                required_trans_fields = ['target', 'acquirer', 'date', 'enterprise_value', 'revenue', 'ev_revenue_multiple']
                for field in required_trans_fields:
                    if field not in transaction or not transaction[field]:
                        validation['empty_fields'].append(f"Transaction #{trans_num} missing {field}")
                    elif '[' in str(transaction[field]):
                        validation['empty_fields'].append(f"Transaction #{trans_num} has placeholder {field}")
    
    return validation

def validate_margin_cost_resilience_slide(slide, content_ir):
    """Validate margin/cost resilience slide - FIXED for correct field names"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', {})
    
    # FIXED: Use the correct field names from your data structure
    required_fields = {
        'title': 'Slide title',
        'cost_management': 'Cost management initiatives',  # FIXED
        'risk_mitigation': 'Risk mitigation strategies'     # FIXED
    }
    
    for field, description in required_fields.items():
        if field not in data:
            validation['missing_fields'].append(f"Missing {description}")
        elif not data[field]:
            validation['empty_fields'].append(f"Empty {description}")
    
    # Validate cost management items
    if 'cost_management' in data:
        cost_mgmt = data['cost_management']
        if isinstance(cost_mgmt, dict) and 'items' in cost_mgmt:
            items = cost_mgmt['items']
            if not items or len(items) == 0:
                validation['empty_fields'].append("Empty cost management items")
            else:
                for i, item in enumerate(items):
                    if not isinstance(item, dict):
                        validation['empty_fields'].append(f"Cost management item #{i+1} is not properly structured")
                    elif not item.get('title') or not item.get('description'):
                        validation['empty_fields'].append(f"Cost management item #{i+1} missing title or description")
    
    # Validate risk mitigation
    if 'risk_mitigation' in data:
        risk_mit = data['risk_mitigation']
        if isinstance(risk_mit, dict):
            if 'main_strategy' not in risk_mit:
                validation['missing_fields'].append("Missing main strategy in risk mitigation")
        
    return validation

def validate_investor_considerations_slide(slide, content_ir):
    """Validate investor considerations slide"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', {})
    
    required_fields = {
        'title': 'Slide title',
        'considerations': 'Investment considerations list',
        'mitigants': 'Risk mitigants list'
    }
    
    for field, description in required_fields.items():
        if field not in data:
            validation['missing_fields'].append(f"Missing {description}")
        elif not data[field] or (isinstance(data[field], list) and len(data[field]) == 0):
            validation['empty_fields'].append(f"Empty {description}")
    
    return validation

def validate_financial_summary_slide(slide, content_ir):
    """Validate financial summary slide"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', {})
    
    required_fields = {
        'title': 'Slide title',
        'key_metrics': 'Key financial metrics',
        'performance_highlights': 'Performance highlights'
    }
    
    for field, description in required_fields.items():
        if field not in data:
            validation['missing_fields'].append(f"Missing {description}")
        elif not data[field]:
            validation['empty_fields'].append(f"Empty {description}")
    
    return validation

def validate_transaction_overview_slide(slide, content_ir):
    """Validate transaction overview slide"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', {})
    
    required_fields = {
        'title': 'Slide title',
        'transaction_structure': 'Transaction structure',
        'key_terms': 'Key transaction terms',
        'timeline': 'Transaction timeline'
    }
    
    for field, description in required_fields.items():
        if field not in data:
            validation['missing_fields'].append(f"Missing {description}")
        elif not data[field]:
            validation['empty_fields'].append(f"Empty {description}")
    
    return validation

def validate_product_service_overview_slide(slide, content_ir):
    """Validate product/service overview slide"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', {})
    
    # Required fields
    required_fields = {
        'title': 'Slide title',
        'products': 'Products list',
        'market_position': 'Market positioning'
    }
    
    for field, description in required_fields.items():
        if field not in data:
            validation['missing_fields'].append(f"Missing {description}")
        elif not data[field]:
            validation['empty_fields'].append(f"Empty {description}")
    
    return validation

def validate_appendix_slide(slide, content_ir):
    """Validate appendix slide"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', {})
    
    if 'title' not in data or not data['title']:
        validation['missing_fields'].append("Missing appendix title")
    
    return validation

def validate_sea_conglomerates_slide(slide, content_ir):
    """Validate SEA conglomerates slide"""
    validation = {'issues': [], 'warnings': [], 'missing_fields': [], 'empty_fields': []}
    
    data = slide.get('data', [])
    
    if not data or len(data) == 0:
        validation['empty_fields'].append("Empty conglomerates data")
    else:
        for i, conglomerate in enumerate(data):
            cong_num = i + 1
            if isinstance(conglomerate, dict):
                required_fields = ['name', 'country', 'description']
                for field in required_fields:
                    if field not in conglomerate or not conglomerate[field]:
                        validation['empty_fields'].append(f"Conglomerate #{cong_num} missing {field}")
    
    return validation

# VALIDATION DISPLAY FUNCTIONS
def display_validation_results(validation_results):
    """Display comprehensive validation results with visual indicators"""
    
    summary = validation_results['summary']
    
    # Create header with summary box
    st.markdown("### 📋 Slide-by-Slide Validation Results")
    
    # Enhanced summary with quality scores
    col1, col2, col3, col4, col5 = st.columns(5)
    with col1:
        st.metric("Total Slides", summary['total_slides'])
    with col2:
        st.metric("Valid Slides", summary['valid_slides'], delta=None if summary['valid_slides'] == summary['total_slides'] else f"-{summary['invalid_slides']}")
    with col3:
        st.metric("Invalid Slides", summary['invalid_slides'], delta=None if summary['invalid_slides'] == 0 else f"+{summary['invalid_slides']}")
    with col4:
        if 'data_quality_score' in validation_results:
            st.metric("Data Quality", f"{validation_results['data_quality_score']:.0f}%")
    with col5:
        if 'completeness_score' in validation_results:
            st.metric("Completeness", f"{validation_results['completeness_score']:.0f}%")
    
    # Overall status with enhanced feedback
    if validation_results['overall_valid']:
        st.success("✅ All slides passed validation! Ready for deck generation.")
        if 'data_quality_score' in validation_results:
            quality_score = validation_results['data_quality_score']
            if quality_score >= 90:
                st.success("🏆 Excellent data quality - matches professional standards!")
            elif quality_score >= 80:
                st.info("👍 Good data quality - ready for production use")
            elif quality_score >= 70:
                st.warning("⚠️ Acceptable quality - minor improvements recommended")
    else:
        st.error(f"❌ {summary['invalid_slides']} slide(s) have validation issues that must be fixed before generating deck.")
    
    # Structure validation results
    if 'structure_validation' in validation_results:
        struct_val = validation_results['structure_validation']
        if struct_val['structure_issues']:
            st.markdown("#### 🗃️ Structure Issues Found")
            st.error("The following structural issues were detected by comparing against professional examples:")
            for issue in struct_val['structure_issues']:
                st.markdown(f"• {issue}")
    
    # Detailed slide results
    if validation_results['slide_validations']:
        st.markdown("#### Detailed Slide Analysis")
        
        for slide_val in validation_results['slide_validations']:
            slide_num = slide_val['slide_number']
            template = slide_val['template']
            is_valid = slide_val['valid']
            
            # Create expandable section for each slide
            status_icon = "✅" if is_valid else "❌"
            warning_icon = " ⚠️" if slide_val['warnings'] else ""
            
            with st.expander(f"Slide {slide_num}: {template} {status_icon}{warning_icon}"):
                
                if not is_valid:
                    # Critical issues
                    if slide_val['issues']:
                        st.markdown("**🚨 Critical Issues:**")
                        for issue in slide_val['issues']:
                            st.markdown(f"  • {issue}")
                    
                    # Missing fields
                    if slide_val['missing_fields']:
                        st.markdown("**📝 Missing Fields:**")
                        for field in slide_val['missing_fields']:
                            st.markdown(f"  • {field}")
                    
                    # Empty fields
                    if slide_val['empty_fields']:
                        st.markdown("**📦 Empty/Placeholder Fields:**")
                        for field in slide_val['empty_fields']:
                            st.markdown(f"  • {field}")
                
                # Warnings (even for valid slides)
                if slide_val['warnings']:
                    st.markdown("**⚠️ Warnings:**")
                    for warning in slide_val['warnings']:
                        st.markdown(f"  • {warning}")
                
                if is_valid and not slide_val['warnings']:
                    st.success("All required content present - no empty boxes expected")
    
    return validation_results['overall_valid']

def create_validation_feedback_for_llm(validation_results):
    """Create specific feedback for the LLM to fix validation issues with example-based guidance"""
    
    if validation_results['overall_valid']:
        return None  # No feedback needed
    
    feedback_sections = []
    feedback_sections.append("❌ VALIDATION FAILED - Your JSONs have empty boxes and missing content that must be fixed before generating the deck.")
    feedback_sections.append("\n🎯 ZERO EMPTY BOXES POLICY VIOLATIONS:")
    
    # Add structure validation feedback first
    if 'structure_validation' in validation_results and validation_results['structure_validation']['structure_issues']:
        feedback_sections.append("\n🗃️ STRUCTURAL ISSUES (compared to professional examples):")
        for issue in validation_results['structure_validation']['structure_issues']:
            feedback_sections.append(f"    - {issue}")
        
        feedback_sections.append("\n📋 STRUCTURE REQUIREMENTS:")
        feedback_sections.append("  Content IR must include these key sections:")
        feedback_sections.append("    - entities: {company: {name: 'Company Name'}}")
        feedback_sections.append("    - management_team: {left_column_profiles: [...], right_column_profiles: [...]}")
        feedback_sections.append("    - strategic_buyers: [{buyer_name, strategic_rationale, fit}, ...]")
        feedback_sections.append("    - financial_buyers: [{buyer_name, strategic_rationale, fit}, ...]")
        
        feedback_sections.append("\n  Each management profile must have:")
        feedback_sections.append("    - role_title: 'Chief Executive Officer'")
        feedback_sections.append("    - experience_bullets: ['bullet 1', 'bullet 2', ...]")
        
        feedback_sections.append("\n  Each buyer must have:")
        feedback_sections.append("    - buyer_name: 'Company Name'")
        feedback_sections.append("    - strategic_rationale: 'reason for acquisition'")
        feedback_sections.append("    - fit: 'High (9/10)' or similar")
    
    # Add slide-specific issues
    for slide_val in validation_results['slide_validations']:
        if not slide_val['valid']:
            slide_num = slide_val['slide_number']
            template = slide_val['template']
            
            feedback_sections.append(f"\nSlide {slide_num} ({template}):")
            
            if slide_val['issues']:
                feedback_sections.append("  🚨 Critical Issues:")
                for issue in slide_val['issues']:
                    feedback_sections.append(f"    - {issue}")
                    
                    # Add specific fix instructions for common issues
                    if "Missing content_ir_key" in issue and template == "buyer_profiles":
                        feedback_sections.append("      FIX: Add 'content_ir_key': 'strategic_buyers' or 'content_ir_key': 'financial_buyers' to the slide object (not in data section)")
                        feedback_sections.append("      EXAMPLE:")
                        feedback_sections.append("      {")
                        feedback_sections.append("        'template': 'buyer_profiles',")
                        feedback_sections.append("        'content_ir_key': 'strategic_buyers',")
                        feedback_sections.append("        'data': {")
                        feedback_sections.append("          'title': 'Strategic Buyers - Global Healthcare Leaders',")
                        feedback_sections.append("          'table_headers': ['Buyer Profile', 'Strategic Rationale', 'Fit']")
                        feedback_sections.append("        }")
                        feedback_sections.append("      }")
            
            if slide_val['missing_fields']:
                feedback_sections.append("  📝 Missing Required Fields:")
                for field in slide_val['missing_fields']:
                    feedback_sections.append(f"    - {field}")
            
            if slide_val['empty_fields']:
                feedback_sections.append("  📦 Empty/Placeholder Content (will create empty boxes):")
                for field in slide_val['empty_fields']:
                    feedback_sections.append(f"    - {field}")
    
    # Add specific buyer_profiles fix instructions with real examples
    has_buyer_issues = any("buyer_profiles" in slide_val['template'] for slide_val in validation_results['slide_validations'] if not slide_val['valid'])
    if has_buyer_issues:
        feedback_sections.append("\n🔧 BUYER_PROFILES SLIDE FIX INSTRUCTIONS:")
        feedback_sections.append("CRITICAL: buyer_profiles slides must reference buyer data using content_ir_key")
        feedback_sections.append("\nCORRECT EXAMPLE - Strategic Buyers:")
        feedback_sections.append('{')
        feedback_sections.append('  "template": "buyer_profiles",')
        feedback_sections.append('  "content_ir_key": "strategic_buyers",')
        feedback_sections.append('  "data": {')
        feedback_sections.append('    "title": "Strategic Buyers - Global Healthcare Leaders",')
        feedback_sections.append('    "table_headers": ["Buyer Profile", "Strategic Rationale", "Key Synergies", "Fit"]')
        feedback_sections.append('  }')
        feedback_sections.append('}')
        
        feedback_sections.append("\nCORRECT EXAMPLE - Financial Buyers:")
        feedback_sections.append('{')
        feedback_sections.append('  "template": "buyer_profiles",')
        feedback_sections.append('  "content_ir_key": "financial_buyers",')
        feedback_sections.append('  "data": {')
        feedback_sections.append('    "title": "Financial Buyers - Global Private Equity",')
        feedback_sections.append('    "table_headers": ["Fund Profile", "Healthcare Strategy", "Fit"]')
        feedback_sections.append('  }')
        feedback_sections.append('}')
        
        feedback_sections.append("\nThe Content IR must have matching arrays:")
        feedback_sections.append('"strategic_buyers": [')
        feedback_sections.append('  {')
        feedback_sections.append('    "buyer_name": "UnitedHealth / Optum",')
        feedback_sections.append('    "strategic_rationale": "SEA market entry with established platform",')
        feedback_sections.append('    "key_synergies": "Data analytics, technology platform",')
        feedback_sections.append('    "fit": "High (9/10)"')
        feedback_sections.append('  }')
        feedback_sections.append(']')
    
    feedback_sections.append(f"\n📊 QUALITY SCORES:")
    if 'data_quality_score' in validation_results:
        feedback_sections.append(f"  Data Quality: {validation_results['data_quality_score']:.0f}% (need 90%+)")
    if 'completeness_score' in validation_results:
        feedback_sections.append(f"  Completeness: {validation_results['completeness_score']:.0f}% (need 90%+)")
    
    feedback_sections.append("\n✅ TO FIX: Please regenerate the JSONs with complete content for all the issues listed above. Follow the professional examples exactly. Every field must have real data, not placeholders or empty values.")
    
    return "\n".join(feedback_sections)

# Load templates and examples for the system prompt
def load_templates_json():
    """Load templates.json for the system prompt"""
    try:
        templates_path = Path("templates.json")
        if templates_path.exists():
            with open(templates_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        else:
            return []
    except Exception as e:
        st.error(f"Error loading templates.json: {e}")
        return []

def load_example_files():
    """Load the example JSON files to include in system prompt"""
    examples = {}
    
    # Try to load complete_content_ir.json
    try:
        content_ir_path = Path("complete_content_ir.json")
        if content_ir_path.exists():
            with open(content_ir_path, 'r', encoding='utf-8') as f:
                examples['content_ir'] = json.load(f)
    except Exception as e:
        print(f"Could not load complete_content_ir.json: {e}")
    
    # Try to load complete_render_plan.json
    try:
        render_plan_path = Path("complete_render_plan.json")
        if render_plan_path.exists():
            with open(render_plan_path, 'r', encoding='utf-8') as f:
                examples['render_plan'] = json.load(f)
    except Exception as e:
        print(f"Could not load complete_render_plan.json: {e}")
    
    # If files don't exist, use the embedded examples
    if 'content_ir' not in examples:
        examples['content_ir'] = {
            "entities": {
                "company": {
                    "name": "SouthernCapital Healthcare"
                }
            },
            "facts": {
                "years": ["2020", "2021", "2022", "2023", "2024E"],
                "revenue_usd_m": [120, 145, 180, 210, 240],
                "ebitda_usd_m": [18, 24, 31, 40, 47],
                "ebitda_margins": [15.0, 16.6, 17.2, 19.0, 19.6]
            },
            "management_team": {
                "left_column_profiles": [
                    {
                        "role_title": "CEO - Michael Chen",
                        "experience_bullets": [
                            "25+ years healthcare industry experience across hospital operations",
                            "Former Regional VP at major international hospital group",
                            "MBA from top-tier business school with healthcare specialization",
                            "Led successful expansion of 40+ healthcare facilities"
                        ]
                    },
                    {
                        "role_title": "CFO - Sarah Martinez",
                        "experience_bullets": [
                            "15+ years finance leadership in healthcare services",
                            "Ex-CFO at publicly-traded healthcare services company",
                            "CPA with proven M&A integration track record",
                            "Successfully completed 8 acquisitions totaling $200M+",
                            "Deep expertise in healthcare reimbursement"
                        ]
                    }
                ],
                "right_column_profiles": [
                    {
                        "role_title": "COO - David Park",
                        "experience_bullets": [
                            "20+ years multi-site healthcare operations experience",
                            "Successfully scaled 50+ clinic locations across SEA",
                            "Lean Six Sigma Master Black Belt certification",
                            "Former Regional Operations Director at international chain"
                        ]
                    }
                ]
            },
            "strategic_buyers": [
                {
                    "buyer_name": "UnitedHealth / Optum",
                    "description": "Global healthcare leader with $350B+ revenue",
                    "strategic_rationale": "SEA market entry with established platform",
                    "key_synergies": "Data analytics, technology platform, corporate relationships",
                    "fit": "High (9/10)"
                }
            ],
            "financial_buyers": [
                {
                    "buyer_name": "Blackstone Growth",
                    "description": "$975B AUM, $40B+ healthcare investments",
                    "strategic_rationale": "Buy-and-build platform strategy across SEA",
                    "key_synergies": "Operational excellence, technology investment",
                    "fit": "Very High (9/10)"
                }
            ]
        }
    
    if 'render_plan' not in examples:
        examples['render_plan'] = {
            "slides": [
                {
                    "template": "business_overview",
                    "data": {
                        "title": "Business & Operational Overview",
                        "description": "Leading integrated healthcare services platform in Southeast Asia with comprehensive medical care across multiple countries. The company operates premium clinic locations serving both individual patients and corporate clients, with established market presence and proven operational excellence in healthcare delivery and patient care management.",
                        "timeline": {
                            "start_year": "2015",
                            "end_year": "2024", 
                            "years_note": "(9+ years of healthcare leadership and expansion)"
                        },
                        "highlights": [
                            "Market-leading network with 35+ premium clinic locations across Singapore, Malaysia, Indonesia, and Philippines",
                            "Strong patient engagement with 125,000+ annual patient visits and exceptional 89% retention rate demonstrating quality care",
                            "Diversified revenue base with 65+ corporate wellness contracts covering major employers and multinational corporations",
                            "Advanced healthcare technology platform with integrated digital health solutions and telemedicine capabilities",
                            "Board-certified medical specialists across multiple disciplines ensuring comprehensive care delivery",
                            "International healthcare accreditation and quality certifications meeting global standards",
                            "Proven scalable business model with consistent growth in patient volume and geographic expansion",
                            "Strong ESG commitment with community health initiatives and sustainable healthcare practices"
                        ],
                        "services": [
                            "Primary Care & Preventive Medicine: Comprehensive health screenings, vaccinations, and preventive care programs",
                            "Specialty Medical Services: Cardiology, orthopedics, dermatology, and other specialized medical treatments",
                            "Diagnostic Imaging & Laboratory: Advanced imaging technology, comprehensive lab testing, and diagnostic services",
                            "Corporate Wellness Programs: Employee health assessments, workplace wellness initiatives, and occupational health",
                            "Digital Health & Telemedicine: Remote consultations, health monitoring apps, and digital patient engagement",
                            "Executive Health Assessments: Comprehensive VIP health packages for corporate executives and high-net-worth individuals",
                            "Emergency & Urgent Care: 24/7 emergency services and urgent care facilities across clinic network",
                            "Health Education & Training: Patient education programs, health workshops, and medical training services"
                        ],
                        "positioning_desc": "The company has established itself as the premier healthcare services provider in Southeast Asia, serving both individual patients and corporate clients with comprehensive medical services, advanced technology platforms, and exceptional care standards that drive high patient satisfaction and retention rates across multiple markets."
                    }
                },
                {
                    "template": "buyer_profiles",
                    "content_ir_key": "strategic_buyers",
                    "data": {
                        "title": "Strategic Buyers - Global Healthcare Leaders",
                        "table_headers": ["Buyer Profile", "Strategic Rationale", "Key Synergies", "Fit"]
                    }
                },
                {
                    "template": "buyer_profiles",
                    "content_ir_key": "financial_buyers", 
                    "data": {
                        "title": "Financial Buyers - Global Private Equity",
                        "table_headers": ["Fund Profile", "Healthcare Strategy", "Value Creation", "Fit"]
                    }
                }
            ]
        }
    
    return examples

# Enhanced validation using real-world examples
def validate_against_examples(content_ir, render_plan, examples):
    """Validate generated JSONs against real-world example structures"""
    validation_results = {
        'content_ir_structure_valid': True,
        'render_plan_structure_valid': True,
        'structure_issues': [],
        'data_quality_score': 0,
        'completeness_score': 0
    }
    
    # Validate Content IR structure against example
    if 'content_ir' in examples:
        example_content_ir = examples['content_ir']
        
        # Check for key sections that should exist
        expected_sections = ['entities', 'management_team', 'strategic_buyers', 'financial_buyers']
        for section in expected_sections:
            if section in example_content_ir and section not in content_ir:
                validation_results['structure_issues'].append(f"Missing key Content IR section: {section}")
                validation_results['content_ir_structure_valid'] = False
        
        # Check management team structure
        if 'management_team' in content_ir:
            mgmt = content_ir['management_team']
            example_mgmt = example_content_ir.get('management_team', {})
            
            for column in ['left_column_profiles', 'right_column_profiles']:
                if column in example_mgmt and column not in mgmt:
                    validation_results['structure_issues'].append(f"Missing management team section: {column}")
                elif column in mgmt and isinstance(mgmt[column], list):
                    # Check profile structure
                    for i, profile in enumerate(mgmt[column]):
                        if 'role_title' not in profile:
                            validation_results['structure_issues'].append(f"Profile {i+1} in {column} missing role_title")
                        if 'experience_bullets' not in profile or not isinstance(profile['experience_bullets'], list):
                            validation_results['structure_issues'].append(f"Profile {i+1} in {column} missing experience_bullets array")
        
        # Check buyer arrays structure
        for buyer_type in ['strategic_buyers', 'financial_buyers']:
            if buyer_type in content_ir and isinstance(content_ir[buyer_type], list):
                for i, buyer in enumerate(content_ir[buyer_type]):
                    required_fields = ['buyer_name', 'strategic_rationale', 'fit']
                    for field in required_fields:
                        if field not in buyer:
                            validation_results['structure_issues'].append(f"{buyer_type} #{i+1} missing required field: {field}")
    
    # Validate Render Plan structure
    if 'render_plan' in examples and 'slides' in render_plan:
        example_slides = examples['render_plan']['slides']
        
        # Check for buyer_profiles slides using content_ir_key
        buyer_slides = [s for s in render_plan['slides'] if s.get('template') == 'buyer_profiles']
        
        for slide in buyer_slides:
            if 'content_ir_key' not in slide:
                validation_results['structure_issues'].append(f"buyer_profiles slide missing content_ir_key")
                validation_results['render_plan_structure_valid'] = False
            elif slide['content_ir_key'] not in content_ir:
                validation_results['structure_issues'].append(f"content_ir_key '{slide['content_ir_key']}' not found in Content IR")
                validation_results['render_plan_structure_valid'] = False
    
    # Calculate quality scores
    total_sections = len(['entities', 'management_team', 'strategic_buyers', 'financial_buyers', 'historical_financials'])
    present_sections = sum(1 for section in ['entities', 'management_team', 'strategic_buyers', 'financial_buyers', 'historical_financials'] if section in content_ir)
    validation_results['completeness_score'] = (present_sections / total_sections) * 100
    
    # Data quality score based on structure compliance
    validation_results['data_quality_score'] = max(0, 100 - (len(validation_results['structure_issues']) * 10))
    
    return validation_results

# Load templates and examples
TEMPLATES = load_templates_json()
EXAMPLES = load_example_files()

# Create example sections for system prompt
def create_examples_text():
    """Create formatted examples text for system prompt"""
    examples_text = ""
    
    if 'content_ir' in EXAMPLES:
        examples_text += "\n\nEXAMPLE CONTENT IR STRUCTURE:\n"
        examples_text += "```json\n"
        examples_text += json.dumps(EXAMPLES['content_ir'], indent=2)
        examples_text += "\n```\n"
    
    if 'render_plan' in EXAMPLES:
        examples_text += "\n\nEXAMPLE RENDER PLAN STRUCTURE:\n"
        examples_text += "```json\n"
        examples_text += json.dumps(EXAMPLES['render_plan'], indent=2)
        examples_text += "\n```\n"
    
    return examples_text

# UPDATED Enhanced System Prompt with CORRECT Field Names
SYSTEM_PROMPT = """
You are a precise, on-task investment banking pitch deck copilot that generates COMPLETE, DOWNLOADABLE JSON files.

🎯 **ZERO EMPTY BOXES POLICY**: Every slide must have complete content - no empty sections, boxes, or placeholder text.

SPECIFIC SLIDE REQUIREMENTS FOR ALL TEMPLATES (UPDATED WITH CORRECT FIELD NAMES):

1. **management_team**:
   - Must have left_column_profiles and right_column_profiles (2-3 profiles each for optimal layout)
   - Each profile needs: role_title WITH NAME, experience_bullets (array of 3-4 concise bullets)
   - CORRECT STRUCTURE: {{"role_title": "CEO - John Smith" OR "Chief Executive Officer - Sarah Johnson", "experience_bullets": ["brief bullet1", "brief bullet2", ...]}}
   - CONTENT LIMITS: Role titles max 50 chars, bullets max 80 chars each to prevent text overlap
   - NAMES REQUIRED: Always include actual executive names in role_title field
   - LAYOUT OPTIMIZATION: Max 3 profiles per column, max 4 bullets per profile for clean presentation

2. **business_overview**:
   - Must have: title, description, highlights (min 6-8 detailed items), services (min 6-8 detailed items), positioning_desc
   - RICH CONTENT DENSITY like iCar Asia example - each highlight must be comprehensive and detailed
   - Highlights format: "Detailed description with specific metrics, numbers, and context" (NOT simple titles)
   - Services format: "Service Category: Detailed description with scope and capabilities" (category + colon + details)
   - Examples of GOOD highlights: "Market-leading position with 130+ outlets across major Indonesian cities and strategic transit hubs", "Strong growth trajectory with 100+ new locations added since 2021 acquisition demonstrating scalable business model"
   - Examples of GOOD services: "Upstream Operations: Oil and gas exploration, development, and production activities across multiple basins", "Digital Solutions: Advanced analytics, AI, and digital transformation programs for operational efficiency"
   - Timeline should include years_note with descriptive context like "(90+ years of operations and global leadership)"
   - All fields must be comprehensive business descriptions, not simple bullet points or placeholders

3. **product_service_footprint**:
   - Must have: title, services array with complete title AND desc for each
   - Services array needs minimum 4 entries with structure: {{"title": "Service Name", "desc": "Description"}}
   - CRITICAL: coverage_table must have 3-4 columns for proper market comparison
   - Coverage table structure: [["Region", "Market Segment", "Assets/Products", "Coverage Details"], ["Saudi Arabia", "Upstream", "Oil fields, refineries", "Ghawar, Safaniya fields"], ...]
   - Example 3-column format: [["Region", "Business Segment", "Market Position"], ["Americas", "Downstream", "Leading refinery network"]]
   - Example 4-column format: [["Region", "Segment", "Assets", "Market Share"], ["Asia", "Petrochemicals", "JV facilities", "15% market share"]]
   - Must include metrics data for operational metrics section
   - NO empty boxes in layout areas - all sections must be populated

4. **buyer_profiles**:
   - Must use content_ir_key to reference buyer data (PREFERRED METHOD)
   - OR use complete table_rows with actual data (fallback method)
   - Each buyer must have complete: buyer_name, strategic_rationale (5-10 words max), key_synergies, fit
   - STRATEGIC RATIONALE: Must be concise 5-10 words explaining strategic value (e.g., "Regional expansion into Southeast Asian markets")
   - FIT: Assess strategic alignment with concise rating/explanation (e.g., "High (9/10)", "Strong synergies", "Good cultural fit")
   - Tables must populate with real data, not be empty
   - Example correct structure:
     ```json
     {{
       "template": "buyer_profiles",
       "content_ir_key": "strategic_buyers", 
       "data": {{
         "title": "Strategic Buyer Profiles",
         "table_headers": ["Buyer Name", "Description", "Strategic Rationale", "Key Synergies", "Fit"]
       }}
     }}
     ```

5. **historical_financial_performance**:
   - Must have chart data with categories, revenue, ebitda arrays (min 3 years each)
   - Must include key_metrics with metrics array
   - Chart structure: {{"categories": ["2020", "2021", ...], "revenue": [120, 145, ...], "ebitda": [18, 24, ...]}}

6. **margin_cost_resilience**:
   - Must have: cost_management with items array, risk_mitigation with main_strategy
   - CORRECT FIELD NAMES: cost_management (not cost_structure), risk_mitigation (not resilience_factors)
   - Structure: {{"cost_management": {{"items": [...]}}, "risk_mitigation": {{"main_strategy": {{...}}}}}}

7. **competitive_positioning**:
   - CRITICAL: Must match iCar Asia format with 5-column assessment table
   - Must have: competitors array, assessment table (5 columns), advantages array, barriers array
   - Competitors structure: [{{"name": "Company", "revenue": 450}}, ...] (for chart)
   - ASSESSMENT TABLE STRUCTURE (5 columns): [["Company", "Market Share", "Tech Platform", "Coverage", "Revenue (M)"], ["Our Company", "⭐⭐⭐⭐", "⭐⭐⭐⭐", "⭐⭐⭐⭐", "$200M"], ...]
   - Star ratings: Use 1-5 stars (⭐) or numeric ratings 1-5 that get converted to stars
   - Revenue column: Include quantitative data like "$200M", "€150M", "¥1.2B"
   - Advantages: Array of concise competitive advantages (4-6 items max)
   - Barriers: Array of market entry barriers (4-6 items max)
   - Content limits: Each advantage/barrier max 80 chars for clean layout
   - CORRECT FIELD NAMES: advantages (not competitive_advantages), assessment (not competitive_analysis)

8. **valuation_overview**:

9. **sea_conglomerates** (Strategic Buyers with Financials):
   - If buyer rows include key financial metrics (e.g., revenue, EBITDA, market cap, net income, margins, ticker, ownership, EV/valuation), **use the `sea_conglomerates` template** instead of `buyer_profiles`.
   - Table structure: 6 columns ["Name", "Country", "Description", "Key shareholders", "Key financials (US$m)", "Contact"]
   - CONTACT COLUMN: Use "Contact" header (not "Moelis contact") - assign Moelis team members or "To be assigned"
   - For content, provide an array of objects with: `name`, `country`, `description` (starts with financials + brief rationale), `key_shareholders`, `key_financials`, `contact`
   - Example item: `{ "name": "Yamazaki Baking Co.", "country": "Japan", "description": "Revenue: US$5.2B • EBITDA: US$520M • Rationale: Japan → SEA expansion", "contact": "John Smith" }`

8'. **valuation_overview (cont.)**:
   - Must have: valuation_data array with proper methodology grouping (ORIGINAL FORMAT REQUIRED)
   - CRITICAL: Include methodology_type field for colored section grouping (precedent_transactions, trading_comps, dcf)
   - Structure: {{"valuation_data": [{{"methodology_type": "precedent_transactions", "methodology": "Precedent Transactions", "commentary": "Based on regional deals...", "enterprise_value": "US$200M", "metric": "EV/Revenue", "22a_multiple": "14.3x", "23e_multiple": "12.8x"}}]}}
   - Required fields per row: methodology_type, methodology, commentary, enterprise_value, metric, 22a_multiple, 23e_multiple
   - Methodology types: "precedent_transactions", "trading_comps", "dcf" (produces colored section tabs)
   - ORIGINAL FORMAT: Colored methodology sections on left side with proper visual grouping

9. **growth_strategy_projections**:
   - Must have: growth_strategy with strategies array, financial_projections
   - May have slide_data wrapper: {{"slide_data": {{"growth_strategy": {{...}}}}}}

10. **precedent_transactions**:
    - Must have: transactions array with target, acquirer, date, enterprise_value, revenue, ev_revenue_multiple
    - Each transaction needs complete data, no placeholders

CONTENT IR STRUCTURE REQUIREMENTS:
- entities: {{"company": {{"name": "Company Name"}}}}
- management_team: {{"left_column_profiles": [...], "right_column_profiles": [...]}}
- strategic_buyers: [{{"buyer_name": "Name", "strategic_rationale": "5-10 word rationale", "fit": "High (9/10)"}}, ...]
- financial_buyers: [{{"buyer_name": "Name", "strategic_rationale": "5-10 word rationale", "fit": "High (9/10)"}}, ...]

VALIDATION BEFORE OUTPUT:
Before generating JSONs, verify each slide will have NO EMPTY BOXES:
- All required fields populated with real data using CORRECT field names
- All arrays have minimum required items
- All chart/table areas have supporting data
- No placeholder text like [COMPANY], [AMOUNT], etc.
- Every content area will render with actual information

If ANY slide would have empty boxes, ask for more information instead of generating incomplete JSONs.

CRITICAL SUCCESS METRICS:
🎯 Generate Content IR with ALL collected data using correct field names
🎯 Generate Render Plan with 8+ diverse slides
🎯 Create files immediately when interview is complete
🎯 Ensure files are production-ready for deck generation

MANDATORY COMPLETION CHECKLIST:
✅ Company name and business description
✅ Investment highlights and value propositions
✅ Business overview (model, operations, positioning)
✅ Product/service footprint (offerings, geography)
✅ Historical financials (3-5 years, revenue, EBITDA, margins)
✅ Margin/cost resilience analysis
✅ Growth strategy with market data and projections
✅ Management team (4-6 profiles with role_title/experience_bullets)
✅ Investor considerations (risks and opportunities)
✅ Competitive positioning (if requested)
✅ Trading precedents (public comps and/or private deals)
✅ Valuation methodologies and assumptions
✅ Strategic buyers (3-4 with rationale)
✅ Financial buyers (3-4 PE firms with rationale)

ENHANCED INTERVIEW FLOW RULES:

1. **COMPLETENESS CHECK BEFORE PROGRESSION**:
   - After EVERY user response, analyze if ALL required information for the current topic is collected
   - If information is missing, ask SPECIFIC follow-up questions about what's missing
   - List exactly what information you still need
   - Only move to next topic when current topic is COMPLETE or user explicitly skips

2. **HANDLE "I DON'T KNOW" RESPONSES**:
   - If user says "I don't know" or provides incomplete information, say:
     "I can search for that information. Let me look it up for you."
   - Use web search to find the missing information
   - Show the user exactly what you found with sources
   - Ask: "I found [specific information]. Is it OK to use this for your pitch deck?"
   - Wait for explicit user consent before using the information
   - NEVER use information without asking permission first

3. **SKIP FUNCTIONALITY**:
   - If user says "skip this slide" or "skip this topic", acknowledge and move to next topic
   - Mark that topic as skipped (do not include in final JSON)
   - Continue with remaining topics

4. **SPECIFIC FOLLOW-UP REQUIREMENTS**:
   - For each topic, have specific required fields
   - Ask targeted questions for missing fields
   - Example: "I have your company name but still need: founding year, legal structure, and primary markets. Can you provide these?"

5. **COMMUNICATION STYLE**:
   - Ask 1-2 focused questions per response
   - Be specific about exactly what information is missing
   - Never provide additional information until you've asked follow-up questions
   - Don't summarize or explain - focus on getting missing data

INTERVIEW SEQUENCE (ask in this order, collecting ALL required info for each):

1. **Company Overview**: 
   Required: Name, business description, founding year, legal structure, core operations, target markets
   
2. **Investment Highlights**: 
   Required: 3-5 key value propositions, competitive advantages, unique selling points
   
3. **Business Model**: 
   Required: Revenue streams, how company makes money, customer segments, pricing model
   
4. **Product/Service Footprint**: 
   Required: Main offerings, geographic presence, market positioning, distribution channels
   
5. **Historical Financials**: 
   Required: 3-5 years of revenue, EBITDA, margins, growth rates, key financial metrics
   
6. **Margin/Cost Resilience**: 
   Required: Cost structure breakdown, margin stability factors, competitive moats
   
7. **Growth Strategy**: 
   Required: Expansion plans, market size data, growth projections, strategic initiatives
   
8. **Management Team**: 
   Required: 4-6 executives with role_title, experience_bullets array for each
   
9. **Investor Considerations**: 
   Required: 3-4 key risks, 3-4 key opportunities, mitigation strategies
   
10. **Competitive Positioning**: (Ask: "Do you want a competitive positioning slide?")
    If yes, required: Main competitors, competitive advantages, market positioning
    
11. **Trading Precedents**: 
    Required: Public comparables OR private transactions (ask preference), multiples, rationale
    
12. **Valuation Overview**: 
    Required: Methodologies to use, key multiples, assumptions, valuation range
    
13. **Strategic Buyers**: 
    Required: 3-4 potential acquirers with buyer_name and specific strategic_rationale for each
    
14. **Financial Buyers**: 
    Required: 3-4 PE firms/sponsors with buyer_name and specific strategic_rationale for each

RESPONSE INTERPRETATION:
- "I don't know" = Offer to search and get consent
- "Skip this slide/topic" = Mark as skipped, move to next
- Partial information = Ask specific follow-ups for missing fields
- Complete information = Move to next topic
- Brief confirmatory responses ("yes", "correct") = Move to next topic if current is complete

AUTO-GENERATION TRIGGER:
When you have collected information for ALL non-skipped items in the completion checklist, immediately respond with:

"Perfect! I now have all the information needed to create your comprehensive pitch deck. Here are your complete, downloadable pitch deck files:

## CONTENT IR JSON:
```json
[INSERT COMPLETE CONTENT IR JSON WITH ALL COLLECTED DATA USING CORRECT FIELD NAMES]
```

## RENDER PLAN JSON:
```json
[INSERT COMPLETE RENDER PLAN JSON WITH ALL SLIDES]
```

These files are now ready for download and can be used directly with your pitch deck generation system!"

CRITICAL JSON GENERATION REQUIREMENTS:
- Generate JSONs for ALL discussed slides (minimum 8-12 slides)
- EXCLUDE any slides that were explicitly skipped
- Use ALL collected information across appropriate slide templates
- Follow the EXACT template structure from the examples below
- Include every piece of data collected during the interview
- NEVER leave placeholder text or empty fields
- USE CORRECT FIELD NAMES as specified above

AVAILABLE SLIDE TEMPLATES:
{json.dumps(TEMPLATES, indent=2)}

EXAMPLE JSON STRUCTURES TO FOLLOW EXACTLY:
{create_examples_text()}

REMEMBER: Focus on getting complete, specific information for each topic. Don't move on until you have all required details or explicit user consent to use searched information. Use the CORRECT field names specified above to match the validation system.
"""

# Helper Functions for Interview Flow and File Generation
def analyze_conversation_progress(messages):
    """Analyze conversation to determine what topics have been covered and what's next"""
    conversation_text = " ".join([msg["content"] for msg in messages if msg["role"] != "system"]).lower()
    
    topics_checklist = {
        "company_overview": {
            "keywords": ["company", "business", "what does", "overview", "operations"],
            "covered": False,
            "skipped": "skip" in conversation_text and any(skip_phrase in conversation_text for skip_phrase in ["skip company", "skip overview"]),
            "next_question": "Now let's discuss your investment highlights. What are the key value propositions and competitive advantages that make your company attractive to investors?"
        },
        "investment_highlights": {
            "keywords": ["investment", "highlights", "value proposition", "competitive advantage", "key strengths"],
            "covered": False,
            "skipped": "skip" in conversation_text and any(skip_phrase in conversation_text for skip_phrase in ["skip investment", "skip highlights"]),
            "next_question": "Let's dive into your business model. How exactly does your company make money? What are your main revenue streams?"
        },
        "business_model": {
            "keywords": ["business model", "revenue", "make money", "revenue streams"],
            "covered": False,
            "skipped": "skip" in conversation_text and any(skip_phrase in conversation_text for skip_phrase in ["skip business model", "skip revenue"]),
            "next_question": "Great! Now I need information about your historical financial performance. Can you provide revenue, EBITDA, and margin data for the last 3-5 years?"
        },
        "historical_financials": {
            "keywords": ["revenue", "financial", "ebitda", "margin", "historical", "years", "growth"],
            "covered": False,
            "skipped": "skip" in conversation_text and any(skip_phrase in conversation_text for skip_phrase in ["skip financial", "skip historical"]),
            "next_question": "Now let's discuss your cost structure and margin resilience. How stable are your margins, and what factors help protect your profitability?"
        },
        "margin_resilience": {
            "keywords": ["margin", "cost", "resilience", "stability", "protect", "profitability"],
            "covered": False,
            "skipped": "skip" in conversation_text and any(skip_phrase in conversation_text for skip_phrase in ["skip margin", "skip cost"]),
            "next_question": "Let's talk about your growth strategy. What are your expansion plans, and do you have market size/growth data I can use for charts?"
        },
        "growth_strategy": {
            "keywords": ["growth", "strategy", "expansion", "market size", "projections", "future"],
            "covered": False,
            "skipped": "skip" in conversation_text and any(skip_phrase in conversation_text for skip_phrase in ["skip growth", "skip strategy"]),
            "next_question": "Now I need information about your management team. Can you provide names, titles, and brief backgrounds for 4-6 key executives?"
        },
        "management_team": {
            "keywords": ["management", "team", "executives", "ceo", "founder", "leadership"],
            "covered": False,
            "skipped": "skip" in conversation_text and any(skip_phrase in conversation_text for skip_phrase in ["skip management", "skip team"]),
            "next_question": "What are the key investor considerations - both risks and opportunities - that potential investors should be aware of?"
        },
        "investor_considerations": {
            "keywords": ["risk", "opportunity", "investor", "considerations", "challenges"],
            "covered": False,
            "skipped": "skip" in conversation_text and any(skip_phrase in conversation_text for skip_phrase in ["skip investor", "skip risk"]),
            "next_question": "Do you want a competitive positioning slide comparing your company to key competitors?"
        },
        "competitive_positioning": {
            "keywords": ["competitive", "competitors", "positioning", "comparison"],
            "covered": False,
            "skipped": "skip" in conversation_text and any(skip_phrase in conversation_text for skip_phrase in ["skip competitive", "skip positioning"]),
            "next_question": "For trading precedents, do you want public comparables, private transactions, or both?"
        },
        "trading_precedents": {
            "keywords": ["trading", "precedents", "comparables", "transactions", "multiples"],
            "covered": False,
            "skipped": "skip" in conversation_text and any(skip_phrase in conversation_text for skip_phrase in ["skip trading", "skip precedents"]),
            "next_question": "Let's cover valuation. What valuation methodologies and multiples should we use, and what are your key assumptions?"
        },
        "valuation": {
            "keywords": ["valuation", "multiple", "methodology", "worth", "assumptions"],
            "covered": False,
            "skipped": "skip" in conversation_text and any(skip_phrase in conversation_text for skip_phrase in ["skip valuation", "skip multiple"]),
            "next_question": "Who are 3-4 potential strategic buyers that might be interested in acquiring your company? Please provide rationale for each."
        },
        "strategic_buyers": {
            "keywords": ["strategic", "buyers", "acquirer", "acquisition", "potential"],
            "covered": False,
            "skipped": "skip" in conversation_text and any(skip_phrase in conversation_text for skip_phrase in ["skip strategic", "skip buyers"]),
            "next_question": "Finally, who are 3-4 potential financial buyers (PE firms or sponsors) that might be interested? Please provide rationale."
        },
        "financial_buyers": {
            "keywords": ["financial", "buyers", "private equity", "pe", "sponsors"],
            "covered": False,
            "skipped": "skip" in conversation_text and any(skip_phrase in conversation_text for skip_phrase in ["skip financial buyers", "skip pe"]),
            "next_question": None  # This is the last topic
        }
    }
    
    # Check which topics have been covered or skipped
    covered_count = 0
    skipped_count = 0
    for topic_name, topic_info in topics_checklist.items():
        if topic_info["skipped"]:
            skipped_count += 1
        elif any(keyword in conversation_text for keyword in topic_info["keywords"]):
            topic_info["covered"] = True
            covered_count += 1
    
    # Find next uncovered and unskipped topic
    next_topic = None
    next_question = None
    for topic_name, topic_info in topics_checklist.items():
        if not topic_info["covered"] and not topic_info["skipped"]:
            next_topic = topic_name
            next_question = topic_info["next_question"]
            break
    
    total_applicable_topics = len(topics_checklist) - skipped_count
    completion_percentage = covered_count / total_applicable_topics if total_applicable_topics > 0 else 1.0
    
    return {
        "topics_covered": covered_count,
        "topics_skipped": skipped_count,
        "total_topics": len(topics_checklist),
        "applicable_topics": total_applicable_topics,
        "completion_percentage": completion_percentage,
        "next_topic": next_topic,
        "next_question": next_question,
        "is_complete": covered_count >= total_applicable_topics - 1  # Allow some flexibility
    }

def check_interview_completion(messages):
    """Check if interview has enough information for JSON generation"""
    conversation_text = " ".join([msg["content"] for msg in messages if msg["role"] != "system"])
    
    required_elements = [
        ("company name", ["company", "business name", "firm"]),
        ("business model", ["business model", "how does", "revenue model", "operations"]),
        ("revenue", ["revenue", "sales", "income", "financial performance"]),
        ("EBITDA", ["EBITDA", "earnings", "profit", "margin"]),
        ("management team", ["management", "team", "CEO", "founder", "executive"]),
        ("growth strategy", ["growth", "strategy", "expansion", "future", "projections"]),
        ("valuation", ["valuation", "multiple", "worth", "value"]),
        ("strategic buyers", ["strategic", "buyer", "acquirer", "acquisition"]),
        ("financial buyers", ["financial buyer", "private equity", "PE", "sponsor"])
    ]
    
    completed_count = 0
    for element_name, keywords in required_elements:
        if any(keyword.lower() in conversation_text.lower() for keyword in keywords):
            completed_count += 1
    
    completion_percentage = completed_count / len(required_elements)
    return completion_percentage >= 0.8, completed_count, len(required_elements)





# --- BEGIN: Auto-convert buyer_profiles with financials → sea_conglomerates ---
import os, re as _re

AUTO_USE_SEA_CONGLOMERATES = os.getenv("AUTO_USE_SEA_CONGLOMERATES", "1") not in ("0","false","False","no","No")

_FINANCE_HINTS = {"revenue","ebitda","market_cap","net_income","profit","earnings","margin","ticker","ownership","assets","liabilities","enterprise_value","ev","valuation"}

def _extract_country_from_name(name: str) -> str:
    # e.g., "Yamazaki Baking Co. (Japan)" -> "Japan"
    if not isinstance(name, str):
        return ""
    m = _re.search(r"\(([^)]+)\)\s*$", name.strip())
    return m.group(1).strip() if m else ""

def _dict_row_has_finance(r: dict) -> bool:
    keys = {k.lower() for k in r.keys()}
    return any(k in keys for k in _FINANCE_HINTS)

def _headers_have_finance(headers) -> bool:
    if not isinstance(headers, list): return False
    hl = [str(h).strip().lower() for h in headers]
    return any(any(hint in h for hint in _FINANCE_HINTS) for h in hl)

def convert_buyer_profiles_to_sea_conglomerates(slide: dict) -> dict:
    """
    If buyer_profiles contains financial fields (dict rows or finance headers),
    convert to sea_conglomerates template with concise description lines.
    """
    if slide.get("template") != "buyer_profiles" or not AUTO_USE_SEA_CONGLOMERATES:
        return slide

    data = slide.get("data", {})
    rows = data.get("table_rows", [])
    headers = data.get("table_headers", [])

    finance_mode = False
    dict_rows = []
    if isinstance(rows, list) and rows and isinstance(rows[0], dict):
        dict_rows = rows
        finance_mode = any(_dict_row_has_finance(r) for r in dict_rows)
    elif _headers_have_finance(headers):
        finance_mode = True

    if not finance_mode:
        return slide  # no conversion

    items = []
    if dict_rows:
        for r in dict_rows:
            name = r.get("buyer_name") or r.get("name","")
            country = r.get("country") or _extract_country_from_name(name)
            parts = []

            # Financials first if present
            for k in ("revenue","ebitda","market_cap","net_income","margin","enterprise_value","valuation","ownership","ticker"):
                v = r.get(k)
                if v not in (None, ""):
                    label = k.replace("_"," ").title()
                    parts.append(f"{label}: {v}")

            # Then rationale/synergies for context
            if r.get("strategic_rationale"):
                parts.append(f"Rationale: {r.get('strategic_rationale')}")
            if r.get("key_synergies"):
                parts.append(f"Synergies: {r.get('key_synergies')}")

            desc = " • ".join(parts) if parts else "—"
            items.append({"name": name, "country": country, "description": desc})
    else:
        # If rows are arrays and headers include finance terms, map by position
        # Build index mapping from headers
        idx = {h.strip().lower(): i for i, h in enumerate(headers) if isinstance(h, str)}
        for r in rows:
            name = r[idx.get("buyer name", 0)] if isinstance(r, list) and len(r)>0 else ""
            country = _extract_country_from_name(name)
            parts = []
            for hint in list(_FINANCE_HINTS):
                pos = None
                # try exact header match or contains
                for h, i in idx.items():
                    if hint in h:
                        pos = i; break
                if pos is not None and isinstance(r, list) and len(r)>pos:
                    v = r[pos]
                    if v not in (None, ""):
                        label = hint.replace("_"," ").title()
                        parts.append(f"{label}: {v}")
            # Try rationale/synergies columns
            for key in ["strategic rationale","rationale","key synergies","synergies"]:
                if key in idx and len(r)>idx[key]:
                    val = r[idx[key]]
                    if val not in (None, ""):
                        parts.append(f"{key.title()}: {val}")
            desc = " • ".join(parts) if parts else "—"
            items.append({"name": name, "country": country, "description": desc})

    # Build the new slide
    new_slide = {
        "template": "sea_conglomerates",
        "data": items
    }
    # Preserve original title as an optional leading descriptor if present
    if isinstance(data, dict) and "title" in data:
        # Some renderers might read title; we prepend a descriptor row
        pass

    return new_slide
# --- END: Auto-convert ---
# --- BEGIN: Normalizers to prevent blank cells and schema drift ---
def normalize_buyer_profiles_slide(slide: dict) -> dict:
    if slide.get("template") != "buyer_profiles":
        return slide
    d = slide.setdefault("data", {})

    headers = d.get("table_headers") or ["Buyer Profile", "Strategic Rationale", "Key Synergies", "Fit"]
    if len(headers) == 4:
        # Keep as 4 columns for buyer_profiles
        headers = headers[:4]
    d["table_headers"] = headers[:5]

    fixed_rows = []
    for r in d.get("table_rows", []):
        if isinstance(r, list):
            r = {
                "buyer_name":          (r[0] if len(r) > 0 else ""),
                "strategic_rationale": (r[1] if len(r) > 1 else ""),
                "key_synergies":       (r[2] if len(r) > 2 else ""),
                "fit":                (r[3] if len(r) > 3 else ""),
                "fit_score":           (r[4] if len(r) > 4 else ""),
            }
        else:
            r = dict(r)
            r["buyer_name"]          = r.get("buyer_name") or r.get("name", "")
            r["strategic_rationale"] = r.get("strategic_rationale") or r.get("rationale", "")
            r["key_synergies"]       = r.get("key_synergies") or r.get("synergies", "")
            r["fit"]                = r.get("fit") or r.get("concerns", "")
            # Note: fit is now the primary field, not fit_score
        fixed_rows.append(r)
    d["table_rows"] = fixed_rows

    d.setdefault("subtitle", d.get("subtitle", ""))
    d.setdefault("company", slide.get("company") or "")
    return slide


def normalize_valuation_overview_slide(slide: dict) -> dict:
    if slide.get("template") != "valuation_overview":
        return slide
    d = slide.setdefault("data", {})
    rows = d.get("valuation_data", [])

    any_22a = False
    any_23e = False
    any_metric = False

    for r in rows:
        meth = (r.get("methodology") or "").lower()
        if not r.get("metric"):
            if "precedent" in meth or "trading" in meth:
                r["metric"] = "EV/Revenue"
            elif "dcf" in meth or "discounted" in meth:
                r["metric"] = "DCF"
        any_metric = any_metric or bool(r.get("metric"))

        if "22a_multiple" not in r:
            r["22a_multiple"] = r.get("22A_multiple") or r.get("FY22_multiple") or "-"
        if "23e_multiple" not in r:
            r["23e_multiple"] = r.get("23E_multiple") or r.get("FY23E_multiple") or "-"

        if not r.get("methodology_type"):
            if "precedent" in meth:
                r["methodology_type"] = "precedent_transactions"
            elif "trading" in meth:
                r["methodology_type"] = "trading_comps"
            elif "dcf" in meth or "discounted" in meth:
                r["methodology_type"] = "dcf"

        any_22a = any_22a or (r.get("22a_multiple") not in ("", None))
        any_23e = any_23e or (r.get("23e_multiple") not in ("", None))

    d["__hide_metric_col"]  = not any_metric
    d["__hide_22a_col"]     = not any_22a
    d["__hide_23e_col"]     = not any_23e
    return slide


def normalize_plan(plan: dict) -> dict:
    try:
        slides_in = plan.get("slides", [])
    except Exception:
        return plan
    slides_out = []
    for s in slides_in:
        # Convert finance-heavy buyer profiles into SEA Conglomerates slide first
        s = convert_buyer_profiles_to_sea_conglomerates(s)
        # Then run standard normalizers
        s = normalize_buyer_profiles_slide(s)
        s = normalize_valuation_overview_slide(s)
        slides_out.append(s)
    plan["slides"] = slides_out
    return plan
    slides_out = []
    for s in slides_in:
        s = normalize_buyer_profiles_slide(s)
        s = normalize_valuation_overview_slide(s)
        slides_out.append(s)
    plan["slides"] = slides_out
    return plan
# --- END: Normalizers ---

def extract_and_validate_jsons(response_text):
    """Extract JSONs and perform comprehensive validation with example-based checking"""
    content_ir, render_plan = extract_jsons_from_response(response_text)
    
    # Normalize for downstream validation and rendering
    if isinstance(render_plan, dict):
        render_plan = normalize_plan(render_plan)
    
    if not content_ir or not render_plan:
        return content_ir, render_plan, {
            'overall_valid': False,
            'summary': {'total_slides': 0, 'valid_slides': 0, 'invalid_slides': 0},
            'critical_issues': ['JSONs not found in response']
        }
    
    # Perform comprehensive validation
    validation_results = validate_individual_slides(content_ir, render_plan)
    
    # Add example-based structure validation
    examples = load_example_files()
    structure_validation = validate_against_examples(content_ir, render_plan, examples)
    
    # Merge structure validation results
    validation_results['structure_validation'] = structure_validation
    validation_results['data_quality_score'] = structure_validation['data_quality_score']
    validation_results['completeness_score'] = structure_validation['completeness_score']
    
    # Add structure issues to critical issues if structure is invalid
    if not structure_validation['content_ir_structure_valid'] or not structure_validation['render_plan_structure_valid']:
        validation_results['critical_issues'].extend(structure_validation['structure_issues'])
        validation_results['overall_valid'] = False
    
    return content_ir, render_plan, validation_results

def create_downloadable_files(content_ir, render_plan, company_name="company"):
    """Create downloadable Content IR and Render Plan files"""
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_company_name = "".join(c for c in company_name if c.isalnum() or c in (' ', '-', '_')).rstrip()
    safe_company_name = safe_company_name.replace(' ', '_')
    
    # Create individual files
    content_ir_filename = f"{safe_company_name}_content_ir_{timestamp}.json"
    render_plan_filename = f"{safe_company_name}_render_plan_{timestamp}.json"
    
    # Format JSON with proper indentation
    content_ir_json = json.dumps(content_ir, indent=2, ensure_ascii=False)
    render_plan_json = json.dumps(render_plan, indent=2, ensure_ascii=False)
    
    return {
        'content_ir_filename': content_ir_filename,
        'content_ir_json': content_ir_json,
        'render_plan_filename': render_plan_filename,
        'render_plan_json': render_plan_json,
        'timestamp': timestamp,
        'company_name': safe_company_name
    }

def create_zip_package(files_data):
    """Create a ZIP package with both JSON files and metadata"""
    zip_buffer = io.BytesIO()
    
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        # Add Content IR file
        zip_file.writestr(files_data['content_ir_filename'], files_data['content_ir_json'])
        
        # Add Render Plan file
        zip_file.writestr(files_data['render_plan_filename'], files_data['render_plan_json'])
        
        # Add README with instructions
        readme_content = f"""# AI-Generated Pitch Deck Files
Company: {files_data['company_name']}
Generated: {files_data['timestamp']}

## Files Included:
1. {files_data['content_ir_filename']} - Contains all content data for slides
2. {files_data['render_plan_filename']} - Defines slide structure and templates

## Usage Instructions:
1. Use these files with your pitch deck generation system
2. Load the Content IR for slide data
3. Load the Render Plan for slide structure
4. Generate your PowerPoint presentation

## File Validation:
- Content IR structure: ✓ Complete
- Render Plan structure: ✓ Complete
- Ready for deck generation: ✓ Yes

Generated by AI Deck Builder - LLM-Powered Pitch Deck Generator
"""
        zip_file.writestr("README.txt", readme_content)
        
        # Add metadata file
        metadata = {
            "generated_at": files_data['timestamp'],
            "company_name": files_data['company_name'],
            "content_ir_file": files_data['content_ir_filename'],
            "render_plan_file": files_data['render_plan_filename'],
            "generator": "AI Deck Builder",
            "version": "1.0"
        }
        zip_file.writestr("metadata.json", json.dumps(metadata, indent=2))
    
    zip_buffer.seek(0)
    return zip_buffer

def show_interview_progress(messages):
    """Show progress indicator for interview completion"""
    progress_info = analyze_conversation_progress(messages)
    
    st.sidebar.subheader("🎯 Interview Progress")
    st.sidebar.progress(progress_info["completion_percentage"])
    st.sidebar.write(f"{progress_info['topics_covered']}/{progress_info['applicable_topics']} topics covered")
    
    if progress_info["topics_skipped"] > 0:
        st.sidebar.write(f"⭐ {progress_info['topics_skipped']} topics skipped")
    
    if progress_info["is_complete"]:
        st.sidebar.success("✅ Ready for JSON generation!")
    else:
        remaining = progress_info['applicable_topics'] - progress_info['topics_covered']
        st.sidebar.info(f"📝 {remaining} topics remaining")
    
    return progress_info["is_complete"]

# Initialize brand extractor
brand_extractor = BrandExtractor()

# LLM Integration Functions - FIXED FOR MESSAGE ALTERNATION
def call_llm_api(messages, model_name, api_key, service="perplexity"):
    """Call LLM API (Perplexity or Claude) with the conversation"""
    try:
        if service == "perplexity":
            return call_perplexity_api(messages, model_name, api_key)
        elif service == "claude":
            return call_claude_api(messages, model_name, api_key)
        else:
            return f"Unknown service: {service}"
    except Exception as e:
        return f"Error calling {service} API: {str(e)}"

def call_perplexity_api(messages, model_name, api_key):
    """Call Perplexity API with the conversation - FIXED for message alternation"""
    try:
        url = "https://api.perplexity.ai/chat/completions"
        
        # Extract system message
        system_message = None
        conversation_messages = []
        
        for msg in messages:
            if msg["role"] == "system":
                system_message = msg["content"]
            elif msg["role"] in ["user", "assistant"]:
                conversation_messages.append(msg)
        
        # Build properly alternating conversation
        # Remove any leading assistant messages (Perplexity needs user first after system)
        while conversation_messages and conversation_messages[0]["role"] == "assistant":
            conversation_messages.pop(0)
        
        # Collapse consecutive same-role messages to enforce alternation
        cleaned_messages = []
        for msg in conversation_messages:
            if cleaned_messages and cleaned_messages[-1]["role"] == msg["role"]:
                # Combine consecutive messages of same role
                cleaned_messages[-1]["content"] = cleaned_messages[-1]["content"].rstrip() + "\n\n" + str(msg.get("content", "")).strip()
            else:
                cleaned_messages.append({
                    "role": msg["role"],
                    "content": str(msg.get("content", "")).strip()
                })
        
        # Build final message array for Perplexity
        final_messages = []
        
        # Add system message if present
        if system_message:
            final_messages.append({"role": "system", "content": system_message})
        
        # Add alternating conversation
        final_messages.extend(cleaned_messages)
        
        # Ensure we don't have empty messages
        final_messages = [msg for msg in final_messages if msg.get("content", "").strip()]
        
        payload = {
            "model": model_name,
            "messages": final_messages,
            "temperature": 0.7,
            "max_tokens": 4000,
            "stream": False
        }
        
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json; charset=utf-8"
        }
        
        # Ensure UTF-8 encoding for Unicode characters (emojis, etc.)
        import json
        json_data = json.dumps(payload, ensure_ascii=False)
        response = requests.post(url, data=json_data.encode('utf-8'), headers=headers)
        
        if response.status_code == 200:
            result = response.json()
            return result.get('choices', [{}])[0].get('message', {}).get('content', 'No response')
        else:
            return f"Perplexity API Error: {response.status_code} - {response.text}"
    
    except Exception as e:
        return f"Error calling Perplexity API: {str(e)}"

def call_claude_api(messages, model_name, api_key):
    """Call Claude API with the conversation"""
    try:
        url = "https://api.anthropic.com/v1/messages"
        
        # Convert messages format for Claude
        claude_messages = []
        system_message = ""
        
        for msg in messages:
            if msg["role"] == "system":
                system_message = msg["content"]
            else:
                claude_messages.append({
                    "role": msg["role"],
                    "content": msg["content"]
                })
        
        payload = {
            "model": model_name,
            "max_tokens": 4000,
            "temperature": 0.7,
            "messages": claude_messages
        }
        
        if system_message:
            payload["system"] = system_message
        
        headers = {
            "x-api-key": api_key,
            "Content-Type": "application/json; charset=utf-8",
            "anthropic-version": "2023-06-01"
        }
        
        # Ensure UTF-8 encoding for Unicode characters (emojis, etc.)
        import json
        json_data = json.dumps(payload, ensure_ascii=False)
        response = requests.post(url, data=json_data.encode('utf-8'), headers=headers)
        
        if response.status_code == 200:
            result = response.json()
            return result.get('content', [{}])[0].get('text', 'No response')
        else:
            return f"Claude API Error: {response.status_code} - {response.text}"
    
    except Exception as e:
        return f"Error calling Claude API: {str(e)}"

# Rest of the app.py code follows with sidebar, main interface, etc.
# (The rest of the code remains the same as in the original app.py)

# Sidebar Configuration
with st.sidebar:
    st.header("🤖 AI Configuration")
    
    # LLM Model Selection
    st.subheader("LLM Model")
    
    # LLM Service Selection
    llm_service = st.radio(
        "LLM Service",
        ["🔍 Perplexity (Recommended)", "🧠 Claude (Anthropic)"],
        help="Choose your preferred LLM service"
    )
    
    if llm_service.startswith("🔍"):
        # Perplexity models - UPDATED with current valid model names
        model_options = [
            "sonar-pro",  # Most capable model (replaces sonar-large-online)
            "sonar",  # Standard model (replaces sonar-small-online)
            "sonar-reasoning",  # For complex reasoning tasks
            "sonar-reasoning-pro",  # Advanced reasoning model
            "sonar-deep-research"  # For comprehensive research
        ]
        selected_model = st.selectbox(
            "Choose Perplexity Model",
            model_options,
            index=0,  # Default to sonar-pro (most capable)
            help="sonar-pro offers the best balance of capability and speed"
        )
        api_service = "perplexity"
    else:
        # Claude models
        model_options = [
            "claude-3-5-sonnet-20241022",
            "claude-3-5-haiku-20241022",
            "claude-3-opus-20240229",
            "claude-3-sonnet-20240229",
            "claude-3-haiku-20240307"
        ]
        selected_model = st.selectbox(
            "Choose Claude Model",
            model_options,
            index=0,  # Default to latest Sonnet
            help="Claude Sonnet offers the best balance of speed and capability"
        )
        api_service = "claude"
    
    # API Key Input
    if api_service == "perplexity":
        api_key = st.text_input(
            "Perplexity API Key",
            type="password",
            help="Enter your Perplexity API key"
        )
    else:
        api_key = st.text_input(
            "Claude API Key",
            type="password",
            help="Enter your Anthropic Claude API key"
        )
    
    if not api_key:
        service_name = "Perplexity" if api_service == "perplexity" else "Claude"
        st.warning(f"⚠️ Please enter your {service_name} API key to use the AI copilot")
    
    st.markdown("---")
    
    # File Status Section
    st.subheader("📁 Generated Files Status")
    
    if st.session_state.get("files_ready", False):
        st.success("✅ Files Ready!")
        files_data = st.session_state.get("files_data", {})
        st.write(f"**Company:** {files_data.get('company_name', 'N/A')}")
        st.write(f"**Generated:** {files_data.get('timestamp', 'N/A')}")
        
        if st.button("🔄 Regenerate Files"):
            st.session_state["files_ready"] = False
            st.session_state.pop("files_data", None)
            st.rerun()
    else:
        st.info("📄 Complete interview to generate files")
    
    st.markdown("---")
    
    # Brand Upload Section with LLM Integration
    st.subheader("🎨 Brand Configuration")
    
    # Add extraction method selector
    extraction_method = st.radio(
        "Brand Extraction Method",
        ["🤖 LLM-Powered (Recommended)", "🔧 Rule-Based (Fallback)"],
        help="LLM extraction analyzes slide content contextually for better brand element detection",
        key="extraction_method"
    )
    
    uploaded_brand = st.file_uploader(
        "Upload Brand Deck (PowerPoint)",
        type=['pptx'],
        help="Upload a PowerPoint file to extract colors, fonts, and styling",
        key="brand_upload"
    )
    
    if uploaded_brand is not None and HAS_PPTX:
        try:
            # Show progress
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            use_llm = extraction_method.startswith("🤖")
            
            if use_llm and api_key:
                # Use LLM extraction
                st.write("🤖 **LLM-Powered Brand Extraction**")
                st.info("💡 AI is analyzing your slides to understand brand context and hierarchy")
                
                status_text.text("🧠 AI analyzing slide content and design patterns...")
                progress_bar.progress(20)
                
                uploaded_brand.seek(0)
                brand_config = brand_extractor.extract_brand_from_pptx(
                    uploaded_brand,
                    use_llm=True,
                    api_key=api_key,
                    model_name=selected_model,
                    api_service=api_service
                )
                
                progress_bar.progress(80)
                status_text.text("✅ AI analysis complete!")
                
            else:
                # Use rule-based extraction
                st.write("🔧 **Rule-Based Brand Extraction**")
                if not api_key:
                    st.info("💡 Add your API key above to enable AI-powered brand extraction")
                
                status_text.text("🔍 Analyzing PowerPoint structure...")
                progress_bar.progress(20)
                
                uploaded_brand.seek(0)
                brand_config = brand_extractor.extract_brand_from_pptx(
                    uploaded_brand,
                    use_llm=False
                )
                
                progress_bar.progress(80)
                status_text.text("✅ Rule-based extraction complete!")
            
            progress_bar.progress(100)
            
            # Store configuration
            st.session_state["brand_config"] = brand_config
            
            # Display results
            colors = brand_config.get('color_scheme', {})
            primary = colors.get('primary')
            
            if hasattr(primary, 'r'):
                # Check if we got custom colors or defaults
                if primary.r == 24 and primary.g == 58 and primary.b == 88:
                    st.warning("⚠️ Using default colors - no distinct brand colors detected")
                    if use_llm:
                        st.info("💡 Try uploading a deck with more prominent brand colors or logos")
                else:
                    st.success("✅ Brand elements extracted successfully!")
            
            # Show extracted colors
            st.write("**🎨 Extracted Brand Colors:**")
            color_cols = st.columns(2)
            color_display_order = ['primary', 'secondary', 'accent', 'text']
            
            for i, name in enumerate(color_display_order):
                if name in colors:
                    color = colors[name]
                    if hasattr(color, 'r'):
                        hex_color = f"#{color.r:02x}{color.g:02x}{color.b:02x}"
                        with color_cols[i % 2]:
                            col1, col2 = st.columns([1, 2])
                            with col1:
                                st.color_picker(
                                    f"{name.title()}",
                                    hex_color,
                                    disabled=True,
                                    key=f"color_{name}"
                                )
                            with col2:
                                st.caption(f"RGB({color.r}, {color.g}, {color.b})")
            
            # Show typography if available
            typography = brand_config.get('typography', {})
            if typography:
                st.write("**🔤 Typography:**")
                col1, col2 = st.columns(2)
                with col1:
                    st.write(f"• **Font:** {typography.get('primary_font', 'Arial')}")
                    st.write(f"• **Title Size:** {typography.get('title_size', 24)}pt")
                with col2:
                    st.write(f"• **Body Size:** {typography.get('body_size', 11)}pt")
            
            # Show LLM analysis details if available
            if use_llm and 'llm_analysis' in brand_config:
                with st.expander("🧠 AI Analysis Details"):
                    analysis = brand_config['llm_analysis']
                    
                    # Brand personality
                    if 'brand_personality' in analysis:
                        personality = analysis['brand_personality']
                        if isinstance(personality, dict) and 'description' in personality:
                            st.write(f"**Brand Style:** {personality['description']}")
                    
                    # Color reasoning
                    if 'color_reasoning' in analysis:
                        st.write("**Color Choices:**")
                        for color_type, reasoning in analysis['color_reasoning'].items():
                            if reasoning:
                                st.write(f"• **{color_type.title()}:** {reasoning}")
                    
                    # Font reasoning
                    if 'font_reasoning' in analysis and analysis['font_reasoning']:
                        st.write(f"**Font Choice:** {analysis['font_reasoning']}")
                    
                    # Design patterns
                    if 'design_patterns' in analysis:
                        patterns = analysis['design_patterns']
                        if isinstance(patterns, dict) and 'description' in patterns:
                            st.write(f"**Design Patterns:** {patterns['description']}")
            
            # Clear progress indicators
            progress_bar.empty()
            status_text.empty()
            
        except Exception as e:
            st.error(f"Brand extraction failed: {str(e)}")
            st.error("Please check your PowerPoint file and try again.")
    
    elif uploaded_brand is not None and not HAS_PPTX:
        st.error("⚠️ Cannot process PowerPoint - python-pptx not installed")
        st.code("pip install python-pptx")
    else:
        st.info("📁 Upload a brand deck to extract colors and fonts")
    
    if "brand_config" not in st.session_state:
        st.session_state["brand_config"] = None
    
    st.markdown("---")
    
    # Other configuration options
    templates_path = st.text_input("templates.json path", value="templates.json")
    company_name = st.text_input("Company name", value="Moelis & Company")
    skip_validate = st.checkbox("Skip validation", value=False)

# Initialize chat history
if "messages" not in st.session_state:
    st.session_state.messages = [
        {"role": "system", "content": SYSTEM_PROMPT}
    ]

if "chat_started" not in st.session_state:
    st.session_state.chat_started = False

# Main App Layout
tab_chat, tab_json, tab_execute = st.tabs(["🤖 AI Copilot", "📄 JSON Editor", "⚙️ Execute"])

with tab_chat:
    st.subheader("🤖 Investment Banking Pitch Deck Copilot")
    
    if not api_key:
        st.error("⚠️ Please enter your API key in the sidebar to start the interview")
    else:
        # Start conversation button
        if not st.session_state.chat_started:
            if st.button("🚀 Start Pitch Deck Interview"):
                st.session_state.chat_started = True
                st.rerun()
        
        # Display chat messages
        if st.session_state.chat_started:
            # Display conversation (skip system message)
            display_messages = [m for m in st.session_state.messages if m["role"] != "system"]
            
            # If no conversation yet, show welcome message in UI only
            if not display_messages:
                with st.chat_message("assistant"):
                    welcome_msg = """Hello! I'm your investment banking pitch deck copilot. I'll conduct a comprehensive interview to gather all the information needed for your pitch deck, then automatically generate the complete JSON structures for you.

**What I'll collect:**
- Company overview & business model
- Financial performance & projections  
- Management team profiles
- Growth strategy & market data
- Valuation & trading precedents
- Strategic & financial buyer targets

**New Enhanced Features:**
- I'll ask specific follow-up questions for missing information
- Say "I don't know" and I'll search for information (with your permission)
- Say "skip this slide" to exclude any topic you don't want
- **Zero Empty Boxes Policy**: All slides will have complete content

Let's start: **What is your company name and give me a brief overview of what your business does?**"""
                    st.markdown(welcome_msg)
            
            # Display existing conversation
            for message in display_messages:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])
            
            # Chat input - ENHANCED with comprehensive validation
            if prompt := st.chat_input("Your response...", key="chat_input"):
                # Add user message
                st.session_state.messages.append({"role": "user", "content": prompt})
                
                # Analyze conversation progress
                progress_info = analyze_conversation_progress(st.session_state.messages)
                
                # Show progress in sidebar
                is_complete = show_interview_progress(st.session_state.messages)
                
                # Check if this was a brief confirmatory response or skip request
                brief_confirmatory = prompt.strip().lower() in ["yes", "correct", "that's right", "sounds good", "ok", "okay", "sure", "right"]
                skip_request = "skip" in prompt.lower() and any(skip_phrase in prompt.lower() for skip_phrase in ["skip this", "skip that", "skip slide", "skip topic"])
                
                if brief_confirmatory and progress_info["next_question"] and not is_complete:
                    # User gave brief confirmation - automatically ask next question
                    ai_response = progress_info["next_question"]
                    st.session_state.messages.append({"role": "assistant", "content": ai_response})
                    st.rerun()
                elif skip_request and progress_info["next_question"]:
                    # User wants to skip current topic
                    ai_response = f"Understood, I'll skip this topic. {progress_info['next_question']}"
                    st.session_state.messages.append({"role": "assistant", "content": ai_response})
                    st.rerun()
                else:
                    # Get normal AI response
                    with st.spinner("🤖 Thinking..."):
                        ai_response = call_llm_api(
                            st.session_state.messages,
                            selected_model,
                            api_key,
                            api_service
                        )
                    
                    # Add AI response to history
                    st.session_state.messages.append({"role": "assistant", "content": ai_response})
                    
                    # Check if JSONs were generated and extract them with comprehensive validation
                    content_ir, render_plan, validation_results = extract_and_validate_jsons(ai_response)
                    
                    if content_ir or render_plan:
                        st.success("🎉 JSON structures generated!")
                        
                        # Display comprehensive validation results
                        is_fully_valid = display_validation_results(validation_results)
                        
                        # If validation failed, provide specific feedback to LLM for retry
                        if not is_fully_valid:
                            st.error("🚨 **Validation Failed** - Files cannot be generated with empty boxes!")
                            
                            # Create specific feedback for the LLM
                            llm_feedback = create_validation_feedback_for_llm(validation_results)
                            
                            if llm_feedback:
                                # Show retry button
                                if st.button("🔄 Auto-Fix Validation Issues", type="primary"):
                                    # Add feedback message for LLM to fix issues
                                    st.session_state.messages.append({"role": "user", "content": llm_feedback})
                                    
                                    with st.spinner("🔄 Fixing validation issues..."):
                                        retry_response = call_llm_api(
                                            st.session_state.messages,
                                            selected_model,
                                            api_key,
                                            api_service
                                        )
                                    
                                    st.session_state.messages.append({"role": "assistant", "content": retry_response})
                                    st.rerun()
                        
                        # If validation passed, create downloadable files
                        else:
                            st.balloons()
                            
                            # Extract company name from content IR
                            company_name_extracted = "Unknown_Company"
                            if content_ir and 'entities' in content_ir and 'company' in content_ir['entities']:
                                company_name_extracted = content_ir['entities']['company'].get('name', 'Unknown_Company')
                            
                            # Create downloadable files
                            files_data = create_downloadable_files(content_ir, render_plan, company_name_extracted)
                            
                            # Store in session state
                            st.session_state["generated_content_ir"] = files_data['content_ir_json']
                            st.session_state["generated_render_plan"] = files_data['render_plan_json']
                            st.session_state["files_ready"] = True
                            st.session_state["files_data"] = files_data
                            
                            # Show download section
                            st.markdown("---")
                            st.subheader("📁 Download Your Pitch Deck Files")
                            
                            # Create download columns
                            download_col1, download_col2, download_col3 = st.columns(3)
                            
                            with download_col1:
                                st.download_button(
                                    "📄 Download Content IR",
                                    data=files_data['content_ir_json'],
                                    file_name=files_data['content_ir_filename'],
                                    mime="application/json",
                                    help="Contains all slide content data"
                                )
                            
                            with download_col2:
                                st.download_button(
                                    "📋 Download Render Plan",
                                    data=files_data['render_plan_json'],
                                    file_name=files_data['render_plan_filename'],
                                    mime="application/json",
                                    help="Defines slide structure and templates"
                                )
                            
                            with download_col3:
                                # Create ZIP package
                                zip_buffer = create_zip_package(files_data)
                                zip_filename = f"{files_data['company_name']}_pitch_deck_files_{files_data['timestamp']}.zip"
                                
                                st.download_button(
                                    "📦 Download Complete Package",
                                    data=zip_buffer,
                                    file_name=zip_filename,
                                    mime="application/zip",
                                    help="ZIP package with both files + README"
                                )
                            
                            # Show next steps
                            st.info("""
                            🎯 **Next Steps:**
                            1. Download the files above
                            2. Use them with your pitch deck generation system
                            3. Or switch to the Execute tab to generate the deck directly
                            """)
                    
                    # If interview seems complete but no JSONs generated, prompt for them
                    elif is_complete and not any("CONTENT IR JSON" in msg["content"] for msg in st.session_state.messages):
                        st.warning("📄 Interview appears complete. Prompting AI to generate JSON files...")
                        
                        completion_prompt = """
I believe we have covered all the necessary information for a comprehensive pitch deck. Please generate the complete Content IR JSON and Render Plan JSON structures now using ALL the information I provided during our conversation.

🎯 **ZERO EMPTY BOXES POLICY** - Requirements:
- Include ALL slides we discussed (minimum 8-10 slides)
- EXCLUDE any slides that were explicitly skipped
- Use every piece of data I provided
- Follow the exact JSON format from your examples
- Create multiple slides of the same type if the data supports it
- Don't skip any information or use placeholder text
- Ensure every field has real content (no empty arrays, null values, or placeholders)
- USE CORRECT FIELD NAMES: role_title/experience_bullets for management, cost_management/risk_mitigation for margins, etc.

Please generate both complete JSON structures now with full validation compliance.
"""
                        st.session_state.messages.append({"role": "user", "content": completion_prompt})
                        
                        with st.spinner("🎯 Generating downloadable JSON files..."):
                            completion_response = call_llm_api(
                                st.session_state.messages,
                                selected_model,
                                api_key,
                                api_service
                            )
                        
                        st.session_state.messages.append({"role": "assistant", "content": completion_response})
                    
                    # Force rerun to display new messages
                    st.rerun()
        
        # Clear conversation button
        if st.session_state.chat_started:
            st.markdown("---")
            col1, col2, col3 = st.columns([1, 1, 2])
            
            with col1:
                if st.button("🔥 Reset Chat"):
                    st.session_state.messages = [{"role": "system", "content": SYSTEM_PROMPT}]
                    st.session_state.chat_started = False
                    st.session_state["files_ready"] = False
                    st.session_state.pop("files_data", None)
                    st.rerun()
            
            with col2:
                if st.button("💾 Export Chat"):
                    chat_export = {
                        "model": selected_model,
                        "messages": st.session_state.messages[1:],  # Exclude system message
                        "timestamp": str(pd.Timestamp.now())
                    }
                    
                    st.download_button(
                        "⬇️ Download Chat History",
                        data=json.dumps(chat_export, indent=2),
                        file_name=f"pitch_deck_interview_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.json",
                        mime="application/json"
                    )

with tab_json:
    st.subheader("📄 JSON Editor")
    st.info("💡 **Tip**: Use the AI Copilot to generate the JSON, then copy it here for manual editing if needed")
    
    # Show file status if files are ready
    if st.session_state.get("files_ready", False):
        files_data = st.session_state.get("files_data", {})
        st.success(f"🎉 Using auto-generated files for {files_data.get('company_name', 'your company')}")
        
        with st.expander("📋 Generated Files Summary"):
            st.write(f"**Content IR:** {files_data.get('content_ir_filename', 'N/A')}")
            st.write(f"**Render Plan:** {files_data.get('render_plan_filename', 'N/A')}")
            st.write(f"**Timestamp:** {files_data.get('timestamp', 'N/A')}")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.caption("Content IR JSON")
        content_ir_str = st.text_area(
            "Content IR",
            value=st.session_state.get("generated_content_ir", "{}"),
            height=400,
            help="The Content IR contains all the data for your pitch deck"
        )
    
    with col2:
        st.caption("Render Plan JSON")
        render_plan_str = st.text_area(
            "Render Plan",
            value=st.session_state.get("generated_render_plan", "{}"),
            height=400,
            help="The Render Plan defines which slides to create and their data mapping"
        )
    
    # Validate manually edited JSONs
    if st.button("🔍 Validate Edited JSONs"):
        try:
            # Clean JSONs before parsing
            cleaned_content_ir = clean_json_string(content_ir_str) if content_ir_str.strip() else "{}"
            cleaned_render_plan = clean_json_string(render_plan_str) if render_plan_str.strip() else "{}"
            
            # Show cleaned JSON preview if different from original
            if cleaned_content_ir != content_ir_str.strip():
                st.info("🔧 Content IR was automatically cleaned for parsing")
                with st.expander("View cleaned Content IR"):
                    st.code(cleaned_content_ir[:500] + "..." if len(cleaned_content_ir) > 500 else cleaned_content_ir)
            
            if cleaned_render_plan != render_plan_str.strip():
                st.info("🔧 Render Plan was automatically cleaned for parsing")
                with st.expander("View cleaned Render Plan"):
                    st.code(cleaned_render_plan[:500] + "..." if len(cleaned_render_plan) > 500 else cleaned_render_plan)
            
            content_ir = json.loads(cleaned_content_ir)
            render_plan = json.loads(cleaned_render_plan)
            
            if content_ir and render_plan:
                validation_results = validate_individual_slides(content_ir, render_plan)
                is_valid = display_validation_results(validation_results)
                
                if is_valid:
                    st.success("✅ Manual edits passed validation!")
                    # Update session state with cleaned versions
                    st.session_state["generated_content_ir"] = json.dumps(content_ir, indent=2)
                    st.session_state["generated_render_plan"] = json.dumps(render_plan, indent=2)
                else:
                    st.error("❌ Manual edits have validation issues")
            else:
                st.warning("⚠️ Please provide both Content IR and Render Plan JSONs")
        except json.JSONDecodeError as e:
            st.error(f"❌ JSON Parse Error: {e}")
            st.error("🔧 Try using the 'Clean JSON' button below to fix formatting issues")
            
            # Add JSON cleaning buttons for manual editing
            col1, col2 = st.columns(2)
            with col1:
                if st.button("🔧 Clean Content IR JSON"):
                    cleaned = clean_json_string(content_ir_str)
                    st.session_state["manual_content_ir_cleaned"] = cleaned
                    st.success("Content IR cleaned! Refresh to see changes.")
            
            with col2:
                if st.button("🔧 Clean Render Plan JSON"):
                    cleaned = clean_json_string(render_plan_str)
                    st.session_state["manual_render_plan_cleaned"] = cleaned
                    st.success("Render Plan cleaned! Refresh to see changes.")
        
        except Exception as e:
            st.error(f"❌ Validation Error: {e}")
    
    # Show cleaned versions if available
    if st.session_state.get("manual_content_ir_cleaned"):
        st.subheader("🔧 Cleaned Content IR")
        st.text_area(
            "Cleaned Content IR JSON",
            value=st.session_state["manual_content_ir_cleaned"],
            height=200,
            help="This is the automatically cleaned version"
        )
        if st.button("✅ Use Cleaned Content IR"):
            st.session_state["generated_content_ir"] = st.session_state["manual_content_ir_cleaned"]
            st.session_state.pop("manual_content_ir_cleaned", None)
            st.success("Cleaned Content IR applied!")
            st.rerun()
    
    if st.session_state.get("manual_render_plan_cleaned"):
        st.subheader("🔧 Cleaned Render Plan")
        st.text_area(
            "Cleaned Render Plan JSON",
            value=st.session_state["manual_render_plan_cleaned"],
            height=200,
            help="This is the automatically cleaned version"
        )
        if st.button("✅ Use Cleaned Render Plan"):
            st.session_state["generated_render_plan"] = st.session_state["manual_render_plan_cleaned"]
            st.session_state.pop("manual_render_plan_cleaned", None)
            st.success("Cleaned Render Plan applied!")
            st.rerun()
    
    # Save to session state
    if st.button("💾 Save JSON to Session"):
        st.session_state["generated_content_ir"] = content_ir_str
        st.session_state["generated_render_plan"] = render_plan_str
        st.success("✅ JSON saved to session. Switch to Execute tab to generate your deck.")

with tab_execute:
    st.subheader("⚙️ Generate Pitch Deck")
    
    # Check if files are ready
    files_ready = st.session_state.get("files_ready", False)
    
    if files_ready:
        files_data = st.session_state.get("files_data", {})
        st.success(f"🎉 Using generated files for {files_data.get('company_name', 'your company')}")
        
        # Show file summary
        with st.expander("📋 Generated Files Summary"):
            st.write(f"**Content IR:** {files_data.get('content_ir_filename', 'N/A')}")
            st.write(f"**Render Plan:** {files_data.get('render_plan_filename', 'N/A')}")
            st.write(f"**Timestamp:** {files_data.get('timestamp', 'N/A')}")
    
    # Get JSON from session state
    content_ir_str = st.session_state.get("generated_content_ir", "{}")
    render_plan_str = st.session_state.get("generated_render_plan", "{}")
    
    # Display JSON previews
    col1, col2 = st.columns(2)
    
    with col1:
        st.caption("Content IR Status")
        try:
            # Clean JSON before parsing
            cleaned_content_ir_str = clean_json_string(content_ir_str)
            content_ir = json.loads(cleaned_content_ir_str)
            
            if content_ir:
                st.success(f"✅ Content IR loaded ({len(str(content_ir))} characters)")
                
                # Show if cleaning was applied
                if cleaned_content_ir_str != content_ir_str.strip():
                    st.info("🔧 JSON was automatically cleaned")
                
                # Show brief summary
                summary = {}
                if "entities" in content_ir:
                    summary["Company"] = content_ir.get("entities", {}).get("company", {}).get("name", "N/A")
                if "management_team" in content_ir:
                    summary["Management Profiles"] = len(content_ir.get("management_team", {}).get("left_column_profiles", [])) + len(content_ir.get("management_team", {}).get("right_column_profiles", []))
                
                st.json(summary)
            else:
                st.warning("⚠️ Empty Content IR")
        except json.JSONDecodeError as e:
            st.error(f"⚠️ Invalid Content IR JSON: {e}")
            if st.button("🔧 Try Auto-Clean Content IR", key="clean_content_ir_exec"):
                cleaned = clean_json_string(content_ir_str)
                st.session_state["generated_content_ir"] = cleaned
                st.rerun()
            content_ir = None
        except Exception as e:
            st.error(f"⚠️ Content IR Error: {e}")
            content_ir = None
    
    with col2:
        st.caption("Render Plan Status")
        try:
            # Clean JSON before parsing
            cleaned_render_plan_str = clean_json_string(render_plan_str)
            render_plan = json.loads(cleaned_render_plan_str)
            
            if render_plan and "slides" in render_plan:
                st.success(f"✅ Render Plan loaded ({len(render_plan['slides'])} slides)")
                
                # Show if cleaning was applied
                if cleaned_render_plan_str != render_plan_str.strip():
                    st.info("🔧 JSON was automatically cleaned")
                
                # Show slide types
                slide_types = [slide.get("template", "unknown") for slide in render_plan["slides"]]
                st.write("**Slide Types:**")
                for i, slide_type in enumerate(slide_types[:10], 1):  # Show first 10
                    st.write(f"{i}. {slide_type}")
                if len(slide_types) > 10:
                    st.write(f"... and {len(slide_types) - 10} more slides")
            else:
                st.warning("⚠️ Empty or invalid Render Plan")
        except json.JSONDecodeError as e:
            st.error(f"⚠️ Invalid Render Plan JSON: {e}")
            if st.button("🔧 Try Auto-Clean Render Plan", key="clean_render_plan_exec"):
                cleaned = clean_json_string(render_plan_str)
                st.session_state["generated_render_plan"] = cleaned
                st.rerun()
            render_plan = None
        except Exception as e:
            st.error(f"⚠️ Render Plan Error: {e}")
            render_plan = None
    
    # Pre-execution validation
    if st.button("🔍 Final Validation Before Generation"):
        if not Path(templates_path).exists():
            st.error(f"⚠️ templates.json not found at {templates_path}")
        elif content_ir is None or render_plan is None:
            st.error("⚠️ Please fix the JSON errors above")
        else:
            try:
                # Comprehensive validation
                validation_results = validate_individual_slides(content_ir, render_plan)
                is_valid = display_validation_results(validation_results)
                
                # Traditional catalog validation (if available)
                catalog = TemplateCatalog.from_file(templates_path)
                if HAS_VALIDATORS and not skip_validate:
                    report = validate_render_plan_against_catalog(content_ir, render_plan, catalog)
                    summary = summarize_issues(report)
                    
                    st.write("**📋 Catalog Validation:**")
                    if report.ok:
                        st.success("✅ Catalog validation passed!")
                    else:
                        st.error("⚠️ Catalog validation issues")
                    st.code(summary)
                else:
                    st.info("ℹ️ Catalog validation skipped")
                
                if is_valid:
                    st.success("🎯 **Ready for deck generation!** All validations passed.")
                else:
                    st.error("🚨 **Cannot generate deck** - Fix validation issues first.")
                    
            except Exception as e:
                st.error(f"⚠️ Validation error: {e}")
    
    # Generate deck
    st.markdown("---")
    out_name = st.text_input("Output filename", value="ai_generated_deck.pptx")
    
    if st.button("🎯 Generate Pitch Deck", type="primary", disabled=(not content_ir or not render_plan)):
        if not Path(templates_path).exists():
            st.error(f"⚠️ templates.json not found at {templates_path}")
        elif content_ir is None or render_plan is None:
            st.error("⚠️ Please fix JSON errors first")
        else:
            # Final validation before generation
            validation_results = validate_individual_slides(content_ir, render_plan)
            
            if not validation_results['overall_valid']:
                st.error("🚨 **Cannot generate deck** - Validation failed!")
                display_validation_results(validation_results)
            else:
                try:
                    # Show progress
                    progress_bar = st.progress(0)
                    status_text = st.empty()
                    
                    status_text.text("📄 Rendering slides...")
                    progress_bar.progress(25)
                    
                    # Get brand configuration
                    brand_config = st.session_state.get("brand_config")
                    
                    # Generate deck
                    # Normalize plan to avoid blank cells / missing fields

                    render_plan = normalize_plan(render_plan)

                    prs, saved_path = execute_plan(
                        plan=render_plan,
                        content_ir=content_ir,
                        templates_path=templates_path,
                        output_path=out_name,
                        company_name=company_name,
                        brand_config=brand_config,
                        debug=True,
                    )
                    
                    progress_bar.progress(75)
                    status_text.text("💾 Preparing download...")
                    
                    # Prepare download
                    buf = io.BytesIO()
                    prs.save(buf)
                    buf.seek(0)
                    
                    progress_bar.progress(100)
                    status_text.text("✅ Deck generated successfully!")
                    
                    # Success message
                    st.balloons()
                    st.success(f"🎉 AI-Generated Pitch Deck Complete!")
                    st.info(f"📊 Generated {len(prs.slides)} slides")
                    if brand_config:
                        st.info("🎨 Custom branding applied")
                    st.info(f"💼 Company: {company_name}")
                    st.success("✅ **Zero Empty Boxes Policy** - All slides have complete content!")
                    
                    # Show slide breakdown
                    if render_plan and "slides" in render_plan:
                        slide_types = [slide.get("template", "unknown") for slide in render_plan["slides"]]
                        with st.expander("📋 Slide Details"):
                            for i, slide_type in enumerate(slide_types, 1):
                                st.write(f"{i}. {slide_type}")
                    
                    # Download button
                    st.download_button(
                        "⬇️ Download Your AI-Generated Pitch Deck",
                        data=buf,
                        file_name=out_name,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        type="primary"
                    )
                    
                    progress_bar.empty()
                    status_text.empty()
                    
                except Exception as e:
                    st.error(f"⚠️ Error generating deck: {str(e)}")
                    st.exception(e)

# Footer
st.markdown("---")
st.markdown("""
<div style='text-align: center; color: #666; font-size: 0.8em;'>
    <p>🤖 <strong>AI Deck Builder</strong> - Powered by LLM AI | Investment Banking Pitch Deck Generator</p>
    <p>💡 <em>Start with the AI Copilot → Download JSON Files → Generate Professional Deck</em></p>
    <p>🎨 <em>Enhanced with Zero Empty Boxes Policy & Comprehensive Slide Validation</em></p>
</div>
""", unsafe_allow_html=True)
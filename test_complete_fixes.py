#!/usr/bin/env python3
"""
Comprehensive test to verify all PowerPoint generation fixes:
1. Competitive positioning chart (numeric conversion)
2. Business overview (no duplication)  
3. Buyer profiles (rationale field mapping and fit format)
4. Complete slide generation without corruption
"""

import json
from pptx import Presentation
from pptx.util import Inches
from slide_templates import (
    render_competitive_positioning_slide,
    render_business_overview_slide,
    render_buyer_profiles_slide
)

def test_competitive_positioning():
    """Test competitive positioning slide with numeric conversion"""
    print("\n🎯 Testing Competitive Positioning Slide...")
    
    # Load test data
    with open('fixed_content_ir.json', 'r') as f:
        content_ir = json.load(f)
    
    # Test data with mixed revenue formats
    slide_data = {
        'title': 'Competitive Positioning',
        'competitors': [
            {'name': 'LlamaIndex', 'revenue': 38.0},  # Float
            {'name': 'LangChain', 'revenue': '30'},   # String number
            {'name': 'CrewAI', 'revenue': '5M'},      # String with M
            {'name': 'OpenAI API', 'revenue': 300},   # Int
            {'name': 'Haystack', 'revenue': '8.5'},   # String decimal
        ],
        'assessment': [
            ['Company', 'Market Focus', 'Connectors', 'Enterprise', 'Factuality'],
            ['LlamaIndex', '⭐⭐⭐⭐⭐', '⭐⭐⭐⭐⭐', '⭐⭐⭐⭐', '⭐⭐⭐⭐⭐'],
            ['LangChain', '⭐⭐⭐⭐', '⭐⭐⭐⭐', '⭐⭐⭐⭐', '⭐⭐⭐⭐'],
            ['CrewAI', '⭐⭐⭐', '⭐⭐⭐', '⭐⭐', '⭐⭐⭐'],
        ],
        'barriers': [
            {'title': 'Retrieval Quality', 'desc': 'Requires deep expertise'},
            {'title': 'Connector Ecosystem', 'desc': 'Comprehensive connectors'},
        ],
        'advantages': [
            {'title': 'Data-centric Design', 'desc': 'Excels at retrieval'},
            {'title': 'Grounded Retrieval', 'desc': 'Reduces hallucination'},
        ]
    }
    
    try:
        prs = render_competitive_positioning_slide(
            data=slide_data,
            content_ir=content_ir,
            prs=None
        )
        prs.save('test_competitive_positioning_fixed.pptx')
        print("✅ Competitive positioning slide generated successfully")
        return True
    except Exception as e:
        print(f"❌ Competitive positioning failed: {e}")
        import traceback
        traceback.print_exc()
        return False

def test_business_overview():
    """Test business overview slide"""
    print("\n🏢 Testing Business Overview Slide...")
    
    slide_data = {
        'title': 'Business Overview',
        'description': 'LlamaIndex is a leading AI infrastructure company enabling developers and enterprises to build knowledge assistants.',
        'timeline': {
            'start_year': 2023,
            'end_year': 2025
        },
        'highlights': [
            'Rapidly scaled from open-source project to 3M+ monthly downloads and 300+ enterprise clients.',
            'Enterprise-grade cloud platform (LlamaCloud) and document parser (LlamaParse) powering knowledge management.',
            'Strategic partnerships with Salesforce, KPMG, and major global investors accelerate innovation.'
        ],
        'services': [
            'LlamaIndex Framework for developer-led agentic RAG and knowledge assistant solutions.',
            'LlamaCloud for secure, scalable, production-grade ingestion and retrieval.',
            'LlamaParse for advanced, high-fidelity extraction from PDFs and PowerPoints.'
        ],
        'positioning_desc': 'LlamaIndex is positioned as the enterprise standard for connecting LLMs to private, complex data.'
    }
    
    try:
        prs = render_business_overview_slide(data=slide_data)
        prs.save('test_business_overview_fixed.pptx')
        print("✅ Business overview slide generated successfully")
        return True
    except Exception as e:
        print(f"❌ Business overview failed: {e}")
        import traceback
        traceback.print_exc()
        return False

def test_buyer_profiles():
    """Test buyer profiles with rationale field mapping"""
    print("\n💼 Testing Buyer Profiles Slide...")
    
    # Load actual data
    with open('fixed_content_ir.json', 'r') as f:
        content_ir = json.load(f)
    
    strategic_buyers = content_ir.get('strategic_buyers', [])
    
    slide_data = {
        'title': 'Strategic Buyer Profiles',
        'table_headers': ['Buyer Name', 'Description', 'Investment Rationale', 'Key Synergies', 'Fit'],
        'table_rows': strategic_buyers
    }
    
    try:
        prs = render_buyer_profiles_slide(data=slide_data)
        prs.save('test_buyer_profiles_fixed.pptx')
        print("✅ Buyer profiles slide generated successfully")
        
        # Verify rationale field mapping
        slide = prs.slides[0]
        for shape in slide.shapes:
            if hasattr(shape, 'table'):
                table = shape.table
                if len(table.rows) > 1 and len(table.columns) >= 3:
                    rationale_cell = table.rows[1].cells[2]  # First data row, rationale column
                    cell_text = rationale_cell.text.strip()
                    print(f"📝 Rationale content: '{cell_text}'")
                    
                    if cell_text and cell_text != 'N/A' and len(cell_text) > 10:
                        print("✅ Rationale field properly populated")
                        return True
                    else:
                        print("❌ Rationale field is empty or insufficient")
                        return False
        
        return True
        
    except Exception as e:
        print(f"❌ Buyer profiles failed: {e}")
        import traceback
        traceback.print_exc()
        return False

def test_complete_generation():
    """Test complete PowerPoint generation using actual data"""
    print("\n📋 Testing Complete PowerPoint Generation...")
    
    try:
        # Load both JSON files
        with open('fixed_content_ir.json', 'r') as f:
            content_ir = json.load(f)
        
        with open('fixed_render_plan.json', 'r') as f:
            render_plan = json.load(f)
        
        print(f"📊 Content IR sections: {len(content_ir)}")
        print(f"📊 Render plan slides: {len(render_plan['slides'])}")
        
        # Check for business overview duplication
        templates = [slide['template'] for slide in render_plan['slides']]
        business_overview_count = templates.count('business_overview')
        print(f"📈 Business overview slides: {business_overview_count}")
        
        if business_overview_count > 1:
            print("❌ Business overview duplication detected!")
            return False
        else:
            print("✅ No business overview duplication")
        
        # Check slide order and count
        print(f"📋 Slide templates: {templates}")
        
        # Import and test slide generation functions
        import slide_templates
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        successful_slides = 0
        failed_slides = []
        
        for i, slide_config in enumerate(render_plan['slides']):
            template = slide_config['template']
            data = slide_config.get('data', {})
            
            print(f"\n🎯 Generating slide {i+1}: {template}")
            
            try:
                # Get the appropriate render function
                render_function = getattr(slide_templates, f'render_{template}_slide', None)
                
                if render_function:
                    # Generate slide - use existing presentation
                    prs = render_function(
                        data=data,
                        prs=prs,
                        content_ir=content_ir
                    )
                    successful_slides += 1
                    print(f"✅ Slide {i+1} ({template}) generated successfully")
                else:
                    print(f"❌ Slide {i+1} ({template}) - No render function found")
                    failed_slides.append(template)
                    
            except Exception as e:
                print(f"❌ Slide {i+1} ({template}) failed: {e}")
                failed_slides.append(template)
        
        # Save complete presentation
        output_file = 'test_complete_fixes_output.pptx'
        prs.save(output_file)
        
        total_slides = len(render_plan['slides'])
        print(f"\n📊 Generation Summary:")
        print(f"✅ Successful slides: {successful_slides}/{total_slides}")
        print(f"❌ Failed slides: {len(failed_slides)}")
        if failed_slides:
            print(f"Failed templates: {failed_slides}")
        
        print(f"💾 Complete presentation saved: {output_file}")
        
        return successful_slides > 0 and len(failed_slides) == 0
        
    except Exception as e:
        print(f"❌ Complete generation failed: {e}")
        import traceback
        traceback.print_exc()
        return False

def main():
    """Run all tests"""
    print("🧪 Starting Comprehensive PowerPoint Generation Tests...")
    print("=" * 60)
    
    results = {}
    
    # Test 1: Competitive positioning (numeric conversion fix)
    results['competitive'] = test_competitive_positioning()
    
    # Test 2: Business overview (no duplication)
    results['business'] = test_business_overview()
    
    # Test 3: Buyer profiles (rationale field mapping)
    results['buyers'] = test_buyer_profiles()
    
    # Test 4: Complete generation
    results['complete'] = test_complete_generation()
    
    # Summary
    print("\n" + "=" * 60)
    print("🏆 TEST RESULTS SUMMARY:")
    print("=" * 60)
    
    passed = sum(results.values())
    total = len(results)
    
    for test_name, result in results.items():
        status = "✅ PASS" if result else "❌ FAIL"
        print(f"{test_name.upper():20} | {status}")
    
    print("-" * 40)
    print(f"OVERALL RESULT: {passed}/{total} tests passed")
    
    if passed == total:
        print("🎉 ALL TESTS PASSED! PowerPoint generation is working correctly.")
        return True
    else:
        print("🚨 Some tests failed. Check the error messages above.")
        return False

if __name__ == "__main__":
    success = main()
    exit(0 if success else 1)
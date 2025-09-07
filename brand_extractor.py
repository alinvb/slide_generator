"""
brand_extractor.py - UPDATED WITH FIXED PERPLEXITY API CALLS
Extract branding elements from uploaded PowerPoint files using LLM analysis
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
import io
import json
import requests
from typing import Dict, Optional, Tuple, List
import logging

class LLMBrandExtractor:
    """LLM-powered brand extraction for superior accuracy"""
    
    def __init__(self, api_key: str, model_name: str, api_service: str = "perplexity"):
        self.api_key = api_key
        self.model_name = model_name
        self.api_service = api_service
        self.logger = logging.getLogger(__name__)
    
    def extract_brand_with_llm(self, pptx_file) -> Dict:
        """Extract brand elements using LLM analysis of PowerPoint content"""
        try:
            print("[LLM BRAND DEBUG] Starting LLM-powered brand extraction...")
            
            # First, extract content from slides for LLM analysis
            slide_content = self._extract_slide_content_for_llm(pptx_file)
            
            if not slide_content['color_samples'] and not slide_content['font_samples']:
                print("[LLM BRAND DEBUG] No content extracted, falling back to defaults")
                return self._get_default_brand_config()
            
            # Use LLM to analyze and extract brand elements
            brand_analysis = self._analyze_brand_with_llm(slide_content)
            
            # Convert LLM analysis to brand config
            brand_config = self._convert_llm_analysis_to_config(brand_analysis)
            
            print("[LLM BRAND DEBUG] LLM extraction completed successfully")
            return brand_config
            
        except Exception as e:
            self.logger.error(f"LLM brand extraction failed: {str(e)}")
            print(f"[LLM BRAND DEBUG] LLM extraction failed: {str(e)}")
            return self._get_default_brand_config()
    
    def _extract_slide_content_for_llm(self, pptx_file) -> Dict:
        """Extract comprehensive content from slides for LLM analysis"""
        try:
            if hasattr(pptx_file, 'read'):
                pptx_file.seek(0)
                prs = Presentation(io.BytesIO(pptx_file.read()))
            else:
                prs = Presentation(pptx_file)
            
            slides_info = []
            all_color_samples = []
            all_font_samples = []
            
            print(f"[LLM BRAND DEBUG] Analyzing {len(prs.slides)} slides for content...")
            
            # Analyze first 8 slides (avoid rId error by using enumerate)
            slides_to_analyze = min(8, len(prs.slides))
            for slide_idx in range(slides_to_analyze):
                slide = prs.slides[slide_idx]
                slide_info = {
                    'slide_number': slide_idx + 1,
                    'text_content': [],
                    'colors_found': [],
                    'fonts_found': [],
                    'shape_count': len(slide.shapes)
                }
                
                for shape in slide.shapes:
                    # Extract text content and context
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        text_content = shape.text_frame.text.strip()
                        if text_content and len(text_content) > 2:  # Skip very short text
                            slide_info['text_content'].append({
                                'text': text_content[:100],  # Limit length
                                'length': len(text_content),
                                'context': self._determine_text_context(text_content, shape)
                            })
                            
                        # Extract font information
                        fonts = self._extract_fonts_from_shape(shape)
                        slide_info['fonts_found'].extend(fonts)
                        all_font_samples.extend(fonts)
                    
                    # Extract colors with rich context
                    colors = self._extract_colors_from_shape_with_context(shape, prs)
                    slide_info['colors_found'].extend(colors)
                    all_color_samples.extend(colors)
                
                slides_info.append(slide_info)
            
            # Filter and deduplicate samples
            try:
                unique_colors = self._filter_unique_colors(all_color_samples)
                unique_fonts = self._filter_unique_fonts(all_font_samples)
                
                print(f"[LLM BRAND DEBUG] Extracted {len(unique_colors)} unique colors and {len(unique_fonts)} font samples")
                
                return {
                    'slides': slides_info,
                    'color_samples': unique_colors[:15],  # Limit for token efficiency
                    'font_samples': unique_fonts[:8],
                    'total_slides': len(prs.slides),
                    'slide_dimensions': {
                        'width': getattr(prs.slide_width, 'inches', 13.33) if hasattr(prs, 'slide_width') else 13.33,
                        'height': getattr(prs.slide_height, 'inches', 7.5) if hasattr(prs, 'slide_height') else 7.5
                    }
                }
            except Exception as e:
                print(f"[LLM BRAND DEBUG] Error filtering samples: {e}")
                return {
                    'slides': slides_info,
                    'color_samples': all_color_samples[:10],  # Use unfiltered colors
                    'font_samples': all_font_samples[:5],
                    'total_slides': len(prs.slides),
                    'slide_dimensions': {
                        'width': getattr(prs.slide_width, 'inches', 13.33) if hasattr(prs, 'slide_width') else 13.33,
                        'height': getattr(prs.slide_height, 'inches', 7.5) if hasattr(prs, 'slide_height') else 7.5
                    }
                }
            
        except Exception as e:
            print(f"[LLM BRAND DEBUG] Content extraction failed: {str(e)}")
            return {'slides': [], 'color_samples': [], 'font_samples': [], 'total_slides': 0}
    
    def _determine_text_context(self, text: str, shape) -> str:
        """Determine the context/role of text content"""
        text_len = len(text)
        
        # Very short text likely titles or headers
        if text_len < 30:
            return 'title_or_header'
        # Medium text likely subheaders
        elif text_len < 100:
            return 'subheader'
        # Long text likely body content
        else:
            return 'body_content'
    
    def _extract_fonts_from_shape(self, shape) -> List[Dict]:
        """Extract font information from a shape"""
        fonts = []
        
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run.font, 'name') and run.font.name:
                            font_size = None
                            if run.font.size:
                                try:
                                    font_size = int(run.font.size.pt)
                                except:
                                    pass
                            
                            font_info = {
                                'name': run.font.name,
                                'size': font_size,
                                'bold': run.font.bold,
                                'italic': run.font.italic,
                                'text_length': len(run.text) if run.text else 0,
                                'context': self._determine_font_context(font_size, run.font.bold)
                            }
                            fonts.append(font_info)
        except Exception as e:
            print(f"[LLM BRAND DEBUG] Font extraction error: {e}")
        
        return fonts
    
    def _determine_font_context(self, size: Optional[int], bold: Optional[bool]) -> str:
        """Determine font usage context"""
        if size and size >= 20:
            return 'title'
        elif size and size >= 14:
            return 'header'
        elif bold:
            return 'emphasis'
        else:
            return 'body'
    
    def _extract_colors_from_shape_with_context(self, shape, prs) -> List[Dict]:
        """Extract colors with rich contextual information"""
        colors = []
        
        try:
            # Fill colors - only if shape has a solid fill
            if (hasattr(shape, 'fill') and 
                hasattr(shape.fill, 'type') and 
                hasattr(shape.fill, 'fore_color') and
                shape.fill.type == 1):  # SOLID fill type
                
                if hasattr(shape.fill.fore_color, 'rgb'):
                    rgb = shape.fill.fore_color.rgb
                    if not self._is_ignore_color(rgb):
                        # Handle RGBColor object (can be unpacked as tuple)
                        if isinstance(rgb, str):
                            # RGBColor is actually a hex string like "0070C0"
                            r = int(rgb[0:2], 16)
                            g = int(rgb[2:4], 16)
                            b = int(rgb[4:6], 16)
                        elif hasattr(rgb, 'r'):
                            r, g, b = rgb.r, rgb.g, rgb.b
                        else:
                            # Handle RGBColor object (can be unpacked as tuple)
                            try:
                                r, g, b = rgb
                            except (ValueError, TypeError):
                                return colors  # Return current colors if we can't process this one
                        
                        colors.append({
                            'type': 'fill',
                            'rgb': [r, g, b],
                            'hex': f"#{r:02x}{g:02x}{b:02x}",
                            'context': self._determine_color_context(shape, prs),
                            'shape_size': self._get_shape_size_category(shape, prs)
                        })
            
            # Text colors
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb'):
                            rgb = run.font.color.rgb
                            if not self._is_ignore_color(rgb):
                                # Handle RGBColor object (can be unpacked as tuple)
                                if isinstance(rgb, str):
                                    # RGBColor is actually a hex string like "0070C0"
                                    r = int(rgb[0:2], 16)
                                    g = int(rgb[2:4], 16)
                                    b = int(rgb[4:6], 16)
                                elif hasattr(rgb, 'r'):
                                    r, g, b = rgb.r, rgb.g, rgb.b
                                else:
                                    # Handle RGBColor object (can be unpacked as tuple)
                                    try:
                                        r, g, b = rgb
                                    except (ValueError, TypeError):
                                        continue  # Skip this run if we can't process the color
                                
                                colors.append({
                                    'type': 'text',
                                    'rgb': [r, g, b],
                                    'hex': f"#{r:02x}{g:02x}{b:02x}",
                                    'context': 'text_content',
                                    'font_size': getattr(run.font.size, 'pt', None) if run.font.size else None,
                                    'is_bold': run.font.bold
                                })
            
            # Line colors
            if hasattr(shape, 'line') and hasattr(shape.line, 'color'):
                if hasattr(shape.line.color, 'rgb'):
                    rgb = shape.line.color.rgb
                    if not self._is_ignore_color(rgb):
                        # Handle RGBColor object (can be unpacked as tuple)
                        if isinstance(rgb, str):
                            # RGBColor is actually a hex string like "0070C0"
                            r = int(rgb[0:2], 16)
                            g = int(rgb[2:4], 16)
                            b = int(rgb[4:6], 16)
                        elif hasattr(rgb, 'r'):
                            r, g, b = rgb.r, rgb.g, rgb.b
                        else:
                            # Handle RGBColor object (can be unpacked as tuple)
                            try:
                                r, g, b = rgb
                            except (ValueError, TypeError):
                                return colors  # Return current colors if we can't process this one
                        
                        colors.append({
                            'type': 'line',
                            'rgb': [r, g, b],
                            'hex': f"#{r:02x}{g:02x}{b:02x}",
                            'context': 'border_or_accent'
                        })
                        
        except Exception as e:
            print(f"[LLM BRAND DEBUG] Color extraction error: {e}")
        
        return colors
    
    def _determine_color_context(self, shape, prs) -> str:
        """Determine color usage context based on shape properties"""
        try:
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                # Large shapes spanning most width = headers
                if hasattr(prs, 'slide_width') and shape.width > prs.slide_width * 0.6:
                    return 'header_or_banner'
                
                # Small shapes = accents or logos
                elif (shape.width < prs.slide_width * 0.1 and 
                      shape.height < prs.slide_height * 0.1):
                    return 'accent_or_logo'
                
                # Medium shapes = content or emphasis
                else:
                    return 'content_emphasis'
        except:
            pass
        
        return 'general_element'
    
    def _get_shape_size_category(self, shape, prs) -> str:
        """Categorize shape by size"""
        try:
            if hasattr(shape, 'width') and hasattr(prs, 'slide_width'):
                width_ratio = shape.width / prs.slide_width
                if width_ratio > 0.6:
                    return 'large'
                elif width_ratio > 0.2:
                    return 'medium'
                else:
                    return 'small'
        except:
            pass
        return 'unknown'
    
    def _is_ignore_color(self, rgb) -> bool:
        """Check if color should be ignored"""
        if not rgb:
            return True
        
        # Handle different RGB color formats
        if isinstance(rgb, str):
            # RGBColor is actually a hex string like "0070C0"
            try:
                if len(rgb) == 6:
                    r = int(rgb[0:2], 16)
                    g = int(rgb[2:4], 16)
                    b = int(rgb[4:6], 16)
                else:
                    return True
            except:
                return True
        elif hasattr(rgb, 'r'):
            r, g, b = rgb.r, rgb.g, rgb.b
        else:
            # Handle RGBColor object (can be unpacked as tuple)
            try:
                r, g, b = rgb
            except (ValueError, TypeError):
                return True
        
        # Ignore pure white, black, and very light colors
        return ((r == 255 and g == 255 and b == 255) or  # White
                (r == 0 and g == 0 and b == 0) or          # Black
                (r > 250 and g > 250 and b > 250))         # Very light
    
    def _filter_unique_colors(self, colors: List[Dict]) -> List[Dict]:
        """Remove duplicate colors and keep most contextually relevant ones"""
        unique_colors = []
        seen_hex = set()
        
        # Sort by context priority
        context_priority = {
            'header_or_banner': 1,
            'content_emphasis': 2,
            'accent_or_logo': 3,
            'text_content': 4,
            'border_or_accent': 5,
            'general_element': 6
        }
        
        try:
            sorted_colors = sorted(colors, key=lambda x: context_priority.get(x.get('context', 'general_element'), 6))
            
            for color in sorted_colors:
                if isinstance(color, dict) and 'hex' in color:
                    if color['hex'] not in seen_hex:
                        unique_colors.append(color)
                        seen_hex.add(color['hex'])
                else:
                    print(f"[LLM BRAND DEBUG] Skipping invalid color format: {color}")
        except Exception as e:
            print(f"[LLM BRAND DEBUG] Error filtering colors: {e}")
            # Fallback: just return the first few colors
            unique_colors = colors[:10] if colors else []
        
        return unique_colors
    
    def _filter_unique_fonts(self, fonts: List[Dict]) -> List[Dict]:
        """Remove duplicate fonts and keep most relevant ones"""
        unique_fonts = []
        seen_names = set()
        
        # Sort by context priority
        context_priority = {'title': 1, 'header': 2, 'emphasis': 3, 'body': 4}
        sorted_fonts = sorted(fonts, key=lambda x: context_priority.get(x.get('context', 'body'), 4))
        
        for font in sorted_fonts:
            if font['name'] not in seen_names:
                unique_fonts.append(font)
                seen_names.add(font['name'])
                if len(unique_fonts) >= 5:  # Limit to top 5 fonts
                    break
        
        return unique_fonts
    
    def _analyze_brand_with_llm(self, slide_content: Dict) -> str:
        """Use LLM to analyze slide content and extract brand elements"""
        
        analysis_prompt = f"""You are an expert brand designer analyzing a PowerPoint presentation to extract brand elements.

PRESENTATION ANALYSIS:
- Total slides: {slide_content['total_slides']}
- Analyzed slides: {len(slide_content['slides'])}

COLOR SAMPLES FOUND ({len(slide_content['color_samples'])} unique colors):
{json.dumps(slide_content['color_samples'], indent=2)}

FONT SAMPLES FOUND ({len(slide_content['font_samples'])} fonts):
{json.dumps(slide_content['font_samples'], indent=2)}

SLIDE CONTENT SUMMARY:
{json.dumps([{{'slide': s['slide_number'], 'shapes': s['shape_count'], 'text_snippets': len(s['text_content'])}} for s in slide_content['slides'][:5]], indent=2)}

TASK: Analyze this data and extract the brand elements. Return ONLY valid JSON in this exact format:

{{
  "primary_color": {{"hex": "#RRGGBB", "rgb": [R, G, B], "reasoning": "why this is the primary brand color"}},
  "secondary_color": {{"hex": "#RRGGBB", "rgb": [R, G, B], "reasoning": "why this is secondary"}},
  "accent_color": {{"hex": "#RRGGBB", "rgb": [R, G, B], "reasoning": "why this is accent"}},
  "text_color": {{"hex": "#RRGGBB", "rgb": [R, G, B], "reasoning": "main text color rationale"}},
  "background_color": {{"hex": "#FFFFFF", "rgb": [255, 255, 255], "reasoning": "background choice"}},
  "primary_font": {{"name": "Font Name", "reasoning": "why this is the main brand font"}},
  "title_font_size": {{"size": 24, "reasoning": "observed title size pattern"}},
  "body_font_size": {{"size": 11, "reasoning": "observed body text size pattern"}},
  "brand_personality": {{"description": "professional/modern/creative/corporate/etc with brief explanation"}},
  "design_patterns": {{"description": "key design elements observed"}}
}}

ANALYSIS GUIDELINES:
1. PRIMARY COLOR: Choose the most prominent color used in headers, banners, or large elements
2. SECONDARY COLOR: Choose a complementary color used for emphasis or accents  
3. ACCENT COLOR: Pick a color used for small details, borders, or highlights
4. TEXT COLOR: Most common text color (usually dark, avoid pure black unless intentional)
5. FONTS: Prioritize fonts used in titles/headers over body text
6. AVOID: Pure white (#FFFFFF), pure black (#000000) unless clearly brand colors
7. CONTEXT MATTERS: Colors in 'header_or_banner' context are more important than 'general_element'
8. CONSISTENCY: Look for colors that appear multiple times across slides

Return ONLY the JSON, no additional text."""

        try:
            print("[LLM BRAND DEBUG] Sending analysis request to LLM...")
            
            messages = [
                {"role": "system", "content": "You are a brand design expert. Return only valid JSON with no additional formatting or text."},
                {"role": "user", "content": analysis_prompt}
            ]
            
            response = self._call_llm_api(messages)
            cleaned_response = self._extract_json_from_response(response)
            
            print(f"[LLM BRAND DEBUG] LLM analysis completed. Response length: {len(cleaned_response)}")
            
            return cleaned_response
            
        except Exception as e:
            print(f"[LLM BRAND DEBUG] LLM analysis failed: {str(e)}")
            return "{}"
    
    def _call_llm_api(self, messages: List[Dict]) -> str:
        """Call the appropriate LLM API based on service"""
        try:
            if self.api_service == "perplexity":
                return self._call_perplexity_api(messages)
            elif self.api_service == "claude":
                return self._call_claude_api(messages)
            else:
                print(f"[LLM BRAND DEBUG] Unknown API service: {self.api_service}")
                return "{}"
        except Exception as e:
            print(f"[LLM BRAND DEBUG] API call failed: {str(e)}")
            return "{}"
    
    def _call_perplexity_api(self, messages: List[Dict]) -> str:
        """Call Perplexity API - FIXED for message alternation"""
        try:
            url = "https://api.perplexity.ai/chat/completions"
            
            # FIX: Ensure proper message alternation for Perplexity
            cleaned_messages = self._build_alternating_messages(messages)
            
            payload = {
                "model": self.model_name,
                "messages": cleaned_messages,
                "temperature": 0.3,  # Lower temperature for more consistent results
                "max_tokens": 4000,  # Increased from 1500 for complete brand analysis
                "stream": False
            }
            
            headers = {
                "Authorization": f"Bearer {self.api_key}",
                "Content-Type": "application/json"
            }
            
            response = requests.post(url, json=payload, headers=headers)
            
            if response.status_code == 200:
                result = response.json()
                return result.get('choices', [{}])[0].get('message', {}).get('content', '{}')
            else:
                print(f"[LLM BRAND DEBUG] Perplexity API Error: {response.status_code} - {response.text}")
                return "{}"
                
        except Exception as e:
            print(f"[LLM BRAND DEBUG] Perplexity API call failed: {str(e)}")
            return "{}"
    
    def _build_alternating_messages(self, messages: List[Dict]) -> List[Dict]:
        """Build properly alternating messages for Perplexity API"""
        # Extract system message
        system_message = None
        conversation_messages = []
        
        for msg in messages:
            if msg["role"] == "system":
                system_message = msg["content"]
            elif msg["role"] in ["user", "assistant"]:
                conversation_messages.append(msg)
        
        # Remove any leading assistant messages (Perplexity needs user first after system)
        while conversation_messages and conversation_messages[0]["role"] == "assistant":
            conversation_messages.pop(0)
        
        # Collapse consecutive same-role messages to enforce alternation
        cleaned_messages = []
        for msg in conversation_messages:
            if cleaned_messages and cleaned_messages[-1]["role"] == msg["role"]:
                # Combine consecutive messages of same role
                cleaned_messages[-1]["content"] = cleaned_messages[-1]["content"].rstrip() + "\n\n" + str(msg.get("content", "")).strip()
            else:
                cleaned_messages.append({
                    "role": msg["role"], 
                    "content": str(msg.get("content", "")).strip()
                })
        
        # Build final message array for Perplexity
        final_messages = []
        
        # Add system message if present
        if system_message:
            final_messages.append({"role": "system", "content": system_message})
        
        # Add alternating conversation
        final_messages.extend(cleaned_messages)
        
        # Ensure we don't have empty messages
        final_messages = [msg for msg in final_messages if msg.get("content", "").strip()]
        
        return final_messages
    
    def _call_claude_api(self, messages: List[Dict]) -> str:
        """Call Claude API"""
        try:
            url = "https://api.anthropic.com/v1/messages"
            
            # Convert messages format for Claude
            claude_messages = []
            system_message = ""
            
            for msg in messages:
                if msg["role"] == "system":
                    system_message = msg["content"]
                else:
                    claude_messages.append({
                        "role": msg["role"],
                        "content": msg["content"]
                    })
            
            payload = {
                "model": self.model_name,
                "max_tokens": 4000,  # Increased from 1500 for complete brand analysis
                "temperature": 0.3,
                "messages": claude_messages
            }
            
            if system_message:
                payload["system"] = system_message
            
            headers = {
                "x-api-key": self.api_key,
                "Content-Type": "application/json",
                "anthropic-version": "2023-06-01"
            }
            
            response = requests.post(url, json=payload, headers=headers)
            
            if response.status_code == 200:
                result = response.json()
                return result.get('content', [{}])[0].get('text', '{}')
            else:
                print(f"[LLM BRAND DEBUG] Claude API Error: {response.status_code} - {response.text}")
                return "{}"
                
        except Exception as e:
            print(f"[LLM BRAND DEBUG] Claude API call failed: {str(e)}")
            return "{}"
    
    def _extract_json_from_response(self, response: str) -> str:
        """Extract and validate JSON from LLM response"""
        try:
            # Find JSON boundaries
            start_idx = response.find('{')
            end_idx = response.rfind('}')
            
            if start_idx != -1 and end_idx != -1 and end_idx > start_idx:
                json_str = response[start_idx:end_idx+1]
                
                # Validate JSON
                parsed = json.loads(json_str)
                return json_str
            else:
                print("[LLM BRAND DEBUG] No valid JSON found in response")
                return "{}"
                
        except json.JSONDecodeError as e:
            print(f"[LLM BRAND DEBUG] JSON parsing failed: {e}")
            return "{}"
    
    def _convert_llm_analysis_to_config(self, analysis_json: str) -> Dict:
        """Convert LLM analysis to brand configuration format"""
        try:
            analysis = json.loads(analysis_json)
            
            def rgb_from_analysis(color_data, default_rgb=[64, 64, 64]):
                if isinstance(color_data, dict) and 'rgb' in color_data:
                    rgb = color_data['rgb']
                    if isinstance(rgb, list) and len(rgb) == 3:
                        return RGBColor(max(0, min(255, rgb[0])), 
                                      max(0, min(255, rgb[1])), 
                                      max(0, min(255, rgb[2])))
                return RGBColor(default_rgb[0], default_rgb[1], default_rgb[2])
            
            brand_config = {
                'color_scheme': {
                    'primary': rgb_from_analysis(analysis.get('primary_color'), [24, 58, 88]),
                    'secondary': rgb_from_analysis(analysis.get('secondary_color'), [181, 151, 91]),
                    'accent': rgb_from_analysis(analysis.get('accent_color'), [64, 64, 64]),
                    'background': rgb_from_analysis(analysis.get('background_color'), [255, 255, 255]),
                    'text': rgb_from_analysis(analysis.get('text_color'), [64, 64, 64]),
                    'light_grey': RGBColor(240, 240, 240),
                    'footer_grey': RGBColor(128, 128, 128)
                },
                'typography': {
                    'primary_font': analysis.get('primary_font', {}).get('name', 'Arial'),
                    'title_size': analysis.get('title_font_size', {}).get('size', 24),
                    'header_size': 14,
                    'body_size': analysis.get('body_font_size', {}).get('size', 11),
                    'small_size': 9
                },
                'header_style': {
                    'type': 'line',
                    'height': 0.05,
                    'color': 'primary',
                    'has_logo': False
                },
                'layout_config': {
                    'title_alignment': 'left',
                    'header_type': 'extracted'
                },
                'llm_analysis': {
                    'brand_personality': analysis.get('brand_personality', {}),
                    'design_patterns': analysis.get('design_patterns', {}),
                    'color_reasoning': {
                        'primary': analysis.get('primary_color', {}).get('reasoning', ''),
                        'secondary': analysis.get('secondary_color', {}).get('reasoning', ''),
                        'accent': analysis.get('accent_color', {}).get('reasoning', '')
                    },
                    'font_reasoning': analysis.get('primary_font', {}).get('reasoning', '')
                }
            }
            
            # Log extracted colors
            print("[LLM BRAND DEBUG] Final extracted brand colors:")
            for name, color in brand_config['color_scheme'].items():
                if hasattr(color, 'r'):
                    print(f"  {name}: RGB({color.r}, {color.g}, {color.b}) = #{color.r:02x}{color.g:02x}{color.b:02x}")
            
            print(f"[LLM BRAND DEBUG] Primary font: {brand_config['typography']['primary_font']}")
            
            return brand_config
            
        except Exception as e:
            print(f"[LLM BRAND DEBUG] Config conversion failed: {str(e)}")
            return self._get_default_brand_config()
    
    def _get_default_brand_config(self) -> Dict:
        """Default brand configuration"""
        return {
            'color_scheme': {
                'primary': RGBColor(24, 58, 88),
                'secondary': RGBColor(181, 151, 91),
                'accent': RGBColor(64, 64, 64),
                'background': RGBColor(255, 255, 255),
                'text': RGBColor(64, 64, 64),
                'light_grey': RGBColor(240, 240, 240),
                'footer_grey': RGBColor(128, 128, 128)
            },
            'typography': {
                'primary_font': 'Arial',
                'title_size': 24,
                'header_size': 14,
                'body_size': 11,
                'small_size': 9
            },
            'header_style': {
                'type': 'line',
                'height': 0.05,
                'color': 'primary',
                'has_logo': False
            },
            'layout_config': {
                'title_alignment': 'left',
                'header_type': 'line'
            }
        }


class BrandExtractor:
    """Enhanced brand extractor with both rule-based and LLM options"""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    def extract_brand_from_pptx(self, pptx_file, use_llm=False, api_key=None, model_name=None, api_service="perplexity") -> Dict:
        """
        Extract brand elements from uploaded PowerPoint file
        
        Args:
            pptx_file: Uploaded file object or file path
            use_llm: Whether to use LLM-powered extraction
            api_key: API key for LLM service
            model_name: LLM model name
            api_service: LLM service ("perplexity" or "claude")
            
        Returns:
            Dict containing color_scheme, typography, and header_style
        """
        if use_llm and api_key and model_name:
            print(f"[BRAND DEBUG] Using LLM-powered extraction with {api_service}")
            llm_extractor = LLMBrandExtractor(api_key, model_name, api_service)
            return llm_extractor.extract_brand_with_llm(pptx_file)
        else:
            print("[BRAND DEBUG] Using rule-based extraction")
            return self._extract_with_rules(pptx_file)
    
    def _extract_with_rules(self, pptx_file) -> Dict:
        """Original rule-based extraction method"""
        try:
            # Handle different input types
            if hasattr(pptx_file, 'read'):
                # File-like object from Streamlit - need to reset position
                pptx_file.seek(0)
                prs = Presentation(io.BytesIO(pptx_file.read()))
            else:
                # File path
                prs = Presentation(pptx_file)
            
            brand_config = {
                'color_scheme': self._extract_colors(prs),
                'typography': self._extract_fonts(prs),
                'header_style': self._extract_header_style(prs),
                'layout_config': {
                    'title_alignment': 'left',
                    'header_type': 'extracted'
                }
            }
            
            self.logger.info("Successfully extracted brand elements from PowerPoint")
            print(f"[BRAND DEBUG] Extracted colors:")
            for k, v in brand_config['color_scheme'].items():
                hex_color = self._rgb_to_hex(v)
                print(f"  {k}: {hex_color}")
            return brand_config
            
        except Exception as e:
            self.logger.error(f"Error extracting brand from PowerPoint: {str(e)}")
            print(f"[BRAND DEBUG] Extraction failed: {str(e)}")
            return self._get_default_brand_config()
    
    # [Include all the existing rule-based methods here for compatibility]
    def _rgb_to_hex(self, rgb_color):
        """Convert RGBColor to hex string for debugging"""
        if hasattr(rgb_color, 'r'):
            return f"#{rgb_color.r:02x}{rgb_color.g:02x}{rgb_color.b:02x}"
        elif isinstance(rgb_color, str):
            # RGBColor is actually a hex string like "0070C0"
            return f"#{rgb_color}"
        else:
            # Handle RGBColor object (can be unpacked as tuple)
            try:
                r, g, b = rgb_color
                return f"#{r:02x}{g:02x}{b:02x}"
            except (ValueError, TypeError):
                return str(rgb_color)
    
    def _extract_colors(self, prs: Presentation) -> Dict:
        """Extract color scheme from the presentation - IMPROVED VERSION"""
        colors = {
            'primary': RGBColor(24, 58, 88),
            'secondary': RGBColor(181, 151, 91),
            'accent': RGBColor(64, 64, 64),
            'background': RGBColor(255, 255, 255),
            'text': RGBColor(64, 64, 64),
            'light_grey': RGBColor(240, 240, 240),
            'footer_grey': RGBColor(128, 128, 128)
        }
        
        print(f"[BRAND DEBUG] Starting color extraction from {len(prs.slides)} slides...")
        
        try:
            extracted_colors = self._extract_colors_from_slides(prs)
            if extracted_colors:
                colors.update(extracted_colors)
                print(f"[BRAND DEBUG] Successfully extracted {len(extracted_colors)} colors from slides")
                return colors
            
            print("[BRAND DEBUG] No custom colors found, using defaults")
        except Exception as e:
            self.logger.warning(f"Could not extract colors, using defaults: {str(e)}")
            print(f"[BRAND DEBUG] Color extraction error: {str(e)}")
        
        return colors
    
    def _extract_colors_from_slides(self, prs: Presentation) -> Optional[Dict]:
        """Extract colors from actual slide elements"""
        colors = {}
        color_candidates = {'primary': [], 'secondary': [], 'accent': [], 'text': []}
        
        print(f"[BRAND DEBUG] Starting color extraction from {len(prs.slides)} slides...")
        
        try:
            for slide_idx, slide in enumerate(prs.slides):
                print(f"[BRAND DEBUG] Processing slide {slide_idx + 1} with {len(slide.shapes)} shapes")
                
                for shape_idx, shape in enumerate(slide.shapes):
                    print(f"[BRAND DEBUG] Processing shape {shape_idx}: {type(shape).__name__}")
                    
                    try:
                        if hasattr(shape, 'fill'):
                            print(f"[BRAND DEBUG] Shape has fill attribute")
                            fill_color = self._get_fill_color(shape.fill)
                            print(f"[BRAND DEBUG] Extracted fill color: {fill_color}")
                            
                            if fill_color and not self._is_default_color(fill_color):
                                print(f"[BRAND DEBUG] Fill color passed default check: {fill_color}")
                                color_role = self._determine_color_role(shape, prs)
                                color_candidates[color_role].append(fill_color)
                                print(f"[BRAND DEBUG] Added fill color {fill_color} to {color_role}")
                            else:
                                print(f"[BRAND DEBUG] Fill color failed default check or is None")
                        
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            print(f"[BRAND DEBUG] Shape has text frame")
                            text_colors = self._extract_text_colors(shape.text_frame)
                            print(f"[BRAND DEBUG] Extracted text colors: {text_colors}")
                            
                            for role, color in text_colors.items():
                                if color and not self._is_default_color(color):
                                    print(f"[BRAND DEBUG] Text color passed default check: {color}")
                                    color_candidates[role].append(color)
                                    print(f"[BRAND DEBUG] Added text color {color} to {role}")
                                else:
                                    print(f"[BRAND DEBUG] Text color failed default check or is None")
                                
                    except Exception as e:
                        print(f"[BRAND DEBUG] Error processing shape {shape_idx}: {str(e)}")
                        continue
                        
                # Limit to first 10 slides to avoid performance issues
                if slide_idx >= 9:
                    break
        except Exception as e:
            print(f"[BRAND DEBUG] Error in color extraction: {str(e)}")
            return None
        
        print(f"[BRAND DEBUG] Color candidates: {color_candidates}")
        
        for role, candidate_colors in color_candidates.items():
            if candidate_colors:
                # Convert extracted color to RGBColor object
                extracted_color = candidate_colors[0]
                if isinstance(extracted_color, tuple) and len(extracted_color) == 3:
                    colors[role] = RGBColor(extracted_color[0], extracted_color[1], extracted_color[2])
                elif isinstance(extracted_color, str) and len(extracted_color) == 6:
                    # Hex string like "FF0000"
                    r = int(extracted_color[0:2], 16)
                    g = int(extracted_color[2:4], 16) 
                    b = int(extracted_color[4:6], 16)
                    colors[role] = RGBColor(r, g, b)
                else:
                    colors[role] = extracted_color  # Keep as-is if already RGBColor
                print(f"[BRAND DEBUG] Selected {role} color: {colors[role]}")
        
        print(f"[BRAND DEBUG] Final extracted colors: {colors}")
        return colors if colors else None
    
    def _get_fill_color(self, fill) -> Optional[str]:
        """Extract fill color from a shape fill - FIXED VERSION"""
        try:
            if not fill:
                return None
                
            print(f"[BRAND DEBUG] Fill type: {getattr(fill, 'type', 'unknown')}")
            
            # Check for solid fill (type 1)
            if hasattr(fill, 'type') and fill.type == 1:
                if hasattr(fill, 'fore_color') and fill.fore_color:
                    if hasattr(fill.fore_color, 'rgb') and fill.fore_color.rgb:
                        print(f"[BRAND DEBUG] Found solid fill color: {fill.fore_color.rgb}")
                        return fill.fore_color.rgb
                    elif hasattr(fill.fore_color, 'theme_color'):
                        # Convert theme color to RGB
                        theme_color = fill.fore_color.theme_color
                        print(f"[BRAND DEBUG] Found theme color: {theme_color}")
                        return self._theme_color_to_hex(theme_color)
            
            # Generic check for any color
            if hasattr(fill, 'fore_color') and fill.fore_color:
                if hasattr(fill.fore_color, 'rgb') and fill.fore_color.rgb:
                    print(f"[BRAND DEBUG] Found generic fore color: {fill.fore_color.rgb}")
                    return fill.fore_color.rgb
                    
        except Exception as e:
            print(f"[BRAND DEBUG] Error extracting fill color: {str(e)}")
        
        return None
    
    def _theme_color_to_hex(self, theme_color):
        """Convert theme color index to hex string"""
        theme_colors = {
            1: "FF0000",      # Red
            2: "00FF00",      # Green
            3: "0000FF",      # Blue
            4: "FFFF00",      # Yellow
            5: "FF00FF",      # Magenta
            6: "00FFFF",       # Cyan
            7: "800000",       # Dark Red
            8: "008000",       # Dark Green
            9: "000080",       # Dark Blue
            10: "808000",      # Dark Yellow
            11: "800080",      # Dark Magenta
            12: "008080",      # Dark Cyan
            13: "000000",      # Black
            14: "808080",      # Gray
            15: "FFFFFF",      # White
        }
        return theme_colors.get(theme_color, "000000")
    
    def _extract_text_colors(self, text_frame) -> Dict[str, RGBColor]:
        """Extract colors from text in a text frame"""
        text_colors = {}
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    try:
                        if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb'):
                            color = run.font.color.rgb
                            if color and not self._is_default_color(color):
                                if run.font.size and hasattr(run.font.size, 'pt'):
                                    font_size_pt = run.font.size.pt
                                    if font_size_pt >= 18:
                                        text_colors['primary'] = color
                                    elif font_size_pt >= 12:
                                        text_colors['text'] = color
                                    else:
                                        text_colors['accent'] = color
                                else:
                                    text_colors['text'] = color
                    except Exception:
                        continue
        except Exception:
            pass
        return text_colors
    
    def _determine_color_role(self, shape, prs) -> str:
        """Determine the role of a color based on shape properties"""
        try:
            if hasattr(shape, 'width') and hasattr(prs, 'slide_width'):
                if shape.width > prs.slide_width * 0.6:
                    return 'primary'
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                if shape.width < prs.slide_width * 0.1 and shape.height < prs.slide_height * 0.1:
                    return 'secondary'
            return 'accent'
        except Exception:
            return 'primary'
    
    def _is_default_color(self, color) -> bool:
        """Check if a color is a default/common color to ignore - ENHANCED"""
        if not color:
            return True
        
        # Handle string-based RGBColor
        if isinstance(color, str):
            # Convert hex string to RGB values
            try:
                if len(color) == 6:
                    r = int(color[0:2], 16)
                    g = int(color[2:4], 16)
                    b = int(color[4:6], 16)
                else:
                    return True
            except:
                return True
        else:
            # Handle RGBColor object (can be unpacked as tuple)
            try:
                r, g, b = color
            except (ValueError, TypeError):
                return True
        
        # More permissive filtering - only filter out pure white, pure black, and very light grays
        return ((r == 255 and g == 255 and b == 255) or  # Pure white
                (r == 0 and g == 0 and b == 0) or        # Pure black
                (r > 250 and g > 250 and b > 250))       # Very light gray/white
    
    def _extract_fonts(self, prs: Presentation) -> Dict:
        """Extract typography settings from the presentation"""
        fonts = {
            'primary_font': 'Arial',
            'title_size': 24,
            'header_size': 14,
            'body_size': 11,
            'small_size': 9
        }
        
        try:
            for slide in prs.slides[:5]:
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.font.name:
                                        font_size_pts = None
                                        if run.font.size:
                                            try:
                                                font_size_pts = int(run.font.size.pt)
                                            except:
                                                pass
                                        
                                        if font_size_pts and font_size_pts >= 20:
                                            fonts['primary_font'] = run.font.name
                                            fonts['title_size'] = font_size_pts
                                            break
                                        elif font_size_pts and font_size_pts >= 14:
                                            fonts['header_size'] = font_size_pts
                                        elif font_size_pts and font_size_pts >= 10:
                                            fonts['body_size'] = font_size_pts
                    except Exception:
                        continue
        except Exception as e:
            self.logger.warning(f"Could not extract fonts, using defaults: {str(e)}")
        
        return fonts
    
    def _extract_header_style(self, prs: Presentation) -> Dict:
        """Extract header/title styling from the presentation"""
        return {
            'type': 'line',
            'height': 0.05,
            'color': 'primary',
            'has_logo': False
        }
    
    def _get_default_brand_config(self) -> Dict:
        """Return default brand configuration when extraction fails"""
        return {
            'color_scheme': {
                'primary': RGBColor(24, 58, 88),
                'secondary': RGBColor(181, 151, 91),
                'accent': RGBColor(64, 64, 64),
                'background': RGBColor(255, 255, 255),
                'text': RGBColor(64, 64, 64),
                'light_grey': RGBColor(240, 240, 240),
                'footer_grey': RGBColor(128, 128, 128)
            },
            'typography': {
                'primary_font': 'Arial',
                'title_size': 24,
                'header_size': 14,
                'body_size': 11,
                'small_size': 9
            },
            'header_style': {
                'type': 'line',
                'height': 0.05,
                'color': 'primary',
                'has_logo': False
            },
            'layout_config': {
                'title_alignment': 'left',
                'header_type': 'line'
            }
        }
"""
Professional slide_templates.py with full-featured renderers
These create actual charts, tables, and sophisticated layouts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from datetime import datetime

# Removed circular import - _apply_standard_header_and_title is now defined locally


def get_brand_styling(brand_config=None, color_scheme=None, typography=None):
    """Extract brand styling or use defaults - FIXED VERSION"""
    print(f"[DEBUG] get_brand_styling called with brand_config: {brand_config is not None}")
    
    if brand_config:
        print(f"[DEBUG] Brand config keys: {list(brand_config.keys())}")
        brand_colors = brand_config.get('color_scheme', {})
        brand_fonts = brand_config.get('typography', {})
        
        print(f"[DEBUG] Brand colors keys: {list(brand_colors.keys())}")
        print(f"[DEBUG] Brand fonts keys: {list(brand_fonts.keys())}")
        
        # Handle different color formats - IMPROVED CONVERSION
        colors = {}
        color_defaults = {
            "primary": RGBColor(24, 58, 88),
            "secondary": RGBColor(181, 151, 91),
            "accent": RGBColor(64, 64, 64),
            "text": RGBColor(64, 64, 64),
            "background": RGBColor(255, 255, 255),
            "light_grey": RGBColor(240, 240, 240),
            "footer_grey": RGBColor(128, 128, 128)
        }
        
        for key, default_color in color_defaults.items():
            if key in brand_colors:
                brand_color = brand_colors[key]
                print(f"[DEBUG] Processing color {key}: {brand_color} (type: {type(brand_color)})")
                
                # Handle different color formats
                if hasattr(brand_color, 'r') and hasattr(brand_color, 'g') and hasattr(brand_color, 'b'):
                    # Already an RGBColor object
                    colors[key] = brand_color
                    print(f"[DEBUG] Using RGBColor: {key} = RGB({brand_color.r}, {brand_color.g}, {brand_color.b})")
                elif isinstance(brand_color, tuple) and len(brand_color) == 3:
                    # Tuple format (r, g, b)
                    colors[key] = RGBColor(*brand_color)
                    print(f"[DEBUG] Converted tuple to RGBColor: {key} = {brand_color}")
                elif isinstance(brand_color, str) and brand_color.startswith('#'):
                    # Hex format #RRGGBB
                    hex_color = brand_color.lstrip('#')
                    r = int(hex_color[0:2], 16)
                    g = int(hex_color[2:4], 16) 
                    b = int(hex_color[4:6], 16)
                    colors[key] = RGBColor(r, g, b)
                    print(f"[DEBUG] Converted hex to RGBColor: {key} = {brand_color} -> RGB({r}, {g}, {b})")
                else:
                    print(f"[DEBUG] Unknown color format for {key}: {brand_color}, using default")
                    colors[key] = default_color
            else:
                colors[key] = default_color
        
        # Handle fonts with better error handling
        font_defaults = {
            "primary_font": 'Arial',
            "title_size": Pt(24),
            "header_size": Pt(14),
            "body_size": Pt(11),
            "small_size": Pt(9)
        }
        
        fonts = {}
        for key, default_value in font_defaults.items():
            if key in brand_fonts:
                brand_value = brand_fonts[key]
                print(f"[DEBUG] Processing font {key}: {brand_value} (type: {type(brand_value)})")
                
                if key == "primary_font":
                    fonts[key] = str(brand_value)
                else:
                    # Handle font sizes
                    if hasattr(brand_value, 'pt'):
                        # Already a Pt object
                        fonts[key] = brand_value
                    elif isinstance(brand_value, (int, float)):
                        # Convert number to Pt
                        fonts[key] = Pt(int(brand_value))
                    else:
                        print(f"[DEBUG] Unknown font size format for {key}: {brand_value}, using default")
                        fonts[key] = default_value
            else:
                fonts[key] = default_value
        
        print(f"[DEBUG] Final colors: primary=RGB({colors['primary'].r},{colors['primary'].g},{colors['primary'].b})")
        print(f"[DEBUG] Final fonts: {fonts['primary_font']}, title={fonts['title_size']}")
        
    else:
        print("[DEBUG] No brand_config, using defaults")
        # Use passed parameters or defaults
        colors = color_scheme or {
            "primary": RGBColor(24, 58, 88),
            "secondary": RGBColor(181, 151, 91),
            "accent": RGBColor(64, 64, 64),
            "text": RGBColor(64, 64, 64),
            "background": RGBColor(255, 255, 255),
            "light_grey": RGBColor(240, 240, 240),
            "footer_grey": RGBColor(128, 128, 128)
        }
        
        fonts = typography or {
            "primary_font": 'Arial',
            "title_size": Pt(24),
            "header_size": Pt(14),
            "body_size": Pt(11),
            "small_size": Pt(9)
        }
    
    return colors, fonts


def _apply_standard_header_and_title(slide, title_text, brand_config=None, company_name="Moelis"):
    """Apply standardized header and title formatting to a slide"""
    # Get brand styling
    colors, fonts = get_brand_styling(brand_config)
    
    # Add title with clean header style
    title_left = Inches(0.5)
    title_top = Inches(0.3)
    title_width = Inches(12.333)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.margin_left = 0
    title_frame.margin_top = 0
    title_frame.margin_right = 0
    title_frame.margin_bottom = 0
    
    title_p = title_frame.paragraphs[0]
    title_p.text = title_text
    title_p.alignment = PP_ALIGN.LEFT
    
    title_run = title_p.runs[0]
    title_run.font.name = fonts["primary_font"]
    title_run.font.size = fonts["title_size"]
    title_run.font.bold = True
    title_run.font.color.rgb = colors["primary"]
    
    # Add blue underline
    underline_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        title_left, Inches(1.0), title_width, Inches(0.05)
    )
    underline_shape.fill.solid()
    underline_shape.fill.fore_color.rgb = colors["primary"]
    underline_shape.line.fill.background()


def get_brand_styling(brand_config=None, color_scheme=None, typography=None):
    """Extract brand styling or use defaults - reusable across all functions"""
    if brand_config:
        colors = {}
        
        # Handle different color formats that might be in brand_config
        if 'color_scheme' in brand_config:
            # Standard format: brand_config['color_scheme']
            brand_colors = brand_config['color_scheme']
            for name, color in brand_colors.items():
                if isinstance(color, tuple) and len(color) == 3:
                    # Convert tuple (r, g, b) to RGBColor
                    r, g, b = color
                    colors[name] = RGBColor(r, g, b)
                elif isinstance(color, str) and color.startswith('#'):
                    # Convert hex to RGBColor
                    hex_color = color.lstrip('#')
                    colors[name] = RGBColor(
                        int(hex_color[0:2], 16),
                        int(hex_color[2:4], 16), 
                        int(hex_color[4:6], 16)
                    )
                elif hasattr(color, 'r'):  # Already RGBColor
                    colors[name] = color
                else:
                    # Fallback to default
                    colors[name] = RGBColor(24, 58, 88)
        
        # Handle the actual format from brand extraction: list of tuples
        elif 'extracted_colors' in brand_config:
            extracted_colors = brand_config['extracted_colors']
            if isinstance(extracted_colors, list):
                for name, color_value in extracted_colors:
                    if isinstance(color_value, str) and color_value.startswith('#'):
                        # Convert hex to RGBColor
                        hex_color = color_value.lstrip('#')
                        colors[name] = RGBColor(
                            int(hex_color[0:2], 16),
                            int(hex_color[2:4], 16), 
                            int(hex_color[4:6], 16)
                        )
                    elif isinstance(color_value, tuple) and len(color_value) == 3:
                        # Convert tuple (r, g, b) to RGBColor
                        r, g, b = color_value
                        colors[name] = RGBColor(r, g, b)
                    else:
                        # Fallback to default
                        colors[name] = RGBColor(24, 58, 88)  # Default blue
        
        # Ensure all required colors are present
        default_colors = {
            "primary": RGBColor(24, 58, 88),
            "secondary": RGBColor(181, 151, 91),
            "accent": RGBColor(64, 64, 64),
            "text": RGBColor(64, 64, 64),
            "background": RGBColor(255, 255, 255),
            "light_grey": RGBColor(240, 240, 240),
            "footer_grey": RGBColor(128, 128, 128)
        }
        
        for name, default_color in default_colors.items():
            if name not in colors:
                colors[name] = default_color
        
        # Handle typography
        brand_fonts = brand_config.get('typography', {})
        fonts = {
            "primary_font": brand_fonts.get('primary_font', 'Arial'),
            "title_size": Pt(brand_fonts.get('title_size', 24)),
            "header_size": Pt(brand_fonts.get('header_size', 14)),
            "body_size": Pt(brand_fonts.get('body_size', 11)),
            "small_size": Pt(brand_fonts.get('small_size', 9))
        }
    else:
        # Use passed parameters or defaults
        colors = color_scheme or {
            "primary": RGBColor(24, 58, 88),
            "secondary": RGBColor(181, 151, 91),
            "accent": RGBColor(64, 64, 64),
            "text": RGBColor(64, 64, 64),
            "background": RGBColor(255, 255, 255),
            "light_grey": RGBColor(240, 240, 240),
            "footer_grey": RGBColor(128, 128, 128)
        }
        
        fonts = typography or {
            "primary_font": 'Arial',
            "title_size": Pt(24),
            "header_size": Pt(14),
            "body_size": Pt(11),
            "small_size": Pt(9)
        }
    
    return colors, fonts


def ensure_prs(prs=None):
    """Return a 16:9 Presentation object without isinstance() pitfalls."""
    from pptx.util import Inches
    # If it's already presentation-like, just return
    try:
        if prs is not None and hasattr(prs, "slides") and hasattr(prs, "slide_width"):
            try:
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
            except Exception:
                pass
            return prs
    except Exception:
        pass

    # If it's a path, try opening
    try:
        if isinstance(prs, (str, bytes)) or getattr(prs, "__fspath__", None):
            from pptx import Presentation as _PresentationFactory
            prs_obj = _PresentationFactory(prs)
        else:
            from pptx import Presentation as _PresentationFactory
            prs_obj = _PresentationFactory()
    except Exception:
        from pptx import Presentation as _PresentationFactory
        prs_obj = _PresentationFactory()

    # 16:9
    try:
        prs_obj.slide_width = Inches(13.333)
        prs_obj.slide_height = Inches(7.5)
    except Exception:
        pass
    return prs_obj


def render_management_team_slide(data=None, color_scheme=None, typography=None, company_name="Your Company", prs=None, brand_config=None, **kwargs):
    """
    Renders a sophisticated management team slide with brand configuration support
    """
    
    # Create or use existing presentation
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs = ensure_prs(prs)
    
    # Get brand styling
    colors, fonts = get_brand_styling(brand_config, color_scheme, typography)
    
    # Add blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # STANDARDIZED: Apply header and title
    title_text = (data or {}).get('title', 'Senior Management Team')
    _apply_standard_header_and_title(slide, title_text, brand_config, company_name)
    
    # Calculate total content to determine adaptive spacing
    left_profiles = (data or {}).get('left_column_profiles', [])
    right_profiles = (data or {}).get('right_column_profiles', [])
    
    # Count total content length and profiles
    total_content_length = 0
    total_profiles = len(left_profiles) + len(right_profiles)
    
    for profile in left_profiles + right_profiles:
        bullets = profile.get('experience_bullets', [])
        for bullet in bullets:
            total_content_length += len(bullet) if bullet else 0
    
    # Determine spacing mode - optimize for 3-column layout with 6 profiles
    is_content_heavy = total_content_length > 1500 or total_profiles >= 6  # Lower threshold for 3-column
    is_very_heavy = total_content_length > 2500 or total_profiles > 6
    
    # Adaptive spacing parameters optimized for 3-column layout
    if is_very_heavy:
        profile_spacing = Inches(0.06)  # Very tight spacing between profiles
        bullet_spacing = Inches(0.04)   # Minimal bullet spacing
        title_height = Inches(0.28)     # Compact title height
        title_gap = Inches(0.3)         # Minimal gap after title
        margin_factor = 0.4             # Significantly reduce margins
    elif is_content_heavy:
        profile_spacing = Inches(0.08)  # Tight spacing between profiles
        bullet_spacing = Inches(0.05)   # Reduced bullet spacing
        title_height = Inches(0.32)     # Reduced title height
        title_gap = Inches(0.35)        # Small gap after title
        margin_factor = 0.6             # Moderately reduce margins
    else:
        profile_spacing = Inches(0.1)   # Compact spacing for 3-column
        bullet_spacing = Inches(0.06)   # Compact bullet spacing
        title_height = Inches(0.35)     # Compact title height
        title_gap = Inches(0.4)         # Moderate gap after title
        margin_factor = 0.8             # Some margin reduction for 3-column
    
    # Function to add management profiles with adaptive spacing
    def add_management_profile(x_pos, y_pos, width, profile_data):
        # Role title with adaptive margins
        title_box = slide.shapes.add_textbox(x_pos, y_pos, width, title_height)
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.margin_left = Inches(0.05 * margin_factor)
        title_frame.margin_right = Inches(0.05 * margin_factor)
        title_frame.margin_top = Inches(0.02 * margin_factor)
        title_frame.margin_bottom = Inches(0.02 * margin_factor)
        title_frame.word_wrap = True
        
        p = title_frame.paragraphs[0]
        # Display both name and role title
        name = profile_data.get('name', '')
        role_title = profile_data.get('role_title', 'Role Title')
        
        if name and role_title:
            p.text = f"{name}\n{role_title}"
        elif name:
            p.text = name
        elif role_title:
            p.text = role_title
        else:
            p.text = 'Name\nRole Title'
            
        p.font.name = fonts["primary_font"]
        p.font.size = fonts["header_size"]
        p.font.color.rgb = colors["primary"]
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        # Experience bullets with adaptive sizing
        current_y = y_pos + title_gap  # Adaptive space after title
        experience_bullets = profile_data.get('experience_bullets', [])
        
        for bullet in experience_bullets:
            # Calculate approximate height for 3-column layout (narrower columns)
            chars_per_line = int(width.inches * (12 if is_content_heavy else 10))  # More conservative estimate for narrower columns
            bullet_text = f"• {bullet}"
            estimated_lines = max(1, len(bullet_text) // chars_per_line + (1 if len(bullet_text) % chars_per_line > 0 else 0))
            
            # Adaptive minimum height optimized for 3-column layout with better line spacing
            min_height = 0.28 if is_very_heavy else (0.32 if is_content_heavy else 0.36)
            bullet_height = max(Inches(min_height), Inches(0.18 * estimated_lines))  # Increased height per line
            
            bullet_box = slide.shapes.add_textbox(x_pos, current_y, width, bullet_height)
            bullet_frame = bullet_box.text_frame
            bullet_frame.clear()
            bullet_frame.margin_left = Inches(0.05 * margin_factor)
            bullet_frame.margin_right = Inches(0.05 * margin_factor)
            bullet_frame.margin_top = Inches(0.015 * margin_factor)
            bullet_frame.margin_bottom = Inches(0.015 * margin_factor)
            bullet_frame.word_wrap = True
            bullet_frame.auto_size = None  # Prevent auto-sizing issues
            
            p = bullet_frame.paragraphs[0]
            p.text = bullet_text
            p.font.name = fonts["primary_font"]
            p.font.size = fonts["body_size"]
            p.font.color.rgb = colors["text"]
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.15 if is_content_heavy else 1.2  # Better line spacing to prevent overlap
            
            # Add adaptive spacing between bullets
            current_y += bullet_height + bullet_spacing
        
        return current_y
    
    # Combine all profiles for 3-column layout
    all_profiles = left_profiles + right_profiles
    
    # Dynamic layout logic: ≤4 profiles = centered 2-column, ≥5 profiles = left-aligned 3-column
    total_profiles_count = len(all_profiles)
    
    if total_profiles_count <= 4:
        # Centered 2-column layout for 4 or fewer profiles
        start_y = Inches(1.45)
        column_width = Inches(5.5)  # Wider columns for 2-column layout
        column_spacing = Inches(0.8)  # More spacing between columns
        
        # Center the 2 columns horizontally on the slide
        total_width = 2 * column_width + column_spacing
        start_x = (Inches(13.333) - total_width) / 2  # Center horizontally
        
        col1_x = start_x
        col2_x = col1_x + column_width + column_spacing
        
        profiles_per_column = 2  # 2 profiles per column
        layout_type = "2-column-centered"
    else:
        # Left-aligned 3-column layout for 5 or more profiles
        start_y = Inches(1.45)
        column_width = Inches(4.0)  # Narrower columns for 3-column layout
        column_spacing = Inches(0.15)  # Tighter spacing for more columns
        
        # Left-aligned positioning
        col1_x = Inches(0.3)
        col2_x = col1_x + column_width + column_spacing
        col3_x = col2_x + column_width + column_spacing
        
        profiles_per_column = 2  # 2 profiles per column (fits up to 6 profiles)
        layout_type = "3-column-left-aligned"
        
    print(f"[DEBUG] Management team: {total_profiles_count} profiles → using {layout_type}")
    
    # Column 1 profiles
    current_y = start_y
    for i in range(min(profiles_per_column, len(all_profiles))):
        profile = all_profiles[i]
        current_y = add_management_profile(col1_x, current_y, column_width, profile)
        if i < profiles_per_column - 1:  # Don't add spacing after last profile in column
            current_y += profile_spacing
    
    # Column 2 profiles
    current_y = start_y
    start_idx = profiles_per_column
    for i in range(start_idx, min(start_idx + profiles_per_column, len(all_profiles))):
        profile = all_profiles[i]
        current_y = add_management_profile(col2_x, current_y, column_width, profile)
        if i < start_idx + profiles_per_column - 1:  # Don't add spacing after last profile in column
            current_y += profile_spacing
    
    # Column 3 profiles (only for 3-column layout with 5+ profiles)
    if total_profiles_count >= 5:
        current_y = start_y
        start_idx = profiles_per_column * 2
        for i in range(start_idx, min(start_idx + profiles_per_column, len(all_profiles))):
            profile = all_profiles[i]
            current_y = add_management_profile(col3_x, current_y, column_width, profile)
            if i < start_idx + profiles_per_column - 1:  # Don't add spacing after last profile in column
                current_y += profile_spacing
    
    # Footer with minimal bottom padding
    from datetime import datetime
    today = datetime.now().strftime("%B %d, %Y")
    
    # Position footer with modest bottom margin for better visual appearance
    footer_y = Inches(6.95)  # Balanced footer position with small bottom margin
    
    # Add footer - "Confidential | [today's date]" on LEFT
    footer_left = slide.shapes.add_textbox(Inches(0.4), footer_y, Inches(6), Inches(0.4))
    footer_left_frame = footer_left.text_frame
    footer_left_frame.clear()
    footer_left_frame.margin_left = Inches(0.05)
    footer_left_frame.margin_right = Inches(0.05)
    
    p = footer_left_frame.paragraphs[0]
    p.text = f"Confidential | {today}"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.LEFT
    footer_left_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add footer - Company name on RIGHT
    footer_right = slide.shapes.add_textbox(Inches(9.5), footer_y, Inches(3.2), Inches(0.4))
    footer_right_frame = footer_right.text_frame
    footer_right_frame.clear()
    footer_right_frame.margin_left = Inches(0.05)
    footer_right_frame.margin_right = Inches(0.05)
    
    p = footer_right_frame.paragraphs[0]
    p.text = company_name if company_name and company_name.strip() else "Moelis"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.RIGHT
    footer_right_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return prs


def render_investor_considerations_slide(data=None, color_scheme=None, typography=None, company_name="Moelis", prs=None, brand_config=None, **kwargs):
    """
    Renders an investor considerations slide with two-column layout (Considerations vs Mitigants)
    """
    
    # Create or use existing presentation
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs = ensure_prs(prs)
    
    # Get brand styling
    colors, fonts = get_brand_styling(brand_config, color_scheme, typography)
    
    # Add blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # STANDARDIZED: Apply header and title
    title_text = (data or {}).get('title', 'Investor Considerations & Mitigating Factors')
    _apply_standard_header_and_title(slide, title_text, brand_config, company_name)
    
    # Add column headers
    # Considerations header
    cons_header = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(5.5), Inches(0.5))
    cons_frame = cons_header.text_frame
    cons_frame.clear()
    p = cons_frame.paragraphs[0]
    p.text = "Considerations"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["body_size"]
    p.font.color.rgb = colors["primary"]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Mitigants header
    mit_header = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.5), Inches(0.5))
    mit_frame = mit_header.text_frame
    mit_frame.clear()
    p = mit_frame.paragraphs[0]
    p.text = "Mitigants"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["body_size"]
    p.font.color.rgb = colors["primary"]
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Content positioning - improved spacing for better readability
    y_start = 1.9
    row_height = 1.1  # Increased from 0.75 to 1.1 for better spacing
    
    # Add considerations and mitigants
    considerations = (data or {}).get('considerations', [])
    mitigants = (data or {}).get('mitigants', [])
    
    max_items = max(len(considerations), len(mitigants))
    
    for i in range(max_items):
        y_pos = y_start + (i * row_height)
        
        # Add consideration if exists
        if i < len(considerations):
            # Question mark circle
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y_pos), Inches(0.3), Inches(0.3))
            circle.fill.solid()
            circle.fill.fore_color.rgb = colors["primary"]
            circle.line.fill.background()
            
            # Question mark text
            q_text = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(0.3), Inches(0.3))
            q_frame = q_text.text_frame
            q_frame.clear()
            p = q_frame.paragraphs[0]
            p.text = "?"
            p.font.name = fonts["primary_font"]
            p.font.size = fonts["body_size"]
            p.font.color.rgb = colors["background"]
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            q_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Consideration text - vertically centered with icon
            cons_text = slide.shapes.add_textbox(Inches(1.1), Inches(y_pos - 0.2), Inches(5.2), Inches(0.8))
            cons_text_frame = cons_text.text_frame
            cons_text_frame.clear()
            p = cons_text_frame.paragraphs[0]
            p.text = considerations[i]
            p.font.name = fonts["primary_font"]
            p.font.size = fonts["body_size"]
            p.font.color.rgb = colors["text"]
            p.font.bold = False
            cons_text_frame.word_wrap = True
            cons_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Add mitigant if exists
        if i < len(mitigants):
            # Info circle
            circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.7), Inches(y_pos), Inches(0.3), Inches(0.3))
            circle2.fill.solid()
            circle2.fill.fore_color.rgb = colors["secondary"]
            circle2.line.fill.background()
            
            # Info icon text
            bulb_text = slide.shapes.add_textbox(Inches(6.7), Inches(y_pos), Inches(0.3), Inches(0.3))
            bulb_frame = bulb_text.text_frame
            bulb_frame.clear()
            p = bulb_frame.paragraphs[0]
            p.text = "i"
            p.font.name = fonts["primary_font"]
            p.font.size = fonts["body_size"]
            p.font.color.rgb = colors["background"]
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            bulb_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Mitigant text - vertically centered with icon
            mit_text = slide.shapes.add_textbox(Inches(7.1), Inches(y_pos - 0.2), Inches(5.7), Inches(0.8))
            mit_text_frame = mit_text.text_frame
            mit_text_frame.clear()
            p = mit_text_frame.paragraphs[0]
            p.text = mitigants[i]
            p.font.name = fonts["primary_font"]
            p.font.size = fonts["body_size"]
            p.font.color.rgb = colors["text"]
            p.font.bold = False
            mit_text_frame.word_wrap = True
            mit_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Get today's date
    today = datetime.now().strftime("%B %d, %Y")
    
    # Add footer - "Confidential | [today's date]" on LEFT
    footer_left = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(6), Inches(0.4))
    footer_left_frame = footer_left.text_frame
    footer_left_frame.clear()
    p = footer_left_frame.paragraphs[0]
    p.text = f"Confidential | {today}"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.LEFT
    footer_left_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add footer - Company name on RIGHT
    footer_right = slide.shapes.add_textbox(Inches(10), Inches(7.0), Inches(3), Inches(0.4))
    footer_right_frame = footer_right.text_frame
    footer_right_frame.clear()
    p = footer_right_frame.paragraphs[0]
    p.text = company_name or "Moelis"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.RIGHT
    footer_right_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return prs


def render_product_service_footprint_slide(data=None, color_scheme=None, typography=None, company_name="Moelis", prs=None, brand_config=None, **kwargs):
    """
    Render a product & service / market footprint slide for investment banking presentations
    """
    
    # FIXED: Extract slide_data from data parameter (matching your system's pattern)
    slide_data = data or {}
    
    # Create presentation if not provided (standard 16:9 dimensions)
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs = ensure_prs(prs)
    
    # Get brand styling
    colors, fonts = get_brand_styling(brand_config, color_scheme, typography)
    
    # Add slide with blank layout
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors["background"]
    
    # STANDARDIZED: Apply header and title
    title_text = slide_data.get('title', 'Product & Service / Market Footprint')
    _apply_standard_header_and_title(slide, title_text, brand_config, company_name)
    
    def add_clean_text(slide, left, top, width, height, text, font_size=14, 
                       color=None, bold=False, align=PP_ALIGN.LEFT, bg_color=None):
        """Add text with consistent styling"""
        if color is None:
            color = colors["text"]
            
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = str(text)  # Convert to string to avoid issues
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = align
            for run in paragraph.runs:
                run.font.name = fonts["primary_font"]
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.bold = bold
        
        if bg_color:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = bg_color
        
        textbox.line.fill.background()
        textbox.shadow.inherit = False
        return textbox
    
    # Left side - Service descriptions with gold icons (ALL FROM DATA)
    services = slide_data.get('services', [])
    
    y_start = Inches(1.4)
    for i, service in enumerate(services):
        y_pos = y_start + Inches(i * 0.85)
        
        # Service icon (gold circle background)
        icon_bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y_pos, Inches(0.3), Inches(0.3))
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = colors["secondary"]
        icon_bg.line.fill.background()
        
        # Service title
        add_clean_text(slide, Inches(1.2), y_pos, Inches(5.5), Inches(0.25), 
                       service.get("title", ""), 12, colors["primary"], True)
        
        # Service description
        add_clean_text(slide, Inches(1.2), y_pos + Inches(0.25), Inches(5.5), Inches(0.55), 
                       service.get("desc", ""), 10, colors["text"])
    
    # Right side - Product/Market Table (ALL FROM DATA)
    table_left = Inches(7.2)
    table_width = Inches(5.5)
    
    # Table title from data
    table_title = slide_data.get('table_title', 'Product & Service Market Coverage')
    add_clean_text(slide, table_left, Inches(1.4), table_width, Inches(0.3), 
                   table_title, 14, colors["primary"], True, PP_ALIGN.CENTER)
    
    # Create table from data
    raw_table_data = slide_data.get('coverage_table', [])
    
    # DEBUG: Print what we received
    print(f"[DEBUG] Raw coverage_table data: {raw_table_data}")
    print(f"[DEBUG] Coverage table type: {type(raw_table_data)}")
    if raw_table_data:
        print(f"[DEBUG] First item type: {type(raw_table_data[0]) if raw_table_data else 'None'}")
        print(f"[DEBUG] Table length: {len(raw_table_data)}")
    
    # FIXED: Handle multiple formats for coverage table (object array, 2D array, and string-encoded arrays)
    table_data = []
    
    # Handle string-encoded arrays (common JSON parsing issue)
    if raw_table_data and isinstance(raw_table_data[0], str):
        print(f"[DEBUG] Detected string-encoded arrays, parsing...")
        try:
            import ast
            # Convert string representations back to actual arrays
            parsed_table_data = []
            for item in raw_table_data:
                if item.startswith('[') and item.endswith(']'):
                    parsed_row = ast.literal_eval(item)
                    parsed_table_data.append(parsed_row)
            
            if parsed_table_data:
                table_data = parsed_table_data
                print(f"[DEBUG] Successfully parsed {len(table_data)} string-encoded rows")
            else:
                raise ValueError("No valid rows parsed from string data")
        
        except Exception as e:
            print(f"[DEBUG] Failed to parse string-encoded arrays: {e}")
            # Create fallback table with the string data as-is
            table_data = [["Data Issue", "Format"], ["String encoded", "arrays detected"]]
    
    elif raw_table_data and isinstance(raw_table_data[0], dict):
        # Object array format: [{"region": "Jakarta", "outlets": "45"}]
        print(f"[DEBUG] Converting coverage_table from object array format")
        
        # Create headers based on first object keys
        if raw_table_data:
            headers = list(raw_table_data[0].keys())
            # Convert to proper case
            formatted_headers = [header.replace('_', ' ').title() for header in headers]
            table_data.append(formatted_headers)
            
            # Add data rows
            for item in raw_table_data:
                row = [str(item.get(key, '')) for key in headers]
                table_data.append(row)
        
        print(f"[DEBUG] Converted {len(table_data)-1} coverage table rows")
        
    elif raw_table_data and isinstance(raw_table_data[0], list):
        # Already in 2D array format
        table_data = raw_table_data
        print(f"[DEBUG] Using coverage_table in 2D array format")
        
        # Check for problematic table data (repeated headers)
        if len(table_data) > 1:
            first_row = table_data[0]
            if all(row == first_row for row in table_data[1:]):
                print(f"[DEBUG] Detected repeated headers in coverage_table, creating fallback message")
                table_data = [
                    ["Region", "Coverage Status"],
                    ["Coverage data", "Will be displayed"],
                    ["when available", "from data source"]
                ]
    else:
        # Create simple fallback if no table data provided
        table_data = [
            ["Region", "Status"],
            ["No data", "Please provide"],
            ["coverage_table", "in JSON data"]
        ]
        print(f"[DEBUG] No coverage_table data found - using minimal fallback")
    
    if table_data:  # Only create table if data exists
        rows, cols = len(table_data), len(table_data[0])
        table_top = Inches(1.8)
        table_height = Inches(2.8)
        
        table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table
        
        # Style table
        for i, row_data in enumerate(table_data):
            for j, cell_text in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = str(cell_text)
                cell.margin_left = Inches(0.05)
                cell.margin_right = Inches(0.05)
                cell.margin_top = Inches(0.05)
                cell.margin_bottom = Inches(0.05)
                
                # Style cell text
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.name = fonts["primary_font"]
                        run.font.size = Pt(9)
                        
                        if i == 0:  # Header row
                            run.font.bold = True
                            run.font.color.rgb = colors["background"]
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = colors["primary"]
                        else:
                            run.font.color.rgb = colors["text"]
                            if i % 2 == 0:  # Alternate row colors
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = colors["light_grey"]
                            else:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = colors["background"]
        
        # ENHANCED: Optimal column width distribution for 3-4 column tables
        if cols > 0:
            total_width = table_width
            for i in range(cols):
                if cols == 3:
                    # 3-column layout: Region (40%), Market (30%), Coverage (30%)
                    if i == 0:
                        table.columns[i].width = Inches(2.2)  # Region column wider
                    else:
                        table.columns[i].width = Inches(1.65)  # Other columns equal
                elif cols == 4:
                    # 4-column layout: Region (35%), Market (25%), Assets (20%), Coverage (20%)
                    if i == 0:
                        table.columns[i].width = Inches(1.9)  # Region column wider
                    elif i == 1:
                        table.columns[i].width = Inches(1.4)  # Market column
                    else:
                        table.columns[i].width = Inches(1.1)  # Assets/Coverage columns
                else:
                    # Default for 2 or 5+ columns
                    if i == 0:  # First column wider for location names
                        table.columns[i].width = Inches(1.2)
                    else:
                        table.columns[i].width = Inches(0.8)
    
    # Key Operational Metrics (ALL FROM DATA)
    metrics_title_top = Inches(4.8)
    metrics_box_top = Inches(5.1)
    
    # Metrics title from data
    metrics_title = slide_data.get('metrics_title', 'Key Operational Metrics')
    add_clean_text(slide, table_left, metrics_title_top, table_width, Inches(0.3), 
                   metrics_title, 14, colors["primary"], True, PP_ALIGN.CENTER)
    
    # Metrics background box - increased height to cover all content
    metrics_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, table_left, metrics_box_top, table_width, Inches(2.0))
    metrics_bg.fill.solid()
    metrics_bg.fill.fore_color.rgb = colors["light_grey"]
    metrics_bg.line.fill.background()
    
    # Get metrics data - handle both formats: key_metrics.metrics and direct metrics
    metrics = {}
    if 'key_metrics' in slide_data and 'metrics' in slide_data['key_metrics']:
        # Format: key_metrics.metrics (array of strings or objects) 
        metrics_array = slide_data['key_metrics']['metrics']
        if isinstance(metrics_array, list):
            # Convert array to dict for processing
            for i, metric in enumerate(metrics_array):
                if isinstance(metric, str):
                    # String format - use as value with generated key
                    metrics[f"metric_{i+1}"] = metric
                elif isinstance(metric, dict):
                    # Object format - use existing structure
                    key = metric.get('key', f"metric_{i+1}")
                    metrics[key] = metric
            print(f"[DEBUG] Using key_metrics.metrics format with {len(metrics)} metrics")
    elif 'metrics' in slide_data:
        # Format: direct metrics (object with key-value pairs)
        metrics = slide_data.get('metrics', {})
        print(f"[DEBUG] Using direct metrics format with {len(metrics)} metrics")
    
    metric_keys = list(metrics.keys())
    
    # FIXED: Add fallback message if no metrics are available
    if not metrics:
        print(f"[DEBUG] No metrics data available for Product & Service slide")
        # Add a "No metrics available" message in the box
        add_clean_text(slide, table_left + Inches(0.3), metrics_box_top + Inches(0.8), Inches(4.9), Inches(0.4), 
                       "Key operational metrics will be displayed here when data is available.", 
                       12, colors["text"], False, PP_ALIGN.CENTER)
    
    # Dynamic metrics layout
    metrics_left_col = table_left + Inches(0.3)
    metrics_right_col = table_left + Inches(2.8)
    col_width = Inches(2.2)
    
    # Left column metrics (first half)
    left_metrics = metric_keys[:len(metric_keys)//2]
    for i, key in enumerate(left_metrics):
        metric_data = metrics[key]
        
        # FIXED: Handle both string and object metric formats
        if isinstance(metric_data, str):
            # String format: convert to object structure
            label = key.replace('_', ' ').title()
            value = metric_data
            print(f"[DEBUG] Converting string metric '{key}': {metric_data}")
        elif isinstance(metric_data, dict):
            # Object format: use existing structure
            label = metric_data.get('label', key.replace('_', ' ').title())
            value = metric_data.get('value', '')
        else:
            # Fallback for other types
            label = key.replace('_', ' ').title()
            value = str(metric_data)
            print(f"[DEBUG] Unknown metric format for '{key}': {type(metric_data)}")
        
        y_offset = Inches(0.2 + i * 0.55)
        add_clean_text(slide, metrics_left_col, metrics_box_top + y_offset, col_width, Inches(0.2), 
                       label, 10, colors["text"])
        add_clean_text(slide, metrics_left_col, metrics_box_top + y_offset + Inches(0.2), col_width, Inches(0.25), 
                       value, 16, colors["primary"], True)
    
    # Right column metrics (second half)
    right_metrics = metric_keys[len(metric_keys)//2:]
    for i, key in enumerate(right_metrics):
        metric_data = metrics[key]
        
        # FIXED: Handle both string and object metric formats
        if isinstance(metric_data, str):
            # String format: convert to object structure
            label = key.replace('_', ' ').title()
            value = metric_data
            print(f"[DEBUG] Converting string metric '{key}': {metric_data}")
        elif isinstance(metric_data, dict):
            # Object format: use existing structure
            label = metric_data.get('label', key.replace('_', ' ').title())
            value = metric_data.get('value', '')
        else:
            # Fallback for other types
            label = key.replace('_', ' ').title()
            value = str(metric_data)
            print(f"[DEBUG] Unknown metric format for '{key}': {type(metric_data)}")
        
        y_offset = Inches(0.2 + i * 0.55)
        add_clean_text(slide, metrics_right_col, metrics_box_top + y_offset, col_width, Inches(0.2), 
                       label, 10, colors["text"])
        add_clean_text(slide, metrics_right_col, metrics_box_top + y_offset + Inches(0.2), col_width, Inches(0.25), 
                       value, 16, colors["primary"], True)
    
    # Get today's date
    today = datetime.now().strftime("%B %d, %Y")
    
    # Add footer - "Confidential | [today's date]" on LEFT
    footer_left = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(6), Inches(0.4))
    footer_left_frame = footer_left.text_frame
    footer_left_frame.clear()
    p = footer_left_frame.paragraphs[0]
    p.text = f"Confidential | {today}"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.LEFT
    footer_left_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add footer - Company name on RIGHT
    footer_right = slide.shapes.add_textbox(Inches(9.5), Inches(7.0), Inches(3.5), Inches(0.4))
    footer_right_frame = footer_right.text_frame
    footer_right_frame.clear()
    p = footer_right_frame.paragraphs[0]
    p.text = company_name or "Moelis"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.RIGHT
    footer_right_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return prs


def render_competitive_positioning_slide(data=None, color_scheme=None, typography=None, company_name="Moelis", prs=None, brand_config=None, **kwargs):
    """
    Render an ENHANCED competitive positioning slide matching iCar Asia format
    Features: 5-column assessment table with star ratings, clean layout, comprehensive data
    """
    
    # Create presentation if not provided
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs = ensure_prs(prs)
    
    # Get brand styling
    colors, fonts = get_brand_styling(brand_config, color_scheme, typography)
    
    # Get brand styling
    colors, fonts = get_brand_styling(brand_config, color_scheme, typography)
    
    # Add slide with blank layout
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors["background"]
    
    def add_clean_text(slide, left, top, width, height, text, font_size=14, 
                       color=None, bold=False, align=PP_ALIGN.LEFT, bg_color=None):
        """Add text with consistent styling"""
        if color is None:
            color = colors["text"]
            
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = align
            for run in paragraph.runs:
                run.font.name = fonts["primary_font"]
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.bold = bold
        
        if bg_color:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = bg_color
        
        textbox.line.fill.background()
        textbox.shadow.inherit = False
        return textbox
    
    # Extract slide data
    slide_data = data or {}
    
    # STANDARDIZED: Apply header and title
    title_text = slide_data.get('title', 'Competitive Positioning')
    _apply_standard_header_and_title(slide, title_text, brand_config, company_name)
    
    # Left side - Revenue Comparison Chart
    add_clean_text(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.3), 
                   "Revenue Comparison vs. Competitors", 14, colors["primary"], True)
    
    # Create bar chart data from user's competitive analysis
    competitors_data = slide_data.get('competitors', [])
    
    # If no competitors provided, create generic placeholder data
    if not competitors_data:
        company_name = content_ir.get('entities', {}).get('company', {}).get('name', 'Company')
        competitors_data = [
            {'name': company_name, 'revenue': 50},  # User's company
            {'name': 'Competitor A', 'revenue': 45},
            {'name': 'Competitor B', 'revenue': 60},
            {'name': 'Competitor C', 'revenue': 35},
            {'name': 'Market Leader', 'revenue': 80}
        ]
    
    chart_data = ChartData()
    chart_data.categories = [comp['name'] for comp in competitors_data]
    chart_data.add_series('Revenue ($M)', [comp['revenue'] for comp in competitors_data])
    
    # Add chart
    chart_left = Inches(0.5)
    chart_top = Inches(1.7)
    chart_width = Inches(6)
    chart_height = Inches(2.5)
    
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, chart_left, chart_top, chart_width, chart_height, chart_data
    )
    
    chart = chart_shape.chart
    
    # Style the chart
    chart.has_legend = False
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.clear()
    
    # Style chart elements
    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    category_axis.tick_labels.font.size = Pt(9)
    category_axis.tick_labels.font.name = fonts["primary_font"]
    
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.tick_labels.font.size = Pt(9)
    value_axis.tick_labels.font.name = fonts["primary_font"]
    
    # FIXED: Dynamic axis scaling for competitor revenue data
    revenue_values = [comp['revenue'] for comp in competitors_data]
    if revenue_values:
        data_max = max(revenue_values)
        # Set axis maximum to 120% of highest value
        dynamic_max = int(data_max * 1.2)
        value_axis.maximum_scale = dynamic_max
        print(f"[DEBUG] Competitive chart axis scaled: max {data_max:,} → {dynamic_max:,}")
    else:
        value_axis.maximum_scale = 500  # Fallback
    
    # Highlight the user's company (first in list or matching company name) in secondary color
    user_company_name = content_ir.get('entities', {}).get('company', {}).get('name', '')
    series = chart.series[0]
    points = series.points
    for i, point in enumerate(points):
        competitor_name = competitors_data[i]['name']
        # Highlight if it's the user's company (exact match) or first competitor (assumed to be user's company)
        if (user_company_name and competitor_name == user_company_name) or i == 0:
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = colors["secondary"]  # Gold/highlight color
        else:
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = colors["primary"]  # Standard blue
    
    # Right side - Competitive Assessment Table
    add_clean_text(slide, Inches(7.5), Inches(1.3), Inches(5.5), Inches(0.3), 
                   "Competitive Assessment", 14, colors["primary"], True)
    
    # Assessment table data - FIXED: Handle both object array and 2D array formats
    raw_assessment = slide_data.get('assessment', [])
    
    print(f"[DEBUG] Raw assessment type: {type(raw_assessment)}, length: {len(raw_assessment) if raw_assessment else 0}")
    if raw_assessment:
        print(f"[DEBUG] First assessment item type: {type(raw_assessment[0])}, value: {raw_assessment[0]}")
    
    # Check if we have the new object format (user data) or old 2D array format
    if raw_assessment and len(raw_assessment) > 0 and isinstance(raw_assessment[0], dict):
        # Convert object format to 2D array format
        print(f"[DEBUG] Converting assessment from object format to table format")
        assessment_data = [
            ["Category", "Our Company", "Competitor A", "Competitor B"]
        ]
        
        for item in raw_assessment[:6]:  # Limit to 6 rows to fit slide
            row = [
                item.get('category', 'N/A'),
                item.get('our_company', 'N/A'), 
                item.get('competitor_a', 'N/A'),
                item.get('competitor_b', 'N/A')
            ]
            assessment_data.append(row)
        
        print(f"[DEBUG] Converted {len(assessment_data)-1} assessment rows")
            
    elif raw_assessment and len(raw_assessment) > 0 and isinstance(raw_assessment[0], list):
        # Already in 2D array format - USE AS-IS
        assessment_data = raw_assessment
        print(f"[DEBUG] Using assessment data in 2D array format, {len(assessment_data)} rows")
        
        # Print first few rows for debugging
        for i, row in enumerate(assessment_data[:3]):
            print(f"[DEBUG] Assessment row {i}: {row}")
    else:
        # Fallback: create proper table from user's actual data or use default
        print(f"[DEBUG] Creating fallback assessment table")
        assessment_data = [
            ["Company", "Market Focus", "Connectors/Indexing", "Enterprise Adoption", "Factuality/Traceability"],
            ["LlamaIndex", "⭐⭐⭐⭐⭐", "⭐⭐⭐⭐⭐", "⭐⭐⭐⭐", "⭐⭐⭐⭐⭐"],
            ["LangChain", "⭐⭐⭐⭐", "⭐⭐⭐⭐", "⭐⭐⭐⭐", "⭐⭐⭐⭐"],
            ["CrewAI", "⭐⭐⭐", "⭐⭐⭐", "⭐⭐", "⭐⭐⭐"],
            ["OpenAI API", "⭐⭐⭐⭐", "⭐⭐", "⭐⭐⭐⭐", "⭐⭐⭐⭐"],
            ["Haystack", "⭐⭐⭐⭐", "⭐⭐⭐⭐", "⭐⭐⭐", "⭐⭐⭐⭐"]
        ]
        print(f"[DEBUG] Using LlamaIndex competitive assessment data")
    
    # Create assessment table
    table_left = Inches(7.5)
    table_top = Inches(1.7)
    table_width = Inches(5.5)
    table_height = Inches(1.5)
    
    rows, cols = len(assessment_data), len(assessment_data[0])
    table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table
    
    # Style assessment table
    for i, row_data in enumerate(assessment_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = fonts["primary_font"]
                    run.font.size = Pt(8)
                    
                    if i == 0:  # Header row
                        run.font.bold = True
                        run.font.color.rgb = colors["background"]
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = colors["primary"]
                    elif i == 1:  # Target company row (highlighted)
                        run.font.color.rgb = colors["primary"]
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = colors["light_grey"]
                    else:
                        run.font.color.rgb = colors["text"]
                        if i % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = colors["light_grey"]
    
    # ENHANCED: Optimal column widths for 5-column iCar Asia format
    if cols >= 5:
        col_widths = [Inches(1.3), Inches(1.0), Inches(1.0), Inches(0.9), Inches(0.8)]  # Company, Market Share, Tech Platform, Coverage, Revenue
        for i, width in enumerate(col_widths[:cols]):
            table.columns[i].width = width
    elif cols == 4:  # New format: Category, Our Company, Competitor A, Competitor B
        col_widths = [Inches(1.5), Inches(1.0), Inches(1.0), Inches(1.0)]
        for i, width in enumerate(col_widths):
            if i < len(table.columns):  # Safety check
                table.columns[i].width = width
    else:
        # Fallback for fewer columns
        equal_width = Inches(5.5) / cols  # table_width / cols
        for i in range(cols):
            table.columns[i].width = equal_width
    
    # ENHANCED: Source note positioned better for 5-column table
    add_clean_text(slide, Inches(6.8), Inches(4.0), Inches(6), Inches(0.2), 
                   "Source: Management estimates, competitor websites, July 2024 [Medium Confidence]", 
                   8, colors["text"])
    
    # Bottom left - Barriers to Entry
    add_clean_text(slide, Inches(0.5), Inches(4.5), Inches(6), Inches(0.3), 
                   "Barriers to Entry", 14, colors["primary"], True)
    
    barriers = slide_data.get('barriers', [
        {"title": "Regulatory Compliance:", "desc": "Stringent healthcare licensing requirements and facility standards"},
        {"title": "Specialist Recruitment:", "desc": "Challenging acquisition of multilingual medical talent"},
        {"title": "Prime Real Estate:", "desc": "Limited availability and high cost of clinic locations"},
        {"title": "Insurance Relationships:", "desc": "Established direct billing partnerships with 35+ insurers"}
    ])
    
    y_start = Inches(4.9)
    for i, barrier in enumerate(barriers):
        y_pos = y_start + Inches(i * 0.35)
        
        # Gold bullet
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y_pos, Inches(0.06), Inches(0.06))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = colors["secondary"]
        bullet.line.fill.background()
        
        # Barrier text
        barrier_text = f"{barrier['title']} {barrier['desc']}"
        add_clean_text(slide, Inches(0.85), y_pos - Inches(0.05), Inches(5.5), Inches(0.3), 
                       barrier_text, 9, colors["text"])
    
    # Bottom right - Company's Unique Advantages
    add_clean_text(slide, Inches(7.5), Inches(3.8), Inches(5.5), Inches(0.3), 
                   "Company's Unique Advantages", 14, colors["primary"], True)
    
    advantages = slide_data.get('advantages', [
        {"title": "International Accreditation:", "desc": "First Hong Kong clinic accredited by Australian Council"},
        {"title": "Multi-specialty Integration:", "desc": "Comprehensive holistic care model spanning physical and mental health"},
        {"title": "In-house Pharmacy:", "desc": "Holder of wholesale pharmacy license with dedicated pharmacy team"},
        {"title": "Teaching Status:", "desc": "Recognized undergraduate and postgraduate teaching unit"}
    ])
    
    y_start = Inches(4.2)
    for i, advantage in enumerate(advantages):
        y_pos = y_start + Inches(i * 0.35)
        
        # Gold bullet
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.7), y_pos, Inches(0.06), Inches(0.06))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = colors["secondary"]
        bullet.line.fill.background()
        
        # Advantage text
        advantage_text = f"{advantage['title']} {advantage['desc']}"
        add_clean_text(slide, Inches(7.85), y_pos - Inches(0.05), Inches(5), Inches(0.3), 
                       advantage_text, 9, colors["text"])
    
    # Source note for chart
    add_clean_text(slide, Inches(0.5), Inches(4.3), Inches(6), Inches(0.15), 
                   "Source: Company analysis, industry reports, 2024 [High Confidence]", 
                   8, colors["text"])
    
    # Get today's date
    today = datetime.now().strftime("%B %d, %Y")
    
    # Add footer - "Confidential | [today's date]" on LEFT
    footer_left = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(4), Inches(0.4))
    footer_left_frame = footer_left.text_frame
    footer_left_frame.clear()
    p = footer_left_frame.paragraphs[0]
    p.text = f"Confidential | {today}"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.LEFT
    footer_left_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add footer - Company name on RIGHT
    footer_right = slide.shapes.add_textbox(Inches(9.5), Inches(7.0), Inches(3.5), Inches(0.4))
    footer_right_frame = footer_right.text_frame
    footer_right_frame.clear()
    p = footer_right_frame.paragraphs[0]
    p.text = company_name or "Moelis"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.RIGHT
    footer_right_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return prs


def render_investor_process_overview_slide(data=None, color_scheme=None, typography=None, company_name="Moelis", prs=None, brand_config=None, **kwargs):
    """
    Renders an investor considerations & process overview slide with 4-quadrant layout
    IMPROVED: Better spacing and text fitting for comprehensive content
    """
    
    # FIXED: Extract slide_data from data parameter
    slide_data = data or {}
    
    # Create or use existing presentation
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs = ensure_prs(prs)
    
    # Get brand styling
    colors, fonts = get_brand_styling(brand_config, color_scheme, typography)
    
    # Add blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # STANDARDIZED: Apply header and title
    title_text = slide_data.get('title', 'Investor Process Overview - Comprehensive Due Diligence')
    _apply_standard_header_and_title(slide, title_text, brand_config, company_name)
    
    # Helper function to add clean text with better wrapping
    def add_clean_text(slide, left, top, width, height, text, font_size=10, 
                       color=colors["text"], bold=False, align=PP_ALIGN.LEFT):
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.margin_left = Inches(0.08)
        text_frame.margin_right = Inches(0.08)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.word_wrap = True
        text_frame.auto_size = None  # Prevent auto-sizing issues
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = align
            paragraph.line_spacing = 1.1  # Better line spacing
            for run in paragraph.runs:
                run.font.name = fonts["primary_font"]
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.bold = bold
        
        textbox.line.fill.background()
        textbox.shadow.inherit = False
        return textbox
    
    # Debug logging
    print(f"[DEBUG] Investor process overview data keys: {list(slide_data.keys())}")
    
    # TOP LEFT: Key Diligence Topics - IMPROVED SPACING
    add_clean_text(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(0.3), 
                   "Key Diligence Topics", 14, colors["primary"], True)
    
    diligence_items = slide_data.get('diligence_topics', [])
    print(f"[DEBUG] Diligence topics count: {len(diligence_items)}")
    
    y_start = Inches(1.85)  # Moved down for more breathing room
    for i, item in enumerate(diligence_items[:5]):  # Allow 5 items now
        y_pos = y_start + Inches(i * 0.42)  # Increased spacing for multi-line content
        
        # Gold bullet - centered vertically with text
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.65), y_pos + Inches(0.08), Inches(0.06), Inches(0.06))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = colors["secondary"]
        bullet.line.fill.background()
        
        # Item text - LARGER TEXT BOX with better alignment
        if isinstance(item, dict):
            item_text = f"{item.get('title', '')}: {item.get('description', '')}"
        else:
            item_text = str(item)
        add_clean_text(slide, Inches(0.8), y_pos, Inches(5.6), Inches(0.28), 
                       item_text, 10, colors["text"])  # Increased font size from 9 to 10
    
    # BOTTOM LEFT: Risk Factors & Mitigants - REPOSITIONED LOWER
    add_clean_text(slide, Inches(0.5), Inches(4.0), Inches(6), Inches(0.3), 
                   "Risk Factors & Mitigants", 14, colors["primary"], True)
    
    risk_factors = slide_data.get('risk_factors', [])
    mitigants = slide_data.get('mitigants', [])
    
    print(f"[DEBUG] Risk factors count: {len(risk_factors)}")
    print(f"[DEBUG] Mitigants count: {len(mitigants)}")
    
    y_start = Inches(4.35)  # Moved down further to avoid overlap
    max_items = max(len(risk_factors), len(mitigants), 5)  # Allow 5 items
    
    for i in range(min(max_items, 5)):  # Max 5 to fit properly
        y_pos = y_start + Inches(i * 0.35)  # Increased spacing for multi-line content
        
        # Risk factor (if exists)
        if i < len(risk_factors):
            # Red circle for risk - centered vertically with text
            risk_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.65), y_pos + Inches(0.06), Inches(0.06), Inches(0.06))
            risk_circle.fill.solid()
            risk_circle.fill.fore_color.rgb = RGBColor(220, 20, 60)  # Red color
            risk_circle.line.fill.background()
            
            # Risk factor text - IMPROVED SIZING with better alignment
            add_clean_text(slide, Inches(0.8), y_pos, Inches(2.6), Inches(0.22), 
                           risk_factors[i], 9, RGBColor(220, 20, 60))  # Increased font size from 8 to 9
        
        # Mitigant (if exists)
        if i < len(mitigants):
            # Green circle for mitigant - centered vertically with text
            mitigant_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.6), y_pos + Inches(0.06), Inches(0.06), Inches(0.06))
            mitigant_circle.fill.solid()
            mitigant_circle.fill.fore_color.rgb = RGBColor(34, 139, 34)  # Green color
            mitigant_circle.line.fill.background()
            
            # Mitigant text - IMPROVED SIZING with better alignment
            add_clean_text(slide, Inches(3.75), y_pos, Inches(2.8), Inches(0.22), 
                           mitigants[i], 9, RGBColor(34, 139, 34))  # Increased font size from 8 to 9
    
    # TOP RIGHT: Investment Highlights - IMPROVED
    synergy_items = slide_data.get('synergy_opportunities', [])
    
    # If no synergy opportunities, show investment highlights instead
    if not synergy_items:
        add_clean_text(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(0.3), 
                       "Investment Highlights", 14, colors["primary"], True)
        
        investment_highlights = slide_data.get('investment_highlights', [])
        
        y_start = Inches(1.85)  # Moved down to match left side
        for i, highlight in enumerate(investment_highlights[:5]):  # Allow 5 highlights
            y_pos = y_start + Inches(i * 0.42)  # Increased spacing for multi-line content
            
            # Gold bullet - centered vertically with text
            bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.15), y_pos + Inches(0.08), Inches(0.06), Inches(0.06))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = colors["secondary"]
            bullet.line.fill.background()
            
            # Highlight text - BETTER SIZING with improved alignment
            add_clean_text(slide, Inches(7.3), y_pos, Inches(5.5), Inches(0.28), 
                           highlight, 10, colors["text"])  # Increased font size from 9 to 10
    else:
        add_clean_text(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(0.3), 
                       "Synergy Opportunities", 14, colors["primary"], True)
        
        y_start = Inches(1.85)  # Moved down to match
        for i, item in enumerate(synergy_items[:5]):
            y_pos = y_start + Inches(i * 0.42)  # Increased spacing for multi-line content
            
            # Gold bullet - centered vertically with text
            bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.15), y_pos + Inches(0.08), Inches(0.06), Inches(0.06))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = colors["secondary"]
            bullet.line.fill.background()
            
            # Item text with better alignment
            if isinstance(item, dict):
                item_text = f"{item.get('title', '')}: {item.get('description', '')}"
            else:
                item_text = str(item)
            add_clean_text(slide, Inches(7.3), y_pos, Inches(5.5), Inches(0.28), 
                           item_text, 10, colors["text"])  # Increased font size from 9 to 10
    
    # BOTTOM RIGHT: Process Next Steps - IMPROVED
    timeline_items = slide_data.get('timeline', [])
    
    if not timeline_items:
        add_clean_text(slide, Inches(7), Inches(4.0), Inches(5.8), Inches(0.3), 
                       "Process Next Steps", 14, colors["primary"], True)
        
        next_steps = slide_data.get('next_steps', [])
        
        y_start = Inches(4.35)  # Moved down to match left side positioning
        for i, item in enumerate(next_steps[:5]):  # Allow 5 steps
            y_pos = y_start + Inches(i * 0.35)  # Increased spacing for multi-line content
            
            # Gold circle marker - centered vertically with text
            marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.15), y_pos + Inches(0.06), Inches(0.06), Inches(0.06))
            marker.fill.solid()
            marker.fill.fore_color.rgb = colors["secondary"]
            marker.line.fill.background()
            
            # Handle both string and dict timeline items
            if isinstance(item, dict):
                combined_text = f"{item.get('date', '')}: {item.get('description', '')}"
            else:
                # Handle string format like "Preparation: 2–3 weeks – NDA, data room, initial buyer list"
                combined_text = str(item)
            add_clean_text(slide, Inches(7.3), y_pos, Inches(5.5), Inches(0.22), 
                           combined_text, 9, colors["text"])  # Increased font size from 8 to 9
    else:
        add_clean_text(slide, Inches(7), Inches(4.0), Inches(5.8), Inches(0.3), 
                       "Transaction Timeline", 14, colors["primary"], True)
        
        y_start = Inches(4.35)  # Match positioning with left side
        for i, item in enumerate(timeline_items[:5]):
            y_pos = y_start + Inches(i * 0.35)  # Increased spacing for multi-line content
            
            # Gold circle marker - centered vertically with text
            marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.15), y_pos + Inches(0.06), Inches(0.06), Inches(0.06))
            marker.fill.solid()
            marker.fill.fore_color.rgb = colors["secondary"]
            marker.line.fill.background()
            
            # Handle both string and dict timeline items
            if isinstance(item, dict):
                combined_text = f"{item.get('date', '')}: {item.get('description', '')}"
            else:
                # Handle string format like "Preparation: 2–3 weeks – NDA, data room, initial buyer list"
                combined_text = str(item)
            add_clean_text(slide, Inches(7.3), y_pos, Inches(5.5), Inches(0.22), 
                           combined_text, 9, colors["text"])  # Increased font size from 8 to 9
    
    # Get today's date
    today = datetime.now().strftime("%B %d, %Y")
    
    # Add footer - "Confidential | [today's date]" on LEFT (light grey)
    footer_left = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(6), Inches(0.4))
    footer_left_frame = footer_left.text_frame
    footer_left_frame.clear()
    p = footer_left_frame.paragraphs[0]
    p.text = f"Confidential | {today}"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.LEFT
    footer_left_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add footer - Company name on RIGHT
    footer_right = slide.shapes.add_textbox(Inches(10), Inches(6.9), Inches(3), Inches(0.4))
    footer_right_frame = footer_right.text_frame
    footer_right_frame.clear()
    p = footer_right_frame.paragraphs[0]
    p.text = company_name or "Moelis"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.RIGHT
    footer_right_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return prs


def render_margin_cost_resilience_slide(data=None, color_scheme=None, typography=None, company_name="Moelis", prs=None, brand_config=None, **kwargs):
    """
    Renders a margin & cost resilience slide with charts and detailed analysis
    """
    
    # Create or use existing presentation
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs = ensure_prs(prs)
    
    # Get brand styling
    colors, fonts = get_brand_styling(brand_config, color_scheme, typography)
    
    # Add blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Helper function to add clean text
    def add_clean_text(slide, left, top, width, height, text, font_size=10, 
                       color=colors["text"], bold=False, align=PP_ALIGN.LEFT, bg_color=None):
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = align
            for run in paragraph.runs:
                run.font.name = fonts["primary_font"]
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.bold = bold
        
        if bg_color:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = bg_color
        
        # Remove shadow/outline
        textbox.line.fill.background()
        textbox.shadow.inherit = False
        
        return textbox
    
    # Debug logging
    print(f"[DEBUG] Margin cost resilience data keys: {list((data or {}).keys())}")
    
    # Extract slide data - handle both data formats
    slide_data = data or {}
    
    # STANDARDIZED: Apply header and title
    title_text = slide_data.get('title', 'Margin & Cost Resilience')
    _apply_standard_header_and_title(slide, title_text, brand_config, company_name)
    
    # EBITDA Margin Trend Chart
    chart_title = slide_data.get('chart_title', 'EBITDA Margin Trend')
    add_clean_text(slide, Inches(1), Inches(1.4), Inches(6), Inches(0.3), 
                   chart_title, 14, colors["primary"], True)
    
    # Create line chart for EBITDA margins
    chart_data_info = slide_data.get('chart_data', {})
    
    # Create chart if we have data
    if chart_data_info:
        try:
            chart_data = ChartData()
            
            categories = chart_data_info.get('categories', ['2020', '2021', '2022', '2023', '2024E'])
            values = chart_data_info.get('values', [15.0, 16.6, 17.2, 19.0, 19.6])
            
            print(f"[DEBUG] Chart categories: {categories}")
            print(f"[DEBUG] Chart values: {values}")
            
            chart_data.categories = categories
            chart_data.add_series('EBITDA Margin %', values)
            
            # Add line chart
            chart_left = Inches(1)
            chart_top = Inches(1.8)
            chart_width = Inches(6)
            chart_height = Inches(2.2)
            
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE_MARKERS, chart_left, chart_top, chart_width, chart_height, chart_data
            )
            
            chart = chart_shape.chart
            
            # Style the chart
            chart.has_legend = False
            chart.chart_title.has_text_frame = True
            chart.chart_title.text_frame.clear()
            
            # Style chart elements
            try:
                category_axis = chart.category_axis
                category_axis.has_major_gridlines = False
                category_axis.tick_labels.font.size = Pt(10)
                category_axis.tick_labels.font.name = fonts["primary_font"]
                
                value_axis = chart.value_axis
                value_axis.has_major_gridlines = True
                value_axis.tick_labels.font.size = Pt(10)
                value_axis.tick_labels.font.name = fonts["primary_font"]
                
                # FIXED: Dynamic axis scaling for margin percentage data
                if values:
                    data_max = max(values)
                    # For percentage data, use reasonable scaling (e.g., round up to nearest 10%)
                    if data_max <= 100:  # Percentage data
                        dynamic_max = int((data_max * 1.2 + 9) // 10 * 10)  # Round up to nearest 10%
                    else:  # Large values
                        dynamic_max = int(data_max * 1.2)
                    value_axis.maximum_scale = dynamic_max
                    print(f"[DEBUG] Margin chart axis scaled: max {data_max}% → {dynamic_max}%")
                else:
                    value_axis.maximum_scale = 50  # Reasonable default for percentages
                
                # Style the line series (secondary color)
                series = chart.series[0]
                series.format.line.color.rgb = colors["secondary"]
                series.format.line.width = Pt(3)
                
                # Style the markers
                for point in series.points:
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = colors["secondary"]
                    point.format.line.color.rgb = colors["secondary"]
            except Exception as e:
                print(f"[DEBUG] Chart styling error: {e}")
                
        except Exception as e:
            print(f"[DEBUG] Chart creation error: {e}")
            # Add fallback text if chart fails
            add_clean_text(slide, Inches(1), Inches(2), Inches(6), Inches(1), 
                           "EBITDA margin trend chart will be displayed here.", 12, colors["text"])
    else:
        # Add fallback if no chart data
        add_clean_text(slide, Inches(1), Inches(2), Inches(6), Inches(1), 
                       "EBITDA margin trend chart will be displayed here.", 12, colors["text"])
    
    # Chart source note
    chart_source = slide_data.get('chart_source', 'Source: Company financials')
    add_clean_text(slide, Inches(5), Inches(4.1), Inches(2.5), Inches(0.15), 
                   chart_source, 8, colors["text"], False, PP_ALIGN.RIGHT)
    
    # Left side - Cost Management Strategies - Enhanced debugging
    cost_section = slide_data.get('cost_management', {})
    has_cost_data = bool(cost_section and cost_section.get('items'))
    print(f"[DEBUG] Cost management data available: {has_cost_data}")
    
    cost_title = cost_section.get('title', 'Cost Management & Efficiency Initiatives')
    add_clean_text(slide, Inches(1), Inches(4.4), Inches(6), Inches(0.3), 
                   cost_title, 14, colors["primary"], True)
    
    # Cost management items - Enhanced with better fallbacks
    if has_cost_data:
        cost_items = cost_section.get('items', [])
        print(f"[DEBUG] Using provided cost management items: {len(cost_items)}")
    else:
        print(f"[DEBUG] Using fallback cost management items")
    
    cost_items = cost_section.get('items', [
        {
            "title": "Operational Efficiency",
            "description": "Streamlined processes and resource optimization initiatives"
        },
        {
            "title": "Technology Investment", 
            "description": "Digital transformation reducing administrative overhead"
        },
        {
            "title": "Supply Chain Management",
            "description": "Centralized procurement and vendor consolidation"
        }
    ])
    
    y_start = Inches(4.8)
    
    for i, item in enumerate(cost_items[:3]):  # Max 3 items to fit properly
        y_pos = y_start + Inches(i * 0.45)  # Reduced spacing from 0.5 to 0.45
        
        # Gold bullet (no shadow)
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y_pos, Inches(0.06), Inches(0.06))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = colors["secondary"]
        bullet.line.fill.background()
        bullet.shadow.inherit = False
        
        # Item title (bold) - reduced height
        add_clean_text(slide, Inches(1.35), y_pos - Inches(0.05), Inches(5.5), Inches(0.15), 
                       item.get('title', ''), 10, colors["primary"], True)
        
        # Item description - reduced height and adjusted position
        add_clean_text(slide, Inches(1.35), y_pos + Inches(0.1), Inches(5.5), Inches(0.28), 
                       item.get('description', ''), 9, colors["text"])
    
    # Right side - Risk Mitigation Strategies - ENHANCED: Handle missing data
    risk_section = slide_data.get('risk_mitigation', {})
    
    # Check if we have risk mitigation data
    has_risk_data = bool(risk_section and risk_section.get('main_strategy'))
    print(f"[DEBUG] Risk mitigation data available: {has_risk_data}")
    print(f"[DEBUG] Risk section keys: {list(risk_section.keys()) if risk_section else 'None'}")
    
    risk_title = risk_section.get('title', 'Risk Mitigation Strategies')
    add_clean_text(slide, Inches(8), Inches(1.4), Inches(5), Inches(0.3), 
                   risk_title, 14, colors["primary"], True)
    
    # Main strategy box (no shadow) - Enhanced with fallback content
    if has_risk_data:
        raw_main_strategy = risk_section.get('main_strategy', {})
        
        # FIXED: Handle both string and object formats for main_strategy
        if isinstance(raw_main_strategy, str):
            # String format: convert to object structure
            main_strategy = {
                'title': 'Risk Mitigation Strategy',
                'description': raw_main_strategy,
                'benefits': ['Operational flexibility', 'Market responsiveness', 'Cost adaptability']
            }
            print(f"[DEBUG] Converted string main_strategy to object format")
        elif isinstance(raw_main_strategy, dict):
            # Already object format
            main_strategy = raw_main_strategy
            print(f"[DEBUG] Using provided risk mitigation data in object format")
        else:
            # Fallback for other types
            main_strategy = {
                'title': 'Risk Management Approach',
                'description': str(raw_main_strategy) if raw_main_strategy else 'Comprehensive risk management strategies in place',
                'benefits': ['Risk monitoring', 'Adaptive strategies', 'Proactive management']
            }
            print(f"[DEBUG] Converted unknown main_strategy type to object format")
    else:
        # Provide comprehensive fallback when no data available
        main_strategy = {
            'title': 'Diversified Revenue Strategy',
            'description': 'Multiple revenue streams and operational efficiency measures provide resilience against market volatility and cost pressures',
            'benefits': ['Reduced earnings volatility', 'Stable cash generation', 'Operational flexibility']
        }
        print(f"[DEBUG] Using fallback risk mitigation strategy")
    
    strategy_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.8), Inches(5), Inches(1.8))
    strategy_box.fill.solid()
    strategy_box.fill.fore_color.rgb = colors["light_grey"]
    strategy_box.line.fill.background()
    strategy_box.shadow.inherit = False  # Remove shadow
    
    # Gold stripe (no shadow)
    gold_stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.8), Inches(0.1), Inches(1.8))
    gold_stripe.fill.solid()
    gold_stripe.fill.fore_color.rgb = colors["secondary"]
    gold_stripe.line.fill.background()
    gold_stripe.shadow.inherit = False  # Remove shadow
    
    # Strategy content
    strategy_title = main_strategy.get('title', 'Revenue Diversification Strategy')
    add_clean_text(slide, Inches(8.3), Inches(1.9), Inches(4.5), Inches(0.25), 
                   strategy_title, 12, colors["secondary"], True)
    
    strategy_desc = main_strategy.get('description', 'Multiple revenue streams provide stability')
    add_clean_text(slide, Inches(8.3), Inches(2.15), Inches(4.5), Inches(0.35), 
                   strategy_desc, 10, colors["text"])
    
    # Benefits
    benefits_title = "Key Benefits:"
    add_clean_text(slide, Inches(8.3), Inches(2.55), Inches(4.5), Inches(0.18), 
                   benefits_title, 10, colors["primary"], True)
    
    benefits = main_strategy.get('benefits', ['Lower volatility', 'Stable cash flows', 'Reduced risk'])
    for i, benefit in enumerate(benefits[:3]):  # Max 3 benefits
        add_clean_text(slide, Inches(8.5), Inches(2.75 + i * 0.13), Inches(4.3), Inches(0.11), 
                       f"• {benefit}", 9, colors["text"])
    
    # Banker's view box (no shadow) - Enhanced fallback
    if has_risk_data and 'banker_view' in risk_section:
        banker_view = risk_section.get('banker_view', {})
        print(f"[DEBUG] Using provided banker's view")
    else:
        # Enhanced fallback banker's view
        banker_view = {
            'title': "BANKER'S VIEW",
            'text': 'Strong operational resilience through diversified revenue streams and disciplined cost management supports sustainable margin expansion and reduced business risk profile.'
        }
        print(f"[DEBUG] Using fallback banker's view")
    
    banker_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(3.8), Inches(5), Inches(0.8))
    banker_box.fill.solid()
    banker_box.fill.fore_color.rgb = RGBColor(240, 255, 240)  # Light green
    banker_box.line.fill.background()
    banker_box.shadow.inherit = False  # Remove shadow
    
    banker_title = banker_view.get('title', "BANKER'S VIEW")
    add_clean_text(slide, Inches(8.2), Inches(3.9), Inches(4.6), Inches(0.15), 
                   banker_title, 10, RGBColor(34, 139, 34), True)
    
    banker_text = banker_view.get('text', 'Strong margin expansion demonstrates operational excellence.')
    add_clean_text(slide, Inches(8.2), Inches(4.05), Inches(4.6), Inches(0.5), 
                   banker_text, 9, colors["text"])
    
    # Bottom summary - FIXED POSITIONING to avoid overlap
    summary = slide_data.get('summary', '')
    if summary:
        # Moved summary higher to avoid overlapping with footer
        add_clean_text(slide, Inches(1), Inches(6.2), Inches(12), Inches(0.5), 
                       summary, 10, colors["primary"])
    
    # Get today's date
    today = datetime.now().strftime("%B %d, %Y")
    
    # Add footer - "Confidential | [today's date]" on LEFT (light grey) - FIXED POSITION
    footer_left = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(6), Inches(0.4))
    footer_left_frame = footer_left.text_frame
    footer_left_frame.clear()
    p = footer_left_frame.paragraphs[0]
    p.text = f"Confidential | {today}"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.LEFT
    footer_left_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add footer - Company name on RIGHT - FIXED POSITION
    footer_right = slide.shapes.add_textbox(Inches(10), Inches(6.9), Inches(3), Inches(0.4))
    footer_right_frame = footer_right.text_frame
    footer_right_frame.clear()
    p = footer_right_frame.paragraphs[0]
    p.text = company_name or "Moelis"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.RIGHT
    footer_right_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return prs


def render_historical_financial_performance_slide(data=None, color_scheme=None, typography=None, company_name="Moelis", prs=None, brand_config=None, **kwargs):
    """
    Renders a historical financial performance slide with chart and metrics
    """
    
    # Create or use existing presentation
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs = ensure_prs(prs)
    
    # Get brand styling
    colors, fonts = get_brand_styling(brand_config, color_scheme, typography)
    
    # Add blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Helper function to add clean text
    def add_clean_text(slide, left, top, width, height, text, font_size=10, 
                       color=colors["text"], bold=False, align=PP_ALIGN.LEFT, bg_color=None):
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = align
            for run in paragraph.runs:
                run.font.name = fonts["primary_font"]
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.bold = bold
        
        if bg_color:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = bg_color
        
        # Remove shadow/outline
        textbox.line.fill.background()
        textbox.shadow.inherit = False
        
        return textbox
    
    # Set white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors["background"]
    
    # STANDARDIZED: Apply header and title
    title_text = (data or {}).get('title', 'Historical Financial Performance (I)')
    _apply_standard_header_and_title(slide, title_text, brand_config, company_name)
    
    # Main chart title
    chart_info = (data or {}).get('chart', {})
    chart_title = chart_info.get('title', 'Company - 5-Year Financial Performance')
    add_clean_text(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.3), 
                   chart_title, 16, colors["primary"], True, PP_ALIGN.CENTER)
    
    # Create combination chart
    chart_data = ChartData()
    categories = chart_info.get('categories', ['2020', '2021', '2022', '2023', '2024'])
    revenue_data = chart_info.get('revenue', [26, 24, 33, 40, 42])
    ebitda_data = chart_info.get('ebitda', [6.5, 5.8, 8.2, 10.0, 11.2])
    
    # DEBUG: Print actual chart data being used
    print(f"[DEBUG] Historical chart categories: {categories}")
    print(f"[DEBUG] Historical chart revenue: {revenue_data}")
    print(f"[DEBUG] Historical chart ebitda: {ebitda_data}")
    
    # Ensure data is numeric (convert strings to numbers if needed)
    try:
        revenue_data = [float(x) for x in revenue_data]
        ebitda_data = [float(x) for x in ebitda_data] 
        print(f"[DEBUG] Successfully converted chart data to numeric")
    except (ValueError, TypeError) as e:
        print(f"[DEBUG] Error converting chart data to numeric: {e}")
        # Keep original data if conversion fails
        pass
    
    chart_data.categories = categories
    chart_data.add_series('Revenue (USD millions)', revenue_data)
    chart_data.add_series('EBITDA (USD millions)', ebitda_data)
    
    # Add chart
    chart_left = Inches(2)
    chart_top = Inches(1.7)
    chart_width = Inches(9)
    chart_height = Inches(2.3)
    
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, chart_left, chart_top, chart_width, chart_height, chart_data
    )
    
    chart = chart_shape.chart
    
    # Style the chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.font.size = Pt(10)
    
    # Style chart elements
    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    category_axis.tick_labels.font.size = Pt(10)
    category_axis.tick_labels.font.name = fonts["primary_font"]
    
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.tick_labels.font.size = Pt(10)
    value_axis.tick_labels.font.name = fonts["primary_font"]
    
    # FIXED: Dynamic axis scaling based on actual data values
    all_values = revenue_data + ebitda_data
    if all_values:
        data_min = min(all_values)
        data_max = max(all_values)
        
        # Set axis maximum to 120% of highest value for proper display
        dynamic_max = int(data_max * 1.2)
        # Ensure minimum is 0 for proper bar display (important for large numbers)
        dynamic_min = 0
        
        value_axis.minimum_scale = dynamic_min
        value_axis.maximum_scale = dynamic_max
        
        print(f"[DEBUG] Chart axis scaled: range {data_min:,} to {data_max:,} → axis 0 to {dynamic_max:,}")
        
        # For very large numbers, ensure PowerPoint can handle the scale
        if dynamic_max > 1000000:  # Over 1 million
            print(f"[DEBUG] Large scale detected, ensuring proper chart rendering")
    else:
        value_axis.minimum_scale = 0
        value_axis.maximum_scale = 45  # Fallback for empty data
    
    # Color the series
    series = chart.series
    
    # Revenue series (primary color)
    revenue_series = series[0]
    revenue_series.format.fill.solid()
    revenue_series.format.fill.fore_color.rgb = colors["primary"]
    for point in revenue_series.points:
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colors["primary"]
    
    # EBITDA series (secondary color)
    ebitda_series = series[1]
    ebitda_series.format.fill.solid()
    ebitda_series.format.fill.fore_color.rgb = colors["secondary"]
    for point in ebitda_series.points:
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colors["secondary"]
    
    # Chart footnote
    chart_footnote = chart_info.get('footnote', '*Historical figures represent estimated performance based on market trends.')
    add_clean_text(slide, Inches(2), Inches(4.1), Inches(9), Inches(0.2), 
                   chart_footnote, 8, colors["text"], False, PP_ALIGN.CENTER)
    
    # Key metrics section (removed shadows from boxes)
    metrics_section = (data or {}).get('key_metrics', {})
    metrics_y = Inches(4.4)
    metrics = metrics_section.get('metrics', [])
    
    # If no metrics provided, create default ones
    if not metrics:
        metrics = [
            {
                'title': 'Patient Growth (CAGR)',
                'value': '12.4%',
                'period': '(2020-2024)',
                'note': '✓ Consistent growth despite pandemic disruptions'
            },
            {
                'title': 'Patient Retention Rate',
                'value': '87%',
                'period': '(2024)',
                'note': '✓ Premium market segment leading indicator'
            },
            {
                'title': 'Avg. Revenue Per Patient',
                'value': '$980',
                'period': 'USD (2024)',
                'note': '↗ +8.2% increase from 2023'
            },
            {
                'title': 'Corporate Contracts',
                'value': '35+',
                'period': '(2024)',
                'note': '● Major financial institutions & MNCs'
            }
        ]
    
    # Calculate positions to fit exactly 4 boxes across slide width
    slide_content_width = Inches(12.5)  # Total usable width
    box_width = Inches(2.8)             # Width of each box
    total_boxes_width = box_width * 4   # Total width of all boxes
    spacing_width = slide_content_width - total_boxes_width  # Remaining space
    box_spacing = spacing_width / 5     # Space between boxes (5 gaps: before, between 3, after)
    
    for i, metric in enumerate(metrics[:4]):  # Max 4 metrics
        x_pos = Inches(0.5) + box_spacing + i * (box_width + box_spacing)
        
        # Metric box (NO SHADOW)
        metric_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, metrics_y, box_width, Inches(1.1))
        metric_bg.fill.solid()
        metric_bg.fill.fore_color.rgb = colors["light_grey"]
        metric_bg.line.fill.background()
        metric_bg.shadow.inherit = False  # Remove shadow
        
        # FIXED: Handle both string and object metric formats
        if isinstance(metric, str):
            # String format: display as simple text block with better visibility
            print(f"[DEBUG] Converting string metric to display format: {metric}")
            add_clean_text(slide, x_pos + Inches(0.05), metrics_y + Inches(0.2), box_width - Inches(0.1), Inches(0.7), 
                           metric, 11, colors["primary"], True)  # Larger font, bold, primary color for visibility
        elif isinstance(metric, dict):
            # Object format: use structured layout
            add_clean_text(slide, x_pos + Inches(0.1), metrics_y + Inches(0.1), box_width - Inches(0.2), Inches(0.2), 
                           metric.get('title', ''), 10, colors["text"], True)
            add_clean_text(slide, x_pos + Inches(0.1), metrics_y + Inches(0.3), box_width - Inches(0.2), Inches(0.25), 
                           metric.get('value', ''), 18, colors["primary"], True)
            add_clean_text(slide, x_pos + Inches(0.1), metrics_y + Inches(0.55), box_width - Inches(0.2), Inches(0.15), 
                           metric.get('period', ''), 9, colors["text"])
            add_clean_text(slide, x_pos + Inches(0.1), metrics_y + Inches(0.75), box_width - Inches(0.2), Inches(0.25), 
                           metric.get('note', ''), 8, RGBColor(34, 139, 34))
        else:
            # Fallback for other types
            print(f"[DEBUG] Unknown metric format: {type(metric)}")
            add_clean_text(slide, x_pos + Inches(0.1), metrics_y + Inches(0.4), box_width - Inches(0.2), Inches(0.3), 
                           str(metric), 10, colors["text"], False)
    
    # Revenue Growth section
    revenue_section = (data or {}).get('revenue_growth', {})
    covid_y = Inches(5.7)
    section_title = revenue_section.get('title', 'Revenue Growth')
    add_clean_text(slide, Inches(1), covid_y, Inches(7), Inches(0.2), 
                   section_title, 12, colors["primary"], True)
    
    # Five bullet points for revenue growth
    revenue_points = revenue_section.get('points', [])
    for i, point in enumerate(revenue_points[:5]):  # Max 5 points
        add_clean_text(slide, Inches(1), covid_y + Inches(0.25 + i * 0.18), Inches(7), Inches(0.16), 
                       f"● {point}", 9, colors["text"])
    
    # Banker's view section (NO SHADOW) - EXPANDED COVERAGE
    banker_view = (data or {}).get('banker_view', {})
    banker_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.2), covid_y - Inches(0.1), Inches(4.8), Inches(1.0))
    banker_box.fill.solid()
    banker_box.fill.fore_color.rgb = RGBColor(240, 255, 240)  # Light green
    banker_box.line.fill.background()  # Remove outline
    banker_box.shadow.inherit = False  # Remove shadow
    
    add_clean_text(slide, Inches(8.4), covid_y, Inches(4.4), Inches(0.15), 
                   banker_view.get('title', "BANKER'S VIEW"), 10, RGBColor(34, 139, 34), True)
    add_clean_text(slide, Inches(8.4), covid_y + Inches(0.2), Inches(4.4), Inches(0.65), 
                   banker_view.get('text', ''), 9, colors["text"])
    
    # Get today's date
    today = datetime.now().strftime("%B %d, %Y")
    
    # Footer
    add_clean_text(slide, Inches(0.5), Inches(6.9), Inches(4), Inches(0.2), 
                   f"Confidential | {today}", 9, colors["footer_grey"])
    
    add_clean_text(slide, Inches(9.5), Inches(6.9), Inches(3.5), Inches(0.2), 
                   f"{company_name} Investment Opportunity    6", 9, colors["footer_grey"], False, PP_ALIGN.RIGHT)
    
    return prs


def render_business_overview_slide(data=None, color_scheme=None, typography=None, company_name="Moelis", prs=None, brand_config=None, **kwargs):
    """
    Render a business & operational overview slide for investment banking presentations
    """
    
    # Create presentation if not provided (standard 16:9 dimensions)
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # Standard 16:9 width
        prs.slide_height = Inches(7.5)    # Standard 16:9 height
    else:
        prs = ensure_prs(prs)
    
    # Get brand styling
    colors, fonts = get_brand_styling(brand_config, color_scheme, typography)
    
    # Add slide with blank layout
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set clean white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors["background"]
    
    def add_clean_text(slide, left, top, width, height, text, font_size=14, 
                       color=None, bold=False, align=PP_ALIGN.LEFT, bg_color=None):
        """Add text with consistent sizing and NO OUTLINES"""
        if color is None:
            color = colors["text"]
            
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.word_wrap = True
        
        # Style the text consistently
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = align
            for run in paragraph.runs:
                run.font.name = fonts["primary_font"]
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.bold = bold
        
        # Background (if needed) - NO OUTLINES
        if bg_color:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = bg_color
        
        # Remove all borders/outlines
        textbox.line.fill.background()
        textbox.shadow.inherit = False
        
        return textbox
    
    # Debug logging
    print(f"[DEBUG] Business overview data keys: {list((data or {}).keys())}")
    
    # Extract slide data - handle both data formats
    slide_data = data or {}
    
    # STANDARDIZED: Apply header and title
    title_text = slide_data.get('title', 'Business & Operational Overview')
    _apply_standard_header_and_title(slide, title_text, brand_config, company_name)
    
    # Company description - FIXED POSITIONING TO AVOID OVERLAP
    company_desc = slide_data.get('description', 'Leading healthcare services provider with comprehensive medical care and operational excellence.')
    print(f"[DEBUG] Company description: {company_desc}")
    
    # FIXED: Reduced width to make room for highlights box and proper text wrapping
    add_clean_text(slide, Inches(0.8), Inches(1.3), Inches(7.0), Inches(1.2), 
                   company_desc, 12, colors["text"])
    
    # Timeline elements - FIXED POSITIONING AND SPACING
    timeline_data = slide_data.get('timeline', {
        'start_year': '2015',
        'end_year': '2024', 
        'years_note': '(9+ years of operation)'
    })
    
    timeline_y = Inches(2.6)  # Moved down slightly to accommodate longer description text
    
    try:
        # Start year circle
        start_year = timeline_data.get('start_year', '2015')
        circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), timeline_y, Inches(0.12), Inches(0.12))
        circle1.fill.solid()
        circle1.fill.fore_color.rgb = colors["secondary"]
        circle1.line.fill.background()
        
        # Timeline line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.12), timeline_y + Inches(0.05), 
                                      Inches(4), Inches(0.02))  # Shorter line
        line.fill.solid()
        line.fill.fore_color.rgb = colors["secondary"]
        line.line.fill.background()
        
        # End year circle
        end_year = timeline_data.get('end_year', '2024')
        circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), timeline_y, Inches(0.12), Inches(0.12))
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = colors["secondary"]
        circle2.line.fill.background()
        
        # Timeline labels - BETTER POSITIONING
        add_clean_text(slide, Inches(0.9), timeline_y - Inches(0.3), Inches(0.5), Inches(0.2), 
                       start_year, 11, colors["primary"], True, PP_ALIGN.CENTER)
        
        add_clean_text(slide, Inches(4.9), timeline_y - Inches(0.3), Inches(0.5), Inches(0.2), 
                       end_year, 11, colors["primary"], True, PP_ALIGN.CENTER)
        
        years_operation = timeline_data.get('years_note', '(9+ years of growth and expansion)')
        add_clean_text(slide, Inches(5.4), timeline_y - Inches(0.1), Inches(3), Inches(0.3), 
                       years_operation, 9, colors["text"])
    except Exception as e:
        print(f"[DEBUG] Timeline creation error: {e}")
    
    # Operational Highlights box - COMPLETELY REPOSITIONED TO AVOID OVERLAP
    try:
        highlights_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.0), Inches(1.3), 
                                               Inches(4.8), Inches(5.8))  # Full height utilization
        highlights_bg.fill.solid()
        highlights_bg.fill.fore_color.rgb = colors["light_grey"]
        highlights_bg.line.fill.background()
        highlights_bg.shadow.inherit = False
        
        # Highlights title
        highlights_title = slide_data.get('highlights_title', 'Key Operational Highlights')
        add_clean_text(slide, Inches(8.2), Inches(1.4), Inches(4.4), Inches(0.3), 
                       highlights_title, 12, colors["primary"], True, PP_ALIGN.CENTER)
        
        # Enhanced highlight items with more detail - SUPPORT FOR RICH CONTENT
        highlights = slide_data.get('highlights', [
            'Market-leading position with 130+ outlets across major Indonesian cities',
            'Strong growth trajectory: 100+ new locations since 2021 acquisition', 
            'Consistent double-digit revenue and EBITDA growth post-acquisition',
            'Premium Japanese-inspired artisan bakery positioning',
            'Diversified revenue streams: B2C retail + institutional B2B clients',
            'Strategic locations in high-traffic transit hubs and commercial centers',
            'Fresh, preservative-free product differentiation',
            'Proven scalable business model with operational excellence'
        ])
        
        # Support more highlights with better spacing
        for i, item in enumerate(highlights[:8]):  # Support up to 8 detailed items
            y_pos = Inches(1.75 + i * 0.6)  # Tighter but readable spacing
            
            # Gold bullet - aligned with text baseline
            bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.3), y_pos + Inches(0.08), Inches(0.04), Inches(0.04))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = colors["secondary"]
            bullet.line.fill.background()
            
            # Text with better wrapping for detailed content
            add_clean_text(slide, Inches(8.42), y_pos, Inches(4.25), Inches(0.55), 
                           item, 9, colors["text"])
    except Exception as e:
        print(f"[DEBUG] Highlights section error: {e}")
    
    # Enhanced Service Lines section with richer content support
    try:
        services_title = slide_data.get('services_title', 'Core Business Lines & Capabilities')
        add_clean_text(slide, Inches(0.8), Inches(3.0), Inches(6.5), Inches(0.3), 
                       services_title, 12, colors["primary"], True)
        
        # Enhanced service items with more descriptive content
        services = slide_data.get('services', [
            'Artisan Sweet & Savory Breads: Premium Japanese-inspired baked goods',
            'Whole Cakes & Celebration Products: Custom orders and special occasions',
            'Fresh Food & Light Meals: Sandwiches, salads, and ready-to-eat items', 
            'Coffee & Beverages: Premium coffee program and specialty drinks',
            'Corporate Catering: B2B institutional sales and corporate partnerships',
            'Retail Merchandise: Branded items and gift products',
            'Franchise Operations: Proven business model for expansion partners',
            'Digital Ordering Platform: Online ordering and delivery capabilities'
        ])
        
        # Enhanced layout supporting more detailed service descriptions
        service_cols = 2  # Split into 2 columns for better readability
        items_per_col = (len(services) + 1) // 2
        
        for i, service in enumerate(services[:8]):  # Support up to 8 detailed services
            if i < items_per_col:  # Left column
                x_pos = Inches(0.9)
                y_pos = Inches(3.4 + i * 0.5)  # Proper spacing for detailed text
            else:  # Right column
                x_pos = Inches(4.2)  # Optimized column positioning
                y_pos = Inches(3.4 + (i - items_per_col) * 0.5)
            
            # Gold bullet - aligned with text baseline
            bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, x_pos, y_pos + Inches(0.08), Inches(0.04), Inches(0.04))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = colors["secondary"]
            bullet.line.fill.background()
            
            # Text with proper width for detailed service descriptions
            add_clean_text(slide, x_pos + Inches(0.12), y_pos, Inches(3.5), Inches(0.45), 
                           service, 9, colors["text"])
    except Exception as e:
        print(f"[DEBUG] Services section error: {e}")
    
    # Strategic Positioning section - REPOSITIONED AND RESIZED
    try:
        positioning_title = slide_data.get('positioning_title', 'Strategic Market Positioning')
        add_clean_text(slide, Inches(0.8), Inches(5.8), Inches(7), Inches(0.3), 
                       positioning_title, 12, colors["primary"], True)
        
        positioning_desc = slide_data.get('positioning_desc', 
            'The company has established itself as the leading premium healthcare provider in Southeast Asia, '
            'serving both individual patients and corporate clients with comprehensive medical services and exceptional care standards.')
        
        # FIXED: Reduced width and increased height for better text flow
        add_clean_text(slide, Inches(0.8), Inches(6.1), Inches(7), Inches(0.8), 
                       positioning_desc, 10, colors["text"])
    except Exception as e:
        print(f"[DEBUG] Positioning section error: {e}")
    
    # Footer with proper formatting - moved down to accommodate content
    footer_top = Inches(7.0)
    current_date = datetime.now().strftime("%B %d, %Y")
    
    # Left footer - Confidential with date
    footer_left = slide.shapes.add_textbox(Inches(0.8), footer_top, Inches(4), Inches(0.2))
    footer_left_frame = footer_left.text_frame
    footer_left_p = footer_left_frame.paragraphs[0]
    footer_left_p.text = f"Confidential | {current_date}"
    footer_left_p.alignment = PP_ALIGN.LEFT
    footer_left_run = footer_left_p.runs[0]
    footer_left_run.font.name = fonts["primary_font"]
    footer_left_run.font.size = fonts["small_size"]
    footer_left_run.font.color.rgb = colors["footer_grey"]
    
    # Right footer - Company name
    footer_right = slide.shapes.add_textbox(Inches(9.5), footer_top, Inches(3.5), Inches(0.2))
    footer_right_frame = footer_right.text_frame
    footer_right_p = footer_right_frame.paragraphs[0]
    footer_right_p.text = f"{company_name}"
    footer_right_p.alignment = PP_ALIGN.RIGHT
    footer_right_run = footer_right_p.runs[0]
    footer_right_run.font.name = fonts["primary_font"]
    footer_right_run.font.size = fonts["small_size"]
    footer_right_run.font.color.rgb = colors["footer_grey"]
    
    return prs


def render_precedent_transactions_slide(data=None, color_scheme=None, typography=None, company_name="Moelis", prs=None, brand_config=None, **kwargs):
    """
    Simple, reliable precedent transactions slide
    """
    
    # FIXED: Extract slide_data from data parameter
    slide_data = data or {}
    
    # Create presentation
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs = ensure_prs(prs)
    
    # Get brand styling
    colors, fonts = get_brand_styling(brand_config, color_scheme, typography)
    
    # Add slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Extract data
    transactions = slide_data.get('transactions', [])
    
    print(f"[DEBUG] Precedent transactions: Found {len(transactions)} transactions")
    
    if not transactions:
        # STANDARDIZED: Apply header and title even if no data
        title_text = slide_data.get('title', 'Precedent Transactions Analysis')
        _apply_standard_header_and_title(slide, title_text, brand_config, company_name)
        
        # Add placeholder message
        placeholder = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(8), Inches(1))
        placeholder_frame = placeholder.text_frame
        placeholder_frame.text = "Transaction analysis will be displayed here when data is available."
        placeholder_para = placeholder_frame.paragraphs[0]
        placeholder_para.font.name = fonts["primary_font"]
        placeholder_para.font.size = Pt(14)
        placeholder_para.alignment = PP_ALIGN.CENTER
        
        return prs
    
    # STANDARDIZED: Apply header and title
    title_text = slide_data.get('title', 'Precedent Transactions Analysis')
    _apply_standard_header_and_title(slide, title_text, brand_config, company_name)
    
    # For now, let's create a simple visual representation using text boxes
    # This eliminates matplotlib corruption issues
    
    # Chart area placeholder
    chart_title = slide.shapes.add_textbox(Inches(2.0), Inches(1.5), Inches(9.0), Inches(0.5))
    chart_title_frame = chart_title.text_frame
    chart_title_frame.text = "EV/Revenue Multiples by Transaction"
    chart_title_para = chart_title_frame.paragraphs[0]
    chart_title_para.font.name = fonts["primary_font"]
    chart_title_para.font.size = Pt(14)
    chart_title_para.font.color.rgb = colors["primary"]
    chart_title_para.font.bold = True
    chart_title_para.alignment = PP_ALIGN.CENTER
    
    # Simple bar representation using rectangles - ADJUSTED POSITIONING
    num_transactions = len(transactions)
    bar_area_left = Inches(1.75)  # Moved slightly left
    bar_area_width = Inches(9.5)  # Increased width
    bar_width = bar_area_width / num_transactions
    bar_top = Inches(2.3)  # Moved down to avoid title overlap
    
    for i, transaction in enumerate(transactions):
        multiple_raw = transaction.get('ev_revenue_multiple', 0)
        
        # Convert multiple to float (handle string formats like "50x", "25.5x", etc.)
        try:
            if isinstance(multiple_raw, str):
                # Remove 'x' and convert to float
                multiple = float(multiple_raw.replace('x', '').replace('X', '').strip())
            else:
                multiple = float(multiple_raw) if multiple_raw else 0
        except (ValueError, TypeError):
            multiple = 0
            
        # Scale bar height properly - max 1.5 inches total
        max_bar_height = Inches(1.5)
        if multiple > 0:
            # Find the max multiple to scale properly
            max_multiple = max([
                float(str(t.get('ev_revenue_multiple', '0')).replace('x', '').replace('X', '').strip() or '0')
                for t in transactions
            ])
            if max_multiple > 0:
                # Scale relative to max, with maximum height of 1.5 inches
                bar_height = Inches((multiple / max_multiple) * 1.5)
            else:
                bar_height = Inches(0.1)
        else:
            bar_height = Inches(0.1)
        
        bar_left = bar_area_left + (bar_width * i) + Inches(0.05)  # Small margin
        bar_actual_width = bar_width - Inches(0.1)  # Space between bars
        
        # Create bar rectangle
        bar_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            bar_left, bar_top + Inches(1.5) - bar_height, bar_actual_width, bar_height
        )
        bar_shape.fill.solid()
        bar_shape.fill.fore_color.rgb = colors["primary"]
        bar_shape.line.fill.background()
        
        # Add value label
        if multiple > 0:
            label_shape = slide.shapes.add_textbox(
                bar_left, bar_top + Inches(1.5) - bar_height - Inches(0.25), 
                bar_actual_width, Inches(0.2)
            )
            label_frame = label_shape.text_frame
            # Safe formatting for the multiple value
            try:
                label_frame.text = f"{multiple:.1f}x"
            except (ValueError, TypeError):
                label_frame.text = f"{multiple}x"
            label_para = label_frame.paragraphs[0]
            label_para.font.name = fonts["primary_font"]
            label_para.font.size = Pt(9)
            label_para.font.bold = True
            label_para.alignment = PP_ALIGN.CENTER
        
        # Add T1, T2, etc. labels
        t_label_shape = slide.shapes.add_textbox(
            bar_left, bar_top + Inches(1.6), bar_actual_width, Inches(0.2)
        )
        t_label_frame = t_label_shape.text_frame
        t_label_frame.text = f"T{i+1}"
        t_label_para = t_label_frame.paragraphs[0]
        t_label_para.font.name = fonts["primary_font"]
        t_label_para.font.size = Pt(9)
        t_label_para.alignment = PP_ALIGN.CENTER
    
    # Create table - REPOSITIONED TO AVOID OVERLAP
    row_labels = ['Date', 'Target', 'Acquirer', 'Country', 'EV ($M)', 'Revenue ($M)', 'EV/Revenue']
    num_rows = len(row_labels)
    
    # Table positioning - moved down to avoid chart overlap
    labels_left = Inches(0.4)
    labels_width = Inches(1.3)
    table_left = Inches(1.75)  # Aligned with chart area
    table_top = Inches(4.8)  # Moved down further to avoid overlap
    table_width = Inches(9.5)  # Match chart width
    row_height = Inches(0.32)  # Slightly increased for better readability
    
    # Create row labels table
    labels_table = slide.shapes.add_table(
        num_rows, 1, 
        labels_left, table_top, 
        labels_width, row_height * num_rows
    ).table
    
    # Create data table
    data_table = slide.shapes.add_table(
        num_rows, num_transactions, 
        table_left, table_top, 
        table_width, row_height * num_rows
    ).table
    
    # Set column widths for better distribution - equal width for all columns
    column_width = table_width / num_transactions
    for col_idx in range(num_transactions):
        data_table.columns[col_idx].width = int(column_width)
    
    # Format row labels
    for i, label in enumerate(row_labels):
        cell = labels_table.cell(i, 0)
        cell.text = label
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors["light_grey"]
        
        para = cell.text_frame.paragraphs[0]
        para.font.name = fonts["primary_font"]
        para.font.size = Pt(12)
        para.font.color.rgb = colors["text"]
        para.font.bold = True
        para.alignment = PP_ALIGN.RIGHT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Format data table
    for col_idx, transaction in enumerate(transactions):
        target = transaction.get('target', 'N/A')
        acquirer = transaction.get('acquirer', 'N/A')
        
        # Smart truncation for company names - keep more characters but break intelligently
        def smart_truncate(name, max_length=18):
            """Intelligently truncate company names"""
            if len(name) <= max_length:
                return name
            
            # Try to break at word boundaries
            words = name.split()
            if len(words) > 1:
                # Keep first word if it's meaningful
                if len(words[0]) <= max_length - 3:
                    return words[0] + '...'
            
            # Otherwise, simple truncation
            return name[:max_length-3] + '...'
        
        target = smart_truncate(target)
        acquirer = smart_truncate(acquirer)
            
        # Convert financial values to proper format with numeric handling
        def convert_financial_value(value, prefix="$"):
            """Convert financial string values like '$10B', '$5.2M' to compact format"""
            if not value or value == 'N/A':
                return 'N/A'
            
            try:
                if isinstance(value, str):
                    # Remove prefix symbols and spaces
                    clean_value = value.replace('$', '').replace(',', '').strip()
                    
                    # Handle suffixes like B, M, K - keep original notation
                    if clean_value.upper().endswith('B'):
                        return f"{prefix}{clean_value[:-1]}B"
                    elif clean_value.upper().endswith('M'):
                        return f"{prefix}{clean_value[:-1]}M"
                    elif clean_value.upper().endswith('K'):
                        return f"{prefix}{clean_value[:-1]}K"
                    else:
                        # For raw numbers, convert to compact notation
                        numeric_value = float(clean_value)
                        if numeric_value >= 1_000_000_000:
                            return f"{prefix}{numeric_value/1_000_000_000:.1f}B"
                        elif numeric_value >= 1_000_000:
                            return f"{prefix}{numeric_value/1_000_000:.0f}M"
                        elif numeric_value >= 1_000:
                            return f"{prefix}{numeric_value/1_000:.0f}K"
                        else:
                            return f"{prefix}{numeric_value:,.0f}"
                else:
                    # Already numeric - convert to compact notation
                    numeric_value = float(value)
                    if numeric_value >= 1_000_000_000:
                        return f"{prefix}{numeric_value/1_000_000_000:.1f}B"
                    elif numeric_value >= 1_000_000:
                        return f"{prefix}{numeric_value/1_000_000:.0f}M"
                    elif numeric_value >= 1_000:
                        return f"{prefix}{numeric_value/1_000:.0f}K"
                    else:
                        return f"{prefix}{numeric_value:,.0f}"
            except (ValueError, TypeError):
                return 'N/A'
        
        def convert_multiple_value(value):
            """Convert multiple string values like '50x', '25.5x' to formatted string"""
            if not value:
                return 'N/A'
            
            try:
                if isinstance(value, str):
                    # Remove 'x' and convert to float
                    clean_value = value.replace('x', '').replace('X', '').strip()
                    numeric_value = float(clean_value)
                    return f"{numeric_value:.1f}x"
                else:
                    return f"{float(value):.1f}x"
            except (ValueError, TypeError):
                return 'N/A'
        
        data_values = [
            transaction.get('date', 'N/A'),
            target,
            acquirer,
            transaction.get('country', 'N/A'),
            convert_financial_value(transaction.get('enterprise_value')),
            convert_financial_value(transaction.get('revenue')),
            convert_multiple_value(transaction.get('ev_revenue_multiple'))
        ]
        
        for row_idx, value in enumerate(data_values):
            cell = data_table.cell(row_idx, col_idx)
            cell.text = str(value)
            
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors["background"]
            
            para = cell.text_frame.paragraphs[0]
            para.font.name = fonts["primary_font"]
            para.font.size = Pt(10)  # Increased from 9 to 10 for better readability
            para.font.color.rgb = colors["text"]
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Adjust cell margins for better spacing
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
    
    # Add footer
    footer_left = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(6), Inches(0.4))
    footer_left_frame = footer_left.text_frame
    footer_left_frame.text = f"Confidential | {datetime.now().strftime('%B %Y')}"
    footer_left_para = footer_left_frame.paragraphs[0]
    footer_left_para.font.name = fonts["primary_font"]
    footer_left_para.font.size = fonts["small_size"]
    footer_left_para.font.color.rgb = colors["footer_grey"]
    
    footer_right = slide.shapes.add_textbox(Inches(7.333), Inches(7.0), Inches(6), Inches(0.4))
    footer_right_frame = footer_right.text_frame
    footer_right_frame.text = company_name
    footer_right_para = footer_right_frame.paragraphs[0]
    footer_right_para.font.name = fonts["primary_font"]
    footer_right_para.font.size = fonts["small_size"]
    footer_right_para.font.color.rgb = colors["footer_grey"]
    footer_right_para.alignment = PP_ALIGN.RIGHT
    
    return prs


def render_valuation_overview_slide(data=None, color_scheme=None, typography=None, company_name="Moelis", prs=None, brand_config=None, **kwargs):
    """
    Renders a clean, simple valuation overview slide with a standard table format
    """
    
    # Create presentation
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs = ensure_prs(prs)
    
    # Get brand styling
    colors, fonts = get_brand_styling(brand_config, color_scheme, typography)
    
    # Add slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Extract data
    if data is None:
        data = {}
    
    title_text = data.get('title', 'Valuation Overview')
    subtitle_text = data.get('subtitle', 'Implied EV/Post IRFS-16 EBITDA')
    valuation_data = data.get('valuation_data', [])
    
    print(f"[DEBUG] Valuation slide - {len(valuation_data)} data rows")
    
    # STANDARDIZED: Apply header and title
    _apply_standard_header_and_title(slide, title_text, brand_config, company_name)
    
    # Add subtitle header
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(11), Inches(0.3))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle_text
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.name = fonts["primary_font"] 
    subtitle_para.font.size = Pt(14)
    subtitle_para.font.bold = True
    subtitle_para.font.color.rgb = colors["primary"]
    
    if not valuation_data:
        # Add message about missing data
        message_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(1))
        message_frame = message_box.text_frame
        message_frame.text = "Valuation data will be displayed here when available."
        message_para = message_frame.paragraphs[0]
        message_para.alignment = PP_ALIGN.CENTER
        message_para.font.name = fonts["primary_font"]
        message_para.font.size = Pt(14)
        message_para.font.color.rgb = colors["text"]
        
        return prs
    
    # Create simple table
    rows = len(valuation_data) + 1  # +1 for header
    cols = 5  # Methodology, Commentary, Enterprise Value, Metric, 22A'/23E
    
    table_left = Inches(0.5)
    table_top = Inches(1.8) 
    table_width = Inches(12.3)
    table_height = Inches(4.5)
    
    # Create table using PowerPoint's native table functionality
    table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table
    
    # Set column widths
    table.columns[0].width = Inches(2.0)  # Methodology
    table.columns[1].width = Inches(4.5)  # Commentary  
    table.columns[2].width = Inches(2.0)  # Enterprise Value
    table.columns[3].width = Inches(1.5)  # Metric
    table.columns[4].width = Inches(2.3)  # 22A'/23E combined
    
    # Header row
    headers = ["Methodology", "Commentary", "Enterprise Value", "Metric", "22A' / 23E (Rev)"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors["primary"]
        
        # Format header text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = fonts["primary_font"]
                run.font.size = Pt(10)
                run.font.bold = True
                run.font.color.rgb = colors["background"]
    
    # Data rows
    for row_idx, row_data in enumerate(valuation_data, 1):
        print(f"[DEBUG] Processing valuation row {row_idx}: {row_data.get('methodology', '')}")
        
        # Methodology
        cell = table.cell(row_idx, 0)
        cell.text = row_data.get('methodology', '')
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors["light_grey"]
        
        # Commentary  
        cell = table.cell(row_idx, 1)
        cell.text = row_data.get('commentary', '')
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors["background"]
        
        # Enterprise Value
        cell = table.cell(row_idx, 2)
        cell.text = row_data.get('enterprise_value', '')
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors["background"]
        
        # Metric
        cell = table.cell(row_idx, 3)
        cell.text = row_data.get('metric', '')
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors["background"]
        
        # 22A' / 23E combined
        cell = table.cell(row_idx, 4)
        a22_val = row_data.get('22a_multiple', '')
        e23_val = row_data.get('23e_multiple', '')
        combined = f"{a22_val} / {e23_val}" if a22_val and e23_val else f"{a22_val}{e23_val}"
        cell.text = combined
        cell.fill.solid()
        cell.fill.fore_color.rgb = colors["background"]
        
        # Format all data cell text
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if col_idx != 1 else PP_ALIGN.LEFT  # Left-align commentary
                for run in paragraph.runs:
                    run.font.name = fonts["primary_font"]
                    run.font.size = Pt(9)
                    run.font.color.rgb = colors["text"]
    
    # Add footer
    footer_left = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(6), Inches(0.4))
    footer_left_frame = footer_left.text_frame
    footer_left_frame.text = f"Confidential | {datetime.now().strftime('%B %Y')}"
    footer_left_para = footer_left_frame.paragraphs[0]
    footer_left_para.font.name = fonts["primary_font"]
    footer_left_para.font.size = Pt(8)
    footer_left_para.font.color.rgb = colors["footer_grey"]
    footer_left_para.alignment = PP_ALIGN.LEFT
    
    footer_right = slide.shapes.add_textbox(Inches(7.333), Inches(7.0), Inches(6), Inches(0.4))
    footer_right_frame = footer_right.text_frame
    footer_right_frame.text = company_name
    footer_right_para = footer_right_frame.paragraphs[0]
    footer_right_para.font.name = fonts["primary_font"]
    footer_right_para.font.size = Pt(8)
    footer_right_para.font.color.rgb = colors["footer_grey"]
    footer_right_para.alignment = PP_ALIGN.RIGHT
    
    return prs

def render_growth_strategy_slide(data=None, color_scheme=None, typography=None, company_name="Moelis", prs=None, brand_config=None, **kwargs):
    """
    Renders a growth strategy & financial projections slide
    """
    
    # Create or use existing presentation
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs = ensure_prs(prs)
    
    # Get brand styling
    colors, fonts = get_brand_styling(brand_config, color_scheme, typography)
    
    # Add blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Helper function
    def add_clean_text(slide, left, top, width, height, text, font_size=10, 
                       color=colors["text"], bold=False, align=PP_ALIGN.LEFT, bg_color=None):
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = align
            for run in paragraph.runs:
                run.font.name = fonts["primary_font"]
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.bold = bold
        
        if bg_color:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = bg_color
        
        textbox.line.fill.background()
        textbox.shadow.inherit = False
        return textbox
    
    # Extract slide data - handle nested structure
    slide_info = data or {}
    if 'slide_data' in slide_info:
        slide_data = slide_info['slide_data']
    else:
        slide_data = slide_info
    
    # STANDARDIZED: Apply header and title
    title_text = slide_data.get('title', 'Growth Strategy & Financial Projections')
    _apply_standard_header_and_title(slide, title_text, brand_config, company_name)
    
    # Growth Strategy section - FIXED: Handle both formats and add fallback content
    growth_strategy = slide_data.get('growth_strategy', {})
    strategies = []
    
    # Check multiple possible data formats
    if growth_strategy:
        strategies = growth_strategy.get('strategies', [])
        print(f"[DEBUG] Growth Strategy: Found {len(strategies)} strategies in growth_strategy.strategies")
    elif 'strategies' in slide_data:
        # Alternative format: strategies directly in slide_data
        strategies = slide_data.get('strategies', [])
        print(f"[DEBUG] Growth Strategy: Found {len(strategies)} strategies in slide_data.strategies")
    elif 'growth_strategies' in slide_data:
        # Another possible format
        strategies = slide_data.get('growth_strategies', [])
        print(f"[DEBUG] Growth Strategy: Found {len(strategies)} strategies in slide_data.growth_strategies")
    
    # Always show the section header
    strategy_title = growth_strategy.get('title', 'Multi-Pronged Growth Strategy') if growth_strategy else 'Multi-Pronged Growth Strategy'
    add_clean_text(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(0.3), 
                   strategy_title, 14, colors["primary"], True)
    
    # Display strategies or fallback message
    if strategies:
        for i, strategy in enumerate(strategies[:6]):  # Max 6 strategies
            y_pos = Inches(1.8 + i * 0.35)
            
            # Gold bullet
            bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y_pos, Inches(0.06), Inches(0.06))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = colors["secondary"]
            bullet.line.fill.background()
            bullet.shadow.inherit = False
            
            # Strategy text
            add_clean_text(slide, Inches(0.85), y_pos - Inches(0.05), Inches(5.5), Inches(0.3), 
                           strategy, 9, colors["text"])
    else:
        # FIXED: Add fallback content when no strategies are available
        print(f"[DEBUG] Growth Strategy: No strategies found, adding fallback message")
        print(f"[DEBUG] Available slide_data keys: {list(slide_data.keys())}")
        
        fallback_strategies = [
            "Strategic growth initiatives will be outlined here",
            "Market expansion and diversification plans", 
            "Operational efficiency improvements",
            "Technology and innovation investments",
            "Sustainability and ESG initiatives"
        ]
        
        for i, strategy in enumerate(fallback_strategies):
            y_pos = Inches(1.8 + i * 0.35)
            
            # Grey bullet for fallback content
            bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y_pos, Inches(0.06), Inches(0.06))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = colors["light_grey"]
            bullet.line.fill.background()
            bullet.shadow.inherit = False
            
            # Strategy text in lighter color
            add_clean_text(slide, Inches(0.85), y_pos - Inches(0.05), Inches(5.5), Inches(0.3), 
                           strategy, 9, colors["footer_grey"])
    
    # Financial Projections Chart - ENHANCED: Better data handling and fallback
    projections = slide_data.get('financial_projections', {})
    
    # Check for alternative data formats
    if not projections:
        # Check for chart data in other possible locations
        if 'chart' in slide_data:
            projections = slide_data['chart']
            print(f"[DEBUG] Growth Strategy: Using chart data from slide_data.chart")
        elif 'projections' in slide_data:
            projections = slide_data['projections']
            print(f"[DEBUG] Growth Strategy: Using projections data from slide_data.projections")
    
    chart_title = 'Revenue & EBITDA Projections'
    
    if projections:
        chart_title = projections.get('chart_title', chart_title) or projections.get('title', chart_title)
        
        # Create chart if we have data
        categories = projections.get('categories', [])
        revenue_data = projections.get('revenue', [])
        ebitda_data = projections.get('ebitda', [])
        
        print(f"[DEBUG] Growth Strategy Chart: {len(categories)} categories, {len(revenue_data)} revenue points, {len(ebitda_data)} ebitda points")
        
        if categories and revenue_data and ebitda_data:
            try:
                chart_data = ChartData()
                chart_data.categories = categories
                chart_data.add_series('Revenue (USD millions)', revenue_data)
                chart_data.add_series('EBITDA (USD millions)', ebitda_data)
                
                # Add chart - MOVED UP and DISABLED LEGEND to prevent overlap
                chart_left = Inches(7.5)
                chart_top = Inches(1.4)  # Moved up from 1.8 to 1.4
                chart_width = Inches(5.5)
                chart_height = Inches(2.2)  # Increased back to 2.2 since no legend
                
                chart_shape = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, chart_left, chart_top, chart_width, chart_height, chart_data
                )
                
                chart = chart_shape.chart
                chart.has_legend = False  # DISABLED legend to prevent overlap
                
                # Add chart title below the chart
                add_clean_text(slide, Inches(7.5), Inches(3.6), Inches(5.5), Inches(0.2), 
                               chart_title, 12, colors["primary"], True)
                
                # Add manual legend as text below chart title
                add_clean_text(slide, Inches(7.5), Inches(3.9), Inches(5.5), Inches(0.2), 
                               "■ Revenue (USD millions)  ■ EBITDA (USD millions)", 10, colors["text"])
                
                # Style chart
                series = chart.series
                # Revenue series (primary color)
                series[0].format.fill.solid()
                series[0].format.fill.fore_color.rgb = colors["primary"]
                
                # EBITDA series (secondary color)
                series[1].format.fill.solid() 
                series[1].format.fill.fore_color.rgb = colors["secondary"]
                
            except Exception as e:
                print(f"[DEBUG] Chart creation error: {e}")
                add_clean_text(slide, Inches(7.5), Inches(2.5), Inches(5.5), Inches(1), 
                               "Financial projections chart will be displayed here.", 12, colors["text"])
        else:
            # Missing chart data - show fallback message
            print(f"[DEBUG] Growth Strategy: Incomplete chart data - categories:{len(categories)}, revenue:{len(revenue_data)}, ebitda:{len(ebitda_data)}")
            add_clean_text(slide, Inches(7.5), Inches(1.4), Inches(5.5), Inches(0.3), 
                           chart_title, 14, colors["primary"], True)
            add_clean_text(slide, Inches(7.5), Inches(2.5), Inches(5.5), Inches(1), 
                           "Financial projections chart will be displayed when data is available.", 12, colors["text"])
    else:
        # No projections data at all - show placeholder
        print(f"[DEBUG] Growth Strategy: No financial projections data found")
        add_clean_text(slide, Inches(7.5), Inches(1.4), Inches(5.5), Inches(0.3), 
                       chart_title, 14, colors["primary"], True)
        add_clean_text(slide, Inches(7.5), Inches(2.5), Inches(5.5), Inches(1), 
                       "Financial projections chart will be displayed when data is available.", 12, colors["text"])
    
    # Key Assumptions - ENHANCED: Handle missing data and alternative formats
    assumptions = slide_data.get('key_assumptions', {})
    assumption_items = []
    
    if assumptions:
        assumption_items = assumptions.get('assumptions', [])
        print(f"[DEBUG] Growth Strategy: Found {len(assumption_items)} assumptions in key_assumptions.assumptions")
    elif 'assumptions' in slide_data:
        # Alternative format: assumptions directly in slide_data
        assumption_items = slide_data.get('assumptions', [])
        print(f"[DEBUG] Growth Strategy: Found {len(assumption_items)} assumptions in slide_data.assumptions")
    
    # Always show the assumptions section header
    assumptions_title = assumptions.get('title', 'Key Planning Assumptions') if assumptions else 'Key Planning Assumptions'
    add_clean_text(slide, Inches(0.5), Inches(4.4), Inches(12.5), Inches(0.3), 
                   assumptions_title, 14, colors["primary"], True)
    
    # Use provided assumptions or fallback
    if not assumption_items:
        # Provide default assumptions if none are available
        assumption_items = [
            "Market growth projections and economic indicators",
            "Regulatory environment and compliance requirements", 
            "Technology adoption and digital transformation ROI",
            "Capital allocation and investment priorities",
            "ESG commitments and sustainability targets",
            "Competitive positioning and market dynamics"
        ]
        print(f"[DEBUG] Growth Strategy: No assumptions found, using {len(assumption_items)} default assumptions")
    
    # Split assumptions into two columns
    left_items = assumption_items[:3]
    right_items = assumption_items[3:]
    
    # Left column
    for i, assumption in enumerate(left_items):
        y_pos = Inches(4.8 + i * 0.3)
        
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y_pos, Inches(0.05), Inches(0.05))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = colors["secondary"]
        bullet.line.fill.background()
        
        add_clean_text(slide, Inches(0.85), y_pos - Inches(0.05), Inches(5.8), Inches(0.25), 
                       assumption, 9, colors["text"])
    
    # Right column
    for i, assumption in enumerate(right_items):
        y_pos = Inches(4.8 + i * 0.3)
        
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), y_pos, Inches(0.05), Inches(0.05))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = colors["secondary"]
        bullet.line.fill.background()
        
        add_clean_text(slide, Inches(7.35), y_pos - Inches(0.05), Inches(5.8), Inches(0.25), 
                       assumption, 9, colors["text"])
    
    # Get today's date
    today = datetime.now().strftime("%B %d, %Y")
    
    # Add footer
    footer_left = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(6), Inches(0.4))
    footer_left_frame = footer_left.text_frame
    footer_left_frame.clear()
    p = footer_left_frame.paragraphs[0]
    p.text = f"Confidential | {today}"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.LEFT
    footer_left_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    footer_right = slide.shapes.add_textbox(Inches(10), Inches(7.0), Inches(3), Inches(0.4))
    footer_right_frame = footer_right.text_frame
    footer_right_frame.clear()
    p = footer_right_frame.paragraphs[0]
    p.text = company_name or "Moelis"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.RIGHT
    footer_right_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return prs

def render_buyer_profiles_slide(data=None, color_scheme=None, typography=None, company_name="Moelis", prs=None, brand_config=None, **kwargs):
    """
    Renders a buyer profiles slide with table layout
    """
    
    # Create or use existing presentation
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs = ensure_prs(prs)
    
    # Get brand styling
    colors, fonts = get_brand_styling(brand_config, color_scheme, typography)
    
    # Add blank slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Helper function
    def add_clean_text(slide, left, top, width, height, text, font_size=10, 
                       color=colors["text"], bold=False, align=PP_ALIGN.LEFT, bg_color=None):
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = align
            for run in paragraph.runs:
                run.font.name = fonts["primary_font"]
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.bold = bold
        
        if bg_color:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = bg_color
        
        textbox.line.fill.background()
        textbox.shadow.inherit = False
        return textbox
    
    # Extract slide data
    slide_data = data or {}
    
    # STANDARDIZED: Apply header and title
    title_text = slide_data.get('title', 'Potential Strategic Buyers')
    _apply_standard_header_and_title(slide, title_text, brand_config, company_name)
    
    # Subtitle if provided
    subtitle = slide_data.get('subtitle', '')
    if subtitle:
        add_clean_text(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.3), 
                       subtitle, 12, colors["accent"], False, PP_ALIGN.LEFT)
    
    # Get table data
    table_rows = slide_data.get('table_rows', [])
    table_headers = slide_data.get('table_headers', ['Buyer Name', 'Description', 'Strategic Rationale', 'Key Synergies', 'Fit'])
    
    if table_rows:
        # Create table
        table_top = Inches(1.6)
        table_left = Inches(0.5)
        table_width = Inches(12.5)
        
        # Calculate rows needed
        num_rows = len(table_rows) + 1  # +1 for header
        num_cols = len(table_headers)
        
        # Adjust height based on number of rows
        row_height = Inches(0.8) if len(table_rows) <= 3 else Inches(0.6)
        table_height = row_height * num_rows
        
        # Create table
        table_shape = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height)
        table = table_shape.table
        
        # Set column widths
        if num_cols == 5:
            col_widths = [Inches(2.8), Inches(2.5), Inches(2.5), Inches(2.2), Inches(1.5)]
        else:
            # Distribute evenly
            col_width = table_width / num_cols
            col_widths = [col_width] * num_cols
        
        for i, width in enumerate(col_widths):
            table.columns[i].width = width
        
        # Format header row
        for j, header in enumerate(table_headers):
            cell = table.cell(0, j)
            cell.text = header
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            
            # Header styling
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors["primary"]
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = fonts["primary_font"]
                    run.font.size = Pt(12)
                    run.font.bold = True
                    run.font.color.rgb = colors["background"]
        
        # Format data rows
        for i, row_data in enumerate(table_rows):
            row_idx = i + 1
            
            # Handle different data formats
            if isinstance(row_data, dict):
                # Convert dict to list based on expected fields
                cell_data = [
                    row_data.get('buyer_name', row_data.get('name', '')),
                    row_data.get('description', ''),
                    row_data.get('strategic_rationale', row_data.get('rationale', '')),
                    row_data.get('key_synergies', row_data.get('synergies', '')),
                    row_data.get('fit', row_data.get('concerns', ''))
                ]
            elif isinstance(row_data, list):
                cell_data = row_data
            else:
                continue
            
            for j, cell_text in enumerate(cell_data):
                if j < num_cols:
                    cell = table.cell(row_idx, j)
                    cell.text = str(cell_text)
                    cell.margin_left = Inches(0.05)
                    cell.margin_right = Inches(0.05)
                    cell.margin_top = Inches(0.05)
                    cell.margin_bottom = Inches(0.05)
                    
                    # Alternate row colors
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = colors["light_grey"]
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = colors["background"]
                    
                    # Text styling
                    cell.text_frame.word_wrap = True
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.LEFT
                        for run in paragraph.runs:
                            run.font.name = fonts["primary_font"]
                            run.font.size = Pt(11)
                            run.font.color.rgb = colors["text"]
                            # Make first column (buyer name) bold
                            if j == 0:
                                run.font.bold = True
    
    else:
        # No data - add placeholder
        add_clean_text(slide, Inches(2), Inches(3), Inches(9), Inches(1), 
                       "Buyer profiles data will be displayed here when available.", 
                       14, colors["text"], False, PP_ALIGN.CENTER)
    
    # Get today's date
    today = datetime.now().strftime("%B %d, %Y")
    
    # Add footer
    footer_left = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(6), Inches(0.4))
    footer_left_frame = footer_left.text_frame
    footer_left_frame.clear()
    p = footer_left_frame.paragraphs[0]
    p.text = f"Confidential | {today}"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.LEFT
    footer_left_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    footer_right = slide.shapes.add_textbox(Inches(10), Inches(7.0), Inches(3), Inches(0.4))
    footer_right_frame = footer_right.text_frame
    footer_right_frame.clear()
    p = footer_right_frame.paragraphs[0]
    p.text = company_name or "Moelis"
    p.font.name = fonts["primary_font"]
    p.font.size = fonts["small_size"]
    p.font.color.rgb = colors["footer_grey"]
    p.alignment = PP_ALIGN.RIGHT
    footer_right_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return prs


def render_sea_conglomerates_slide(data=None, color_scheme=None, typography=None, company_name="Moelis", prs=None, brand_config=None, **kwargs):
    """
    Render SEA conglomerates slide with brand configuration support
    """
    
    # FIXED: Extract slide_data from data parameter - handle all possible formats
    slide_data = []
    if isinstance(data, list):
        # Direct list of conglomerates
        slide_data = data
    elif isinstance(data, dict):
        # Check multiple possible keys where the actual data might be
        slide_data = (
            data.get('sea_conglomerates', []) or  # Primary: {"sea_conglomerates": [...]} from structure fix
            data.get('data', []) or  # Fallback: {"data": [...]}
            data.get('conglomerates', []) or  # Alternative: {"conglomerates": [...]}
            []
        )
        # If data is a dict but contains the conglomerate fields directly, wrap it
        if not slide_data and 'name' in data and 'country' in data:
            slide_data = [data]
    
    print(f"[DEBUG] SEA conglomerates: Found {len(slide_data)} companies")
    print(f"[DEBUG] First company data: {slide_data[0] if slide_data else 'No data'}")
    print(f"[DEBUG] Input data type: {type(data)}")
    print(f"[DEBUG] Input data keys: {list(data.keys()) if isinstance(data, dict) else 'Not a dict'}")
    print(f"[DEBUG] Using fallback data: {not slide_data}")
    
    # Default data if none provided - FALLBACK ONLY (Oil & Gas focused)
    if not slide_data:
        slide_data = [
            {
                "name": "Reliance Industries Limited",
                "country": "India",
                "description": "Large diversified conglomerate with significant oil refining and petrochemicals operations",
                "key_shareholders": "Mukesh Ambani family trust and public investors",
                "key_financials": "Revenue: US$104B, Strong energy sector presence",
                "contact": "Managing Director - Energy Sector"
            },
            {
                "name": "Mitsubishi Corporation",
                "country": "Japan", 
                "description": "Global trading and investment conglomerate with extensive energy and natural resources portfolio",
                "key_shareholders": "Mitsubishi Group companies and institutional investors",
                "key_financials": "Revenue: US$156B+, Major energy trading operations",
                "contact": "Executive Director - Energy Trading"
            },
            {
                "name": "China National Petroleum Corporation",
                "country": "China",
                "description": "State-owned oil and gas corporation with integrated upstream and downstream operations",
                "key_shareholders": "Chinese state-owned enterprise",
                "key_financials": "Revenue: US$480B+, Leading global oil producer",
                "contact": "Managing Director - International Operations"
            },
            {
                "name": "Adnoc Group", 
                "country": "UAE",
                "description": "National oil company with diversified energy portfolio and strategic international partnerships",
                "key_shareholders": "UAE government and sovereign funds",
                "key_financials": "Revenue: US$60B+, Expanding global footprint",
                "contact": "Executive Director - Corporate Development"
            }
        ]
    
    # Create or use existing presentation
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs = ensure_prs(prs)
    
    # Get brand styling
    colors, fonts = get_brand_styling(brand_config, color_scheme, typography)
    
    # Pagination: Split data into chunks of 4
    max_entries_per_slide = 4
    slide_chunks = []
    
    for i in range(0, len(slide_data), max_entries_per_slide):
        chunk = slide_data[i:i + max_entries_per_slide]
        slide_chunks.append(chunk)
    
    # Create slides for each chunk
    for slide_index, chunk_data in enumerate(slide_chunks):
        # Add blank slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title with pagination info
        title_text = "SEA Conglomerate Strategic Buyers"
        if len(slide_chunks) > 1:
            title_text += f" (cont'd)" if slide_index > 0 else ""
        
        # STANDARDIZED: Apply header and title
        _apply_standard_header_and_title(slide, title_text, brand_config, company_name)
        
        # Table dimensions and positioning
        table_left = Inches(0.5)
        table_top = Inches(1.5)
        table_width = Inches(12.333)
        table_height = Inches(5.0)
        
        # Calculate number of rows (header + data rows)
        num_rows = len(chunk_data) + 1
        num_cols = 6
        
        # Create table
        table_shape = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height)
        table = table_shape.table
        
        # Column headers
        headers = ["Name", "Country", "Description", "Key shareholders", "Key financials (US$m)", "Contact"]
        
        # Set column widths to match original
        col_widths = [Inches(1.8), Inches(1.0), Inches(4.2), Inches(2.2), Inches(1.8), Inches(1.333)]
        for i, width in enumerate(col_widths):
            table.columns[i].width = width
        
        # Format header row - DARK BLUE BACKGROUND
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            # DARK BLUE BACKGROUND
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors["primary"]
            
            # Header text
            cell.text = header
            cell.text_frame.margin_left = Inches(0.1)
            cell.text_frame.margin_right = Inches(0.1)
            cell.text_frame.margin_top = Inches(0.05)
            cell.text_frame.margin_bottom = Inches(0.05)
            
            # Header formatting - WHITE TEXT on dark blue background
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = fonts["primary_font"]
                paragraph.font.size = fonts["header_size"]
                paragraph.font.bold = True
                paragraph.font.color.rgb = colors["background"]  # WHITE text on dark blue
                paragraph.alignment = PP_ALIGN.CENTER
        
        # Populate data rows - WHITE BACKGROUND
        for row_idx, company in enumerate(chunk_data, 1):
            # Define the data for each column - handle different data structures
            row_data = [
                company.get('name', ''),
                company.get('country', ''),
                company.get('description', company.get('healthcare_focus', '')),  # Fallback to healthcare_focus
                company.get('key_shareholders', 'N/A'),
                company.get('key_financials', company.get('revenue', '')),  # Fallback to revenue
                company.get('contact', company.get('moelis_contact', 'To be assigned'))
            ]
            
            for col_idx, data in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(data)
                
                # WHITE BACKGROUND
                cell.fill.solid()
                cell.fill.fore_color.rgb = colors["background"]
                
                # Cell margins
                cell.text_frame.margin_left = Inches(0.1)
                cell.text_frame.margin_right = Inches(0.1)
                cell.text_frame.margin_top = Inches(0.05)
                cell.text_frame.margin_bottom = Inches(0.05)
                
                # Text wrapping
                cell.text_frame.word_wrap = True
                
                # Text formatting - BLACK TEXT
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = fonts["primary_font"]
                    paragraph.font.size = fonts["body_size"]
                    paragraph.font.bold = False
                    paragraph.font.color.rgb = colors["text"]
                    paragraph.alignment = PP_ALIGN.LEFT
        
        # Set row heights
        header_row_height = Inches(0.6)
        data_row_height = Inches(1.1)
        
        # Set header row height
        table.rows[0].height = header_row_height
        
        # Set data row heights
        for i in range(1, num_rows):
            table.rows[i].height = data_row_height
        
        # Get today's date
        today = datetime.now().strftime("%B %d, %Y")
        
        # Add footer
        footer_left = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(6), Inches(0.4))
        footer_left_frame = footer_left.text_frame
        footer_left_frame.clear()
        p = footer_left_frame.paragraphs[0]
        p.text = f"Confidential | {today}"
        p.font.name = fonts["primary_font"]
        p.font.size = fonts["small_size"]
        p.font.color.rgb = colors["footer_grey"]
        p.alignment = PP_ALIGN.LEFT
        footer_left_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        footer_right = slide.shapes.add_textbox(Inches(10), Inches(7.0), Inches(3), Inches(0.4))
        footer_right_frame = footer_right.text_frame
        footer_right_frame.clear()
        p = footer_right_frame.paragraphs[0]
        p.text = company_name or "Moelis"
        p.font.name = fonts["primary_font"]
        p.font.size = fonts["small_size"]
        p.font.color.rgb = colors["footer_grey"]
        p.alignment = PP_ALIGN.RIGHT
        footer_right_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return prs
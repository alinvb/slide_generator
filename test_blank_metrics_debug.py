#!/usr/bin/env python3

"""
Debug test to replicate the specific issue shown in the user's image:
1. Blank Key Operational Metrics box  
2. Table with repeated placeholder text
"""

from pptx import Presentation
from slide_templates import render_product_service_footprint_slide

def test_problematic_scenarios():
    """Test scenarios that might cause the issues seen in user's image"""
    print("=== Testing Problematic Data Scenarios ===")
    
    # Scenario 1: Empty metrics
    print("\n1. Testing with empty metrics...")
    empty_metrics_data = {
        "title": "Product & Service Footprint",
        "services": [
            {"title": "Crude Oil Production", "desc": "World's largest producer"},
            {"title": "Natural Gas & NGLs", "desc": "Extraction and processing"},
        ],
        "coverage_table": [
            ["Region", "Key Operations", "Major Facilities"],
            ["Region", "Key Operations", "Major Facilities"],  # Repeated headers (problematic)
            ["Region", "Key Operations", "Major Facilities"],  # Repeated headers (problematic)
        ],
        "metrics": {}  # EMPTY METRICS - This could cause blank box
    }
    
    prs = Presentation()
    prs = render_product_service_footprint_slide(data=empty_metrics_data, prs=prs)
    prs.save("test_empty_metrics.pptx")
    print("✓ Empty metrics test saved")
    
    # Scenario 2: Missing metrics key entirely
    print("\n2. Testing with missing metrics key...")
    no_metrics_data = {
        "title": "Product & Service Footprint", 
        "services": [
            {"title": "Crude Oil Production", "desc": "World's largest producer"},
        ],
        "coverage_table": [
            ["Region", "Key Operations", "Major Facilities/Partners"],
            ["Region", "Key Operations", "Major Facilities/Partners"],  # Problematic repeated content
            ["Region", "Key Operations", "Major Facilities/Partners"],
        ]
        # NO METRICS KEY AT ALL
    }
    
    prs = Presentation()
    prs = render_product_service_footprint_slide(data=no_metrics_data, prs=prs)
    prs.save("test_no_metrics_key.pptx")
    print("✓ No metrics key test saved")
    
    # Scenario 3: Malformed table data (repeated headers)
    print("\n3. Testing with malformed table data...")
    malformed_table_data = {
        "title": "Product & Service Footprint",
        "services": [
            {"title": "Crude Oil Production", "desc": "World's largest producer"},
            {"title": "Natural Gas & NGLs", "desc": "Extraction and processing"},
        ],
        "coverage_table": [
            ["Region", "Key Operations", "Major Facilities/Partners"],
            ["Region", "Key Operations", "Major Facilities/Partners"],  # This is the issue!
            ["Region", "Key Operations", "Major Facilities/Partners"],  # Repeated instead of data
            ["Region", "Key Operations", "Major Facilities/Partners"],  # This creates the placeholder look
        ],
        "metrics": {
            "production": {"label": "Daily Production", "value": "12.8M bbl/day"},
            "reserves": {"label": "Proven Reserves", "value": "267B barrels"}
        }
    }
    
    prs = Presentation()  
    prs = render_product_service_footprint_slide(data=malformed_table_data, prs=prs)
    prs.save("test_malformed_table.pptx")
    print("✓ Malformed table test saved")
    
    # Scenario 4: String metrics instead of objects (common JSON issue)
    print("\n4. Testing with string metrics (before JSON fixing)...")
    string_metrics_data = {
        "title": "Product & Service Footprint",
        "services": [
            {"title": "Crude Oil Production", "desc": "World's largest producer"},
        ],
        "coverage_table": [
            ["Region", "Key Operations", "Major Facilities/Partners"],
            ["Saudi Arabia", "Upstream Production", "Ghawar, Safaniya"],
        ],
        "metrics": {
            "production": "12.8M bbl/day",  # STRING instead of object
            "reserves": "267B barrels",     # STRING instead of object  
            "capacity": "5.4M bbl/day"      # STRING instead of object
        }
    }
    
    prs = Presentation()
    prs = render_product_service_footprint_slide(data=string_metrics_data, prs=prs)
    prs.save("test_string_metrics.pptx")
    print("✓ String metrics test saved")
    
    # Scenario 5: Completely empty data (worst case)
    print("\n5. Testing with completely empty data...")
    empty_data = {
        "title": "Product & Service Footprint"
        # Everything else missing
    }
    
    prs = Presentation()
    prs = render_product_service_footprint_slide(data=empty_data, prs=prs)  
    prs.save("test_completely_empty.pptx")
    print("✓ Completely empty test saved")
    
    return True

def test_json_fixer_on_problematic_data():
    """Test if JSON fixer handles the problematic scenarios"""
    print("\n=== Testing JSON Fixer on Problematic Data ===")
    
    try:
        from json_data_fixer import fix_product_service_footprint_data
        
        # Test string metrics fixing
        problematic_data = {
            "metrics": {
                "production": "12.8M bbl/day",  # Should be converted to object
                "reserves": "267B barrels"
            }
        }
        
        print("Before fixing:", problematic_data["metrics"])
        fixed_data = fix_product_service_footprint_data(problematic_data)
        print("After fixing:", fixed_data["metrics"])
        
        # Check if it was fixed correctly
        first_metric = list(fixed_data["metrics"].values())[0]
        if isinstance(first_metric, dict) and "label" in first_metric and "value" in first_metric:
            print("✓ JSON fixer correctly converts string metrics to objects")
        else:
            print("❌ JSON fixer did not fix string metrics properly")
        
        return True
        
    except Exception as e:
        print(f"❌ JSON fixer test failed: {e}")
        return False

def main():
    """Run problematic scenario tests"""
    print("Testing problematic scenarios that could cause blank metrics and repeated table text...")
    
    try:
        success1 = test_problematic_scenarios()
        success2 = test_json_fixer_on_problematic_data()
        
        if success1 and success2:
            print(f"\n🔍 PROBLEMATIC SCENARIOS TESTED!")
            print("Generated test files to check:")
            print("- test_empty_metrics.pptx (empty metrics object)")
            print("- test_no_metrics_key.pptx (missing metrics key entirely)")  
            print("- test_malformed_table.pptx (repeated headers instead of data)")
            print("- test_string_metrics.pptx (string values instead of objects)")
            print("- test_completely_empty.pptx (minimal data)")
            print(f"\n💡 Check these files to see which scenario matches your issue!")
            
        else:
            print("\n❌ Some tests failed.")
    except Exception as e:
        print(f"\n💥 Error: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()
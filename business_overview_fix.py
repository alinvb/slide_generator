"""
Fixed Business Overview Slide Template
Resolves text truncation, overlapping elements, and formatting issues
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

def render_business_overview_slide_fixed(data=None, color_scheme=None, typography=None, company_name="Moelis", prs=None, brand_config=None, **kwargs):
    """
    FIXED Business Overview slide with proper text handling and layout
    """
    
    # Create presentation if not provided 
    if prs is None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    
    # Get brand styling (using defaults for now)
    colors = {
        "background": RGBColor(255, 255, 255),
        "primary": RGBColor(31, 56, 100), 
        "secondary": RGBColor(218, 165, 32),
        "text": RGBColor(64, 64, 64),
        "light_grey": RGBColor(248, 249, 250),
        "footer_grey": RGBColor(128, 128, 128)
    }
    
    fonts = {
        "primary_font": "Calibri",
        "small_size": Pt(8)
    }
    
    # Add slide with blank layout
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)
    
    # Set white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colors["background"]
    
    def add_robust_text(slide, left, top, width, height, text, font_size=12, 
                       color=None, bold=False, align=PP_ALIGN.LEFT):
        """Enhanced text function with robust handling"""
        if color is None:
            color = colors["text"]
            
        # Handle text content safely
        text = str(text) if text else ""
        if len(text) > 400:  # Limit text length
            text = text[:397] + "..."
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()  # Clear default paragraph
        
        # Configure text frame
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05) 
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        
        # Add paragraph and set text
        paragraph = text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = align
        
        # Format text
        for run in paragraph.runs:
            run.font.name = fonts["primary_font"]
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold
        
        # Remove borders
        textbox.line.fill.background()
        textbox.shadow.inherit = False
        
        return textbox
    
    # Extract slide data
    slide_data = data or {}
    
    # Header - Company name and slide title
    title_text = slide_data.get('title', 'Business Overview')
    add_robust_text(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.4), 
                   title_text, 18, colors["primary"], True, PP_ALIGN.LEFT)
    
    # Company description - FIXED positioning and sizing
    company_desc = slide_data.get('description', 
        'Leading integrated company providing comprehensive services with operational excellence.')
    
    # Split long descriptions into manageable chunks
    if len(company_desc) > 200:
        # First part 
        add_robust_text(slide, Inches(0.5), Inches(1.0), Inches(7.5), Inches(0.8), 
                       company_desc[:200] + "...", 11, colors["text"])
        
        # Second part if needed
        if len(company_desc) > 200:
            add_robust_text(slide, Inches(0.5), Inches(1.7), Inches(7.5), Inches(0.6), 
                           "..." + company_desc[200:], 11, colors["text"])
    else:
        add_robust_text(slide, Inches(0.5), Inches(1.0), Inches(7.5), Inches(1.0), 
                       company_desc, 11, colors["text"])
    
    # Timeline section - REPOSITIONED
    timeline_data = slide_data.get('timeline', {
        'start_year': '2015', 
        'end_year': '2024'
    })
    
    timeline_y = Inches(2.4)
    
    try:
        # Start circle
        start_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), timeline_y, 
                                            Inches(0.1), Inches(0.1))
        start_circle.fill.solid()
        start_circle.fill.fore_color.rgb = colors["secondary"]
        start_circle.line.fill.background()
        
        # Timeline line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), timeline_y + Inches(0.04), 
                                    Inches(3), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = colors["secondary"] 
        line.line.fill.background()
        
        # End circle
        end_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.8), timeline_y, 
                                          Inches(0.1), Inches(0.1))
        end_circle.fill.solid()
        end_circle.fill.fore_color.rgb = colors["secondary"]
        end_circle.line.fill.background()
        
        # Year labels
        add_robust_text(slide, Inches(0.7), timeline_y - Inches(0.25), Inches(0.3), Inches(0.2), 
                       timeline_data.get('start_year', '2015'), 10, colors["primary"], True, PP_ALIGN.CENTER)
        
        add_robust_text(slide, Inches(3.7), timeline_y - Inches(0.25), Inches(0.3), Inches(0.2), 
                       timeline_data.get('end_year', '2024'), 10, colors["primary"], True, PP_ALIGN.CENTER)
        
    except Exception as e:
        print(f"Timeline creation error: {e}")
    
    # Key highlights section - RIGHT SIDE
    highlights_title = slide_data.get('highlights_title', 'Key Highlights')
    add_robust_text(slide, Inches(8.5), Inches(1.0), Inches(4.2), Inches(0.3), 
                   highlights_title, 14, colors["primary"], True, PP_ALIGN.CENTER)
    
    # Background box for highlights
    highlights_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), Inches(1.3), 
                                         Inches(4.6), Inches(4.8))
    highlights_bg.fill.solid()
    highlights_bg.fill.fore_color.rgb = colors["light_grey"]
    highlights_bg.line.fill.background()
    
    # Highlight items
    highlights = slide_data.get('highlights', [
        'Market leader in sector',
        'Strong financial performance', 
        'Experienced management team',
        'Global operational presence',
        'Advanced technology platform'
    ])
    
    for i, highlight in enumerate(highlights[:6]):  # Max 6 items
        y_pos = Inches(1.5 + i * 0.7)
        
        # Bullet point
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), y_pos + Inches(0.05), 
                                       Inches(0.04), Inches(0.04))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = colors["secondary"]
        bullet.line.fill.background()
        
        # Highlight text - PROPER SIZING
        add_robust_text(slide, Inches(8.6), y_pos, Inches(4.1), Inches(0.5), 
                       highlight, 9, colors["text"])
    
    # Services section - BOTTOM LEFT
    services_title = slide_data.get('services_title', 'Core Services')
    add_robust_text(slide, Inches(0.5), Inches(2.9), Inches(4), Inches(0.3), 
                   services_title, 12, colors["primary"], True)
    
    services = slide_data.get('services', [
        'Primary service offering',
        'Secondary service line', 
        'Additional capabilities'
    ])
    
    for i, service in enumerate(services[:3]):  # Max 3 services to fit
        y_pos = Inches(3.3 + i * 0.5)
        
        # Service bullet
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y_pos + Inches(0.05), 
                                       Inches(0.03), Inches(0.03))
        bullet.fill.solid() 
        bullet.fill.fore_color.rgb = colors["secondary"]
        bullet.line.fill.background()
        
        # Service text
        add_robust_text(slide, Inches(0.8), y_pos, Inches(3.8), Inches(0.4), 
                       service, 9, colors["text"])
    
    # Strategic positioning - BOTTOM SECTION
    positioning_title = slide_data.get('positioning_title', 'Strategic Positioning')
    add_robust_text(slide, Inches(0.5), Inches(5.0), Inches(6), Inches(0.3), 
                   positioning_title, 12, colors["primary"], True)
    
    positioning_desc = slide_data.get('positioning_desc', 
        'The company maintains a strong competitive position through operational excellence and strategic market presence.')
    
    add_robust_text(slide, Inches(0.5), Inches(5.3), Inches(7.5), Inches(1.0), 
                   positioning_desc, 10, colors["text"])
    
    # Footer
    footer_y = Inches(6.8)
    current_date = datetime.now().strftime("%B %d, %Y")
    
    # Left footer
    add_robust_text(slide, Inches(0.5), footer_y, Inches(4), Inches(0.2), 
                   f"Confidential | {current_date}", 8, colors["footer_grey"])
    
    # Right footer  
    add_robust_text(slide, Inches(9), footer_y, Inches(4), Inches(0.2), 
                   company_name, 8, colors["footer_grey"], False, PP_ALIGN.RIGHT)
    
    return prs


# Integration function to replace the problematic slide template
def get_fixed_business_overview_renderer():
    """Return the fixed business overview slide renderer"""
    return render_business_overview_slide_fixed
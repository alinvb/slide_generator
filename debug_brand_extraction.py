#!/usr/bin/env python3
"""
Quick debug script to test brand extraction without Streamlit interference
"""

import io
from brand_extractor import BrandExtractor
from pptx import Presentation
from pptx.dml.color import RGBColor

def create_test_pptx():
    """Create a simple test PowerPoint with actual brand colors"""
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    
    prs = Presentation()
    
    # Add a slide with some colored shapes
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank slide
    
    # Add a title with custom color
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Test Brand Deck"
    
    # Set title text color to blue
    for paragraph in title_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 112, 192)  # Blue
            run.font.size = Pt(24)
    
    # Add a colored rectangle (brand primary color)
    rect = slide.shapes.add_shape(
        1, Inches(2), Inches(2), Inches(6), Inches(1)  # Rectangle shape
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(255, 102, 0)  # Orange
    
    # Add another shape with secondary color
    rect2 = slide.shapes.add_shape(
        1, Inches(2), Inches(4), Inches(6), Inches(1)
    )
    rect2.fill.solid() 
    rect2.fill.fore_color.rgb = RGBColor(0, 176, 80)  # Green
    
    # Save to BytesIO
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

def test_brand_extraction():
    """Test the brand extraction directly"""
    print("Creating test PowerPoint...")
    test_pptx = create_test_pptx()
    
    print("Initializing BrandExtractor...")
    extractor = BrandExtractor()
    
    print("Testing rule-based extraction...")
    try:
        result = extractor.extract_brand_from_pptx(test_pptx, use_llm=False)
        
        print("Extraction successful!")
        print("Colors found:")
        for name, color in result.get('color_scheme', {}).items():
            if hasattr(color, 'r'):
                print(f"  {name}: RGB({color.r}, {color.g}, {color.b})")
            else:
                print(f"  {name}: {color}")
                
        print("Typography:")
        for name, value in result.get('typography', {}).items():
            print(f"  {name}: {value}")
            
        return True
        
    except Exception as e:
        print(f"Extraction failed: {e}")
        import traceback
        print(traceback.format_exc())
        return False

if __name__ == "__main__":
    # Need to import Pt for the test
    from pptx.util import Pt
    test_brand_extraction()
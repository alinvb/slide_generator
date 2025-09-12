from slide_templates import render_precedent_transactions_slide
from pptx import Presentation
import json

def test_precedent_transactions_rendering():
    """Test the actual rendering of precedent transactions with real data"""
    
    # Sample data that matches what extract_slide_data produces
    test_data = {
        "title": "Precedent Transactions",
        "transactions": [
            {
                "target": "Netflix Competitor A",
                "acquirer": "Media Giant Corp",
                "date": "2023",
                "country": "USA",
                "enterprise_value": "$15.2B",
                "revenue": "$3.4B", 
                "ev_revenue_multiple": "4.5x"
            },
            {
                "target": "Streaming Service B", 
                "acquirer": "Tech Company",
                "date": "2022",
                "country": "UK",
                "enterprise_value": "$8.1B",
                "revenue": "$1.8B",
                "ev_revenue_multiple": "4.5x"
            },
            {
                "target": "Content Platform C",
                "acquirer": "Entertainment Corp",
                "date": "2023", 
                "country": "Canada",
                "enterprise_value": "$5.3B",
                "revenue": "$1.2B",
                "ev_revenue_multiple": "4.4x"
            }
        ]
    }
    
    # Test brand config (minimal)
    brand_config = {
        "templates": {
            "corporate": {
                "color_scheme": {
                    "primary": {"r": 88, "g": 28, "b": 135},
                    "secondary": {"r": 147, "g": 51, "b": 234},
                    "accent": {"r": 196, "g": 181, "b": 253},
                }
            }
        }
    }
    
    print("=== TESTING PRECEDENT TRANSACTIONS RENDERING ===")
    print(f"Input data structure: {json.dumps(test_data, indent=2)}")
    
    # Test rendering 
    try:
        prs = render_precedent_transactions_slide(
            data=test_data, 
            company_name="Netflix",
            brand_config=brand_config,
            template_name="corporate"
        )
        
        print("✅ Rendering succeeded!")
        print(f"Presentation slides count: {len(prs.slides)}")
        
        # Save test file
        prs.save("test_precedent_transactions.pptx")
        print("✅ Test file saved as 'test_precedent_transactions.pptx'")
        
        # Check if placeholder content was used
        slide = prs.slides[0]
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame.text:
                text_content.append(shape.text_frame.text)
        
        # Look for placeholder text
        placeholder_found = any('placeholder' in text.lower() or 'will be displayed' in text.lower() for text in text_content)
        
        if placeholder_found:
            print("⚠️ ISSUE DETECTED: Placeholder content found in rendered slide!")
            print("Slide text content:")
            for i, text in enumerate(text_content):
                print(f"  Text {i+1}: {text}")
        else:
            print("✅ No placeholder content detected")
            print("✅ Transaction data appears to be rendered properly")
    
    except Exception as e:
        print(f"❌ Rendering failed: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    test_precedent_transactions_rendering()

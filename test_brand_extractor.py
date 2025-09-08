#!/usr/bin/env python3

import sys
import io
from brand_extractor import BrandExtractor

def test_brand_extractor():
    """Test the brand extractor functionality"""
    print("🧪 Testing Brand Extractor...")
    
    # Initialize brand extractor
    try:
        brand_extractor = BrandExtractor()
        print("✅ BrandExtractor initialized successfully")
    except Exception as e:
        print(f"❌ Failed to initialize BrandExtractor: {e}")
        return False
    
    # Test with a simple PPTX creation
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        
        # Create a simple test presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
        
        # Add a title with custom color
        title = slide.shapes.title
        title.text = "Test Brand Slide"
        
        # Try to set some colors
        if hasattr(title, 'text_frame'):
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)  # Red
        
        # Save to bytes
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        print("✅ Test PPTX created successfully")
        
        # Test rule-based extraction
        print("🔧 Testing rule-based extraction...")
        brand_config = brand_extractor.extract_brand_from_pptx(pptx_bytes, use_llm=False)
        
        print(f"🎨 Extracted brand config:")
        print(f"   Color scheme keys: {list(brand_config.get('color_scheme', {}).keys())}")
        print(f"   Typography keys: {list(brand_config.get('typography', {}).keys())}")
        print(f"   Header style: {brand_config.get('header_style', {})}")
        
        # Check if colors are extracted
        colors = brand_config.get('color_scheme', {})
        if colors:
            print("✅ Colors extracted successfully!")
            for name, color in colors.items():
                print(f"   {name}: {color}")
            return True
        else:
            print("❌ No colors extracted!")
            return False
            
    except Exception as e:
        print(f"❌ Test failed with error: {e}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    success = test_brand_extractor()
    sys.exit(0 if success else 1)
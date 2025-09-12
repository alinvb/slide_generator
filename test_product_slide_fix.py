#!/usr/bin/env python3
"""
Test to debug the product slide data population issue
"""

import sys
import json
from slide_templates import render_product_service_footprint_slide
from pptx import Presentation

def test_product_slide_rendering():
    """Test the product slide renderer with sample data"""
    
    # Sample data that should work
    test_data = {
        "title": "Product & Service Footprint",
        "services": [
            {
                "title": "Cloud Infrastructure",
                "desc": "Scalable cloud hosting and infrastructure services for enterprise clients"
            },
            {
                "title": "Software Development",
                "desc": "Custom software solutions and application development"
            },
            {
                "title": "Data Analytics",
                "desc": "Business intelligence and data analytics platforms"
            }
        ],
        "table_title": "Geographic Market Coverage",
        "coverage_table": [
            ["Region", "Market Share", "Revenue %"],
            ["North America", "45%", "52%"],
            ["Europe", "30%", "28%"],
            ["Asia Pacific", "25%", "20%"]
        ],
        "metrics": {
            "total_customers": 1250,
            "geographic_reach": "15 countries",
            "service_lines": 3
        }
    }
    
    print("🧪 Testing Product Service Footprint Slide Renderer...")
    print(f"📊 Test data structure: {json.dumps(test_data, indent=2)}")
    
    try:
        # Create presentation
        prs = Presentation()
        
        # Test the renderer
        result_prs = render_product_service_footprint_slide(
            data=test_data,
            company_name="Test Company",
            prs=prs,
            template_name="professional"
        )
        
        print(f"✅ Renderer executed successfully!")
        print(f"📄 Slides in presentation: {len(result_prs.slides)}")
        
        # Save test output
        result_prs.save("/home/user/webapp/test_product_slide_output.pptx")
        print(f"💾 Test presentation saved successfully!")
        
        return True
        
    except Exception as e:
        print(f"❌ Error in product slide renderer: {e}")
        import traceback
        traceback.print_exc()
        return False

def test_empty_data_handling():
    """Test how the renderer handles empty or missing data"""
    
    print("\n🧪 Testing Empty Data Handling...")
    
    test_cases = [
        {},  # Empty dict
        {"services": []},  # Empty services
        {"services": [{"title": "Test", "desc": ""}]},  # Partial data
        None  # None data
    ]
    
    for i, test_data in enumerate(test_cases):
        print(f"\n📋 Test case {i+1}: {test_data}")
        
        try:
            prs = Presentation()
            result_prs = render_product_service_footprint_slide(
                data=test_data,
                company_name="Test Company",
                prs=prs,
                template_name="professional"
            )
            print(f"✅ Test case {i+1} handled successfully")
            
        except Exception as e:
            print(f"❌ Test case {i+1} failed: {e}")

if __name__ == "__main__":
    print("🔍 Product Slide Renderer Debug Test")
    print("=" * 50)
    
    success = test_product_slide_rendering()
    test_empty_data_handling()
    
    if success:
        print("\n✅ All tests completed! Check the output file to verify rendering.")
    else:
        print("\n❌ Tests found issues that need to be fixed.")
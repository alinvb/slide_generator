"""
adapters.py - patched and fixed
Robust adapter that maps template keys to your concrete renderer functions
and calls them with the right signature: fn(data=..., prs=..., company_name=...).
Accepts plan/content/content_ir forms.
Now supports brand configuration and standardized formatting.
"""
from typing import Any, Dict, List, Optional
try:
    from pptx import Presentation  # type: ignore
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
except Exception:
    Presentation = None  # type: ignore
    RGBColor = None
    Pt = None
    PP_ALIGN = None

import importlib

# Import your renderers module (must be importable on PYTHONPATH)
slide_templates = importlib.import_module("slide_templates")

# ---- renderer map ----

def _get_attr(name: str):
    return getattr(slide_templates, name, None)

RENDERER_MAP = {
    # core
    "management_team": _get_attr("render_management_team_slide"),
    "historical_financial_performance": _get_attr("render_historical_financial_performance_slide"),
    "margin_cost_resilience": _get_attr("render_margin_cost_resilience_slide"),
    "investor_considerations": _get_attr("render_investor_considerations_slide"),
    "investor_process_overview": _get_attr("render_investor_process_overview_slide"),
    "product_service_footprint": _get_attr("render_product_service_footprint_slide"),
    "business_overview": _get_attr("render_business_overview_slide"),
    # analysis
    "competitive_positioning": _get_attr("render_competitive_positioning_slide"),
    "precedent_transactions": _get_attr("render_precedent_transactions_slide"),
    "valuation_overview": _get_attr("render_valuation_overview_slide"),
    # buyers
    "sea_conglomerates": _get_attr("render_sea_conglomerates_slide"),
    "buyer_profiles": _get_attr("render_buyer_profiles_slide"),
    # growth
    "growth_strategy_projections": _get_attr("render_growth_strategy_slide"),
}

# ---- brand configuration helpers ----

def _convert_brand_colors(brand_config: Optional[Dict]) -> Optional[Dict]:
    """Convert brand configuration colors to RGBColor objects"""
    if not brand_config or not brand_config.get('color_scheme'):
        return None
    
    color_scheme = brand_config['color_scheme']
    converted_colors = {}
    
    for name, color in color_scheme.items():
        if isinstance(color, str) and color.startswith('#'):
            # Convert hex to RGBColor
            hex_color = color.lstrip('#')
            converted_colors[name] = RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16), 
                int(hex_color[4:6], 16)
            )
        elif isinstance(color, tuple) and len(color) == 3:
            # Convert tuple (r, g, b) to RGBColor
            r, g, b = color
            converted_colors[name] = RGBColor(r, g, b)
        elif hasattr(color, 'r'):  # Already RGBColor
            converted_colors[name] = color
        else:
            # Fallback to default
            converted_colors[name] = RGBColor(24, 58, 88)  # Default blue
    
    return converted_colors

def _standardize_typography(brand_config: Optional[Dict]) -> Optional[Dict]:
    """Standardize typography configuration"""
    if not brand_config or not brand_config.get('typography'):
        return None
    
    typography = brand_config['typography']
    standardized = {}
    
    for name, value in typography.items():
        if isinstance(value, str) and 'Pt(' in value:
            # Extract point size from string like "Pt(24)"
            try:
                size_str = value.split('Pt(')[1].split(')')[0]
                standardized[name] = Pt(int(size_str))
            except:
                standardized[name] = value
        elif hasattr(value, 'pt'):  # Already Pt object
            standardized[name] = value
        else:
            standardized[name] = value
    
    return standardized

def _apply_standard_header_and_title(slide, title_text: str, brand_config: Optional[Dict] = None, 
                                   company_name: str = "Moelis"):
    """Apply standardized header and left-aligned title to all slides"""
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    
    # Get brand styling or use defaults
    if brand_config:
        colors = _convert_brand_colors(brand_config) or {}
        typography = _standardize_typography(brand_config) or {}
        primary_color = colors.get('primary', RGBColor(24, 58, 88))
        title_font = typography.get('primary_font', 'Arial')
        title_size = typography.get('title_size', Pt(24))
    else:
        primary_color = RGBColor(24, 58, 88)  # Default blue
        title_font = 'Arial'
        title_size = Pt(24)
    
    # STANDARDIZED: Left-aligned title (no center alignment anywhere)
    title_left = Inches(0.5)
    title_top = Inches(0.3)
    title_width = Inches(12)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.margin_left = Inches(0)
    title_frame.margin_top = Inches(0)
    title_frame.margin_right = Inches(0)
    title_frame.margin_bottom = Inches(0)
    
    title_p = title_frame.paragraphs[0]
    title_p.text = title_text
    title_p.alignment = PP_ALIGN.LEFT  # ALWAYS LEFT ALIGNED
    
    title_run = title_p.runs[0]
    title_run.font.name = title_font
    title_run.font.size = title_size
    title_run.font.bold = True
    title_run.font.color.rgb = primary_color
    
    # STANDARDIZED: Blue underline (consistent across all slides)
    underline_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        title_left, Inches(1.1), Inches(12.3), Inches(0.05)
    )
    underline_shape.fill.solid()
    underline_shape.fill.fore_color.rgb = primary_color
    underline_shape.line.fill.background()
    
    return title_box, underline_shape

# ---- helpers ----

def _ensure_prs(prs=None):
    """Return a python-pptx Presentation. If 'prs' is None or not a Presentation, create a new one."""
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed. Please `pip install python-pptx`.")
    try:
        from pptx import Presentation as _P
    except Exception:
        # If import fails oddly, just create new using already imported Presentation
        _P = Presentation  # type: ignore
    if prs is None or (hasattr(_P, "__call__") and not hasattr(prs, "slides")):
        return Presentation()
    # Some callers pass truthy sentinels; verify it looks like a pptx.Presentation
    if not hasattr(prs, "slides") or not hasattr(prs, "slide_layouts"):
        return Presentation()
    return prs

def _coerce_plan(plan: Optional[Dict]=None, content: Optional[Any]=None, content_ir: Optional[Any]=None) -> Dict[str, List[Dict]]:
    """
    Normalize different inputs into a dict with a 'slides' list of {template, data}.
    """
    src = plan if plan is not None else (content if content is not None else content_ir)
    if src is None:
        raise ValueError("No plan/content/content_ir provided.")
    # Accept already-correct shape
    if isinstance(src, dict) and "slides" in src and isinstance(src["slides"], list):
        return src  # type: ignore
    # Accept list of slides directly
    if isinstance(src, list):
        return {"slides": src}  # type: ignore
    # Some wrappers store slides under other key
    for key in ("plan", "render_plan", "content", "content_ir"):
        if isinstance(src, dict) and key in src and isinstance(src[key], list):
            return {"slides": src[key]}
    raise ValueError("Unrecognized plan/content shape. Expect dict with 'slides' or a list of slide dicts.")

def _safe_call(renderer, data: Dict, prs, company_name: str, content: Dict = None, brand_config: Optional[Dict] = None):
    """
    Call renderer with the correct parameter signature for your slide templates.
    Your renderers expect: renderer(slide_data, color_scheme=None, typography=None, company_name="Moelis", prs=None)
    Now passes brand configuration when available.
    """
    if renderer is None:
        # Fallback: ensure we at least add a blank slide so user sees progress
        from pptx.util import Inches
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        tf = tb.text_frame
        tf.text = f"[Adapter notice] No renderer found for this template."
        return prs

    try:
        # Extract brand configuration components
        color_scheme = None
        typography = None
        
        if brand_config:
            color_scheme = _convert_brand_colors(brand_config)
            typography = _standardize_typography(brand_config)
        
        # Your renderers expect slide_data as POSITIONAL parameter, not keyword
        # Pass brand configuration if available
        result = renderer(
            data, 
            color_scheme=color_scheme,
            typography=typography, 
            company_name=company_name, 
            prs=prs,
            brand_config=brand_config  # Pass full brand config for header standardization
        )
        
    except TypeError as te:
        try:
            # Try without brand_config in case some renderers don't support it yet
            result = renderer(data, color_scheme=None, typography=None, company_name=company_name, prs=prs)
        except TypeError as te2:
            try:
                # Try minimal signature
                result = renderer(data, prs=prs)
            except Exception as e:
                # On renderer failure, add a diagnostic slide
                from pptx.util import Inches, Pt
                from pptx.enum.text import PP_ALIGN
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(11.5), Inches(1.5))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.text = f"Renderer error for {renderer.__name__ if hasattr(renderer, '__name__') else 'unknown'}: {e}"
                for r in p.runs:
                    r.font.size = Pt(14)
                p.alignment = PP_ALIGN.LEFT
                print(f"Error calling renderer: {e}")
                return prs
    except Exception as e:
        # On renderer failure, add a diagnostic slide
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(11.5), Inches(1.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = f"Renderer error: {e}"
        for r in p.runs:
            r.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        print(f"Error in renderer: {e}")
        return prs

    # Most of your renderers return the same prs; if not, preserve new prs
    return result if hasattr(result, "slides") else prs

# ---- main entrypoint ----

def render_plan_to_pptx(
    plan: Optional[Dict]=None,
    content: Optional[Any]=None,
    content_ir: Optional[Any]=None,
    prs=None,
    company_name: str = "Moelis",
    brand_config: Optional[Dict] = None,  # NEW: Brand configuration
    **_ignore_kwargs,
):
    """
    Renders all slides in the normalized plan into a python-pptx Presentation.
    Returns the Presentation.
    Extra kwargs are ignored for forward compatibility.
    Now supports brand configuration for consistent styling.
    """
    prs = _ensure_prs(prs)
    plan_obj = _coerce_plan(plan=plan, content=content, content_ir=content_ir)

    slides = plan_obj.get("slides", [])
    if not isinstance(slides, list):
        raise ValueError("Plan 'slides' must be a list.")

    # Convert content to dict if provided
    content_dict = {}
    if isinstance(content, dict):
        content_dict = content
    elif isinstance(content_ir, dict):
        content_dict = content_ir

    print(f"[DEBUG] Processing {len(slides)} slides")
    if brand_config:
        print(f"[DEBUG] Using custom brand configuration")
    
    for idx, item in enumerate(slides, start=1):
        if not isinstance(item, dict):
            print(f"[DEBUG] Slide {idx}: Not a dict, skipping")
            continue
            
        template = item.get("template")
        data = item.get("data", {})
        
        print(f"[DEBUG] Slide {idx}: template='{template}', data keys={list(data.keys()) if isinstance(data, dict) else 'not dict'}")
        
        # If slide references content by ID, resolve it
        content_id = item.get("content_id")
        if content_id and content_dict:
            content_data = content_dict.get(content_id, {})
            if content_data:
                # Merge content data with any existing data (existing data takes precedence)
                merged_data = {**content_data, **data}
                data = merged_data
                print(f"[DEBUG] Slide {idx}: Merged content for ID '{content_id}', new data keys={list(data.keys())}")
        
        if template is None:
            print(f"[DEBUG] Slide {idx}: No template, skipping")
            continue

        renderer = RENDERER_MAP.get(template)
        if renderer is None:
            print(f"[DEBUG] Slide {idx}: No renderer found for template '{template}'")
        else:
            print(f"[DEBUG] Slide {idx}: Found renderer for '{template}': {renderer.__name__ if hasattr(renderer, '__name__') else str(renderer)}")
            
        prs = _safe_call(renderer, data, prs, company_name, content_dict, brand_config)

    print(f"[DEBUG] Finished processing. Total slides in presentation: {len(prs.slides)}")
    return prs
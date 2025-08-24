"""
Enhanced Brand Extractor with Theme Color Support
This version can extract and convert PowerPoint theme colors to RGB values
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.dml import MSO_COLOR_TYPE
import io
from typing import Dict, Optional, Tuple, List
import logging

class EnhancedBrandExtractor:
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    def extract_brand_from_pptx(self, pptx_file) -> Dict:
        """
        Extract brand elements including theme colors from uploaded PowerPoint file
        """
        try:
            # Handle different input types
            if hasattr(pptx_file, 'read'):
                pptx_file.seek(0)
                prs = Presentation(io.BytesIO(pptx_file.read()))
            else:
                prs = Presentation(pptx_file)
            
            print(f"[THEME DEBUG] Starting enhanced extraction from {len(prs.slides)} slides...")
            
            # Extract theme colors first (most reliable)
            theme_colors = self._extract_theme_colors_comprehensive(prs)
            
            # Extract direct RGB colors as backup
            rgb_colors = self._extract_direct_rgb_colors(prs)
            
            # Extract colors from actual usage on slides
            usage_colors = self._extract_colors_from_usage(prs)
            
            # Combine all color sources (theme takes priority)
            final_colors = self._combine_color_sources(theme_colors, rgb_colors, usage_colors)
            
            brand_config = {
                'color_scheme': final_colors,
                'typography': self._extract_fonts(prs),
                'header_style': self._extract_header_style(prs),
                'layout_config': {
                    'title_alignment': 'left',
                    'header_type': 'extracted'
                }
            }
            
            print(f"[THEME DEBUG] Final extracted colors:")
            for name, color in final_colors.items():
                if hasattr(color, 'r'):
                    print(f"  {name}: RGB({color.r}, {color.g}, {color.b}) = #{color.r:02x}{color.g:02x}{color.b:02x}")
            
            return brand_config
            
        except Exception as e:
            print(f"[THEME DEBUG] Enhanced extraction failed: {str(e)}")
            import traceback
            traceback.print_exc()
            return self._get_default_brand_config()
    
    def _extract_theme_colors_comprehensive(self, prs: Presentation) -> Dict[str, RGBColor]:
        """Extract theme colors using multiple approaches"""
        theme_colors = {}
        
        print("[THEME DEBUG] Attempting theme color extraction...")
        
        # Method 1: Extract from slide masters
        theme_colors.update(self._extract_from_slide_masters(prs))
        
        # Method 2: Extract from theme part (if accessible)
        theme_colors.update(self._extract_from_theme_part(prs))
        
        # Method 3: Extract by analyzing theme color usage
        theme_colors.update(self._extract_theme_usage_colors(prs))
        
        return theme_colors
    
    def _extract_from_slide_masters(self, prs: Presentation) -> Dict[str, RGBColor]:
        """Extract colors from slide masters"""
        colors = {}
        
        try:
            print("[THEME DEBUG] Checking slide masters...")
            
            if hasattr(prs, 'slide_masters'):
                for master_idx, master in enumerate(prs.slide_masters):
                    print(f"[THEME DEBUG] Slide master {master_idx + 1}:")
                    
                    # Check master background
                    if hasattr(master, 'background') and hasattr(master.background, 'fill'):
                        bg_color = self._extract_color_from_fill(master.background.fill)
                        if bg_color:
                            colors['background'] = bg_color
                            print(f"  Background from master: RGB({bg_color.r}, {bg_color.g}, {bg_color.b})")
                    
                    # Check shapes in master
                    for shape in master.shapes:
                        shape_colors = self._extract_all_colors_from_shape(shape, f"Master shape")
                        colors.update(shape_colors)
                        
        except Exception as e:
            print(f"[THEME DEBUG] Slide master extraction error: {e}")
        
        return colors
    
    def _extract_from_theme_part(self, prs: Presentation) -> Dict[str, RGBColor]:
        """Extract colors directly from theme XML"""
        colors = {}
        
        try:
            print("[THEME DEBUG] Attempting theme part extraction...")
            
            # Access theme through presentation part
            if hasattr(prs, 'part'):
                presentation_part = prs.part
                
                # Look for theme in related parts
                if hasattr(presentation_part, 'related_parts'):
                    for part_key, part in presentation_part.related_parts.items():
                        print(f"[THEME DEBUG] Checking part: {type(part).__name__}")
                        
                        # Check if this is a theme part
                        if 'theme' in str(type(part)).lower():
                            print(f"[THEME DEBUG] Found theme part: {part}")
                            theme_colors = self._parse_theme_part(part)
                            colors.update(theme_colors)
                            
        except Exception as e:
            print(f"[THEME DEBUG] Theme part extraction error: {e}")
        
        return colors
    
    def _parse_theme_part(self, theme_part) -> Dict[str, RGBColor]:
        """Parse theme part to extract colors"""
        colors = {}
        
        try:
            # Try to access theme element
            if hasattr(theme_part, 'element'):
                theme_element = theme_part.element
                print(f"[THEME DEBUG] Theme element: {theme_element}")
                
                # Look for color scheme in theme
                # This is highly dependent on the XML structure
                for child in theme_element:
                    print(f"[THEME DEBUG] Theme child: {child.tag if hasattr(child, 'tag') else child}")
                    
                    # Look for color scheme elements
                    if 'colorScheme' in str(child.tag) or 'clrScheme' in str(child.tag):
                        print("[THEME DEBUG] Found color scheme element!")
                        colors.update(self._parse_color_scheme_element(child))
                        
        except Exception as e:
            print(f"[THEME DEBUG] Theme part parsing error: {e}")
        
        return colors
    
    def _parse_color_scheme_element(self, color_scheme_element) -> Dict[str, RGBColor]:
        """Parse color scheme XML element"""
        colors = {}
        
        try:
            # Standard PowerPoint theme color names
            theme_color_mapping = {
                'accent1': 'primary',
                'accent2': 'secondary', 
                'accent3': 'accent',
                'dk1': 'text',
                'lt1': 'background',
                'dk2': 'text',
                'lt2': 'light_grey'
            }
            
            for child in color_scheme_element:
                color_name = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                print(f"[THEME DEBUG] Found theme color: {color_name}")
                
                if color_name in theme_color_mapping:
                    rgb_color = self._parse_color_element(child)
                    if rgb_color:
                        mapped_name = theme_color_mapping[color_name]
                        colors[mapped_name] = rgb_color
                        print(f"[THEME DEBUG] Mapped {color_name} -> {mapped_name}: RGB({rgb_color.r}, {rgb_color.g}, {rgb_color.b})")
                        
        except Exception as e:
            print(f"[THEME DEBUG] Color scheme parsing error: {e}")
        
        return colors
    
    def _parse_color_element(self, color_element) -> Optional[RGBColor]:
        """Parse individual color element to RGB"""
        try:
            # Look for RGB values in various formats
            for child in color_element:
                child_tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                
                # Direct RGB color
                if child_tag == 'srgbClr' and 'val' in child.attrib:
                    hex_val = child.attrib['val']
                    r = int(hex_val[0:2], 16)
                    g = int(hex_val[2:4], 16)
                    b = int(hex_val[4:6], 16)
                    return RGBColor(r, g, b)
                
                # System color
                elif child_tag == 'sysClr' and 'lastClr' in child.attrib:
                    hex_val = child.attrib['lastClr']
                    r = int(hex_val[0:2], 16)
                    g = int(hex_val[2:4], 16)
                    b = int(hex_val[4:6], 16)
                    return RGBColor(r, g, b)
                    
        except Exception as e:
            print(f"[THEME DEBUG] Color element parsing error: {e}")
        
        return None
    
    def _extract_theme_usage_colors(self, prs: Presentation) -> Dict[str, RGBColor]:
        """Extract theme colors by analyzing how they're used on slides"""
        colors = {}
        theme_color_usage = {}
        
        try:
            print("[THEME DEBUG] Analyzing theme color usage...")
            
            for slide_idx, slide in enumerate(prs.slides[:5]):  # Check first 5 slides
                print(f"[THEME DEBUG] Slide {slide_idx + 1}:")
                
                for shape in slide.shapes:
                    # Check for theme color references
                    theme_refs = self._get_theme_color_references(shape)
                    for theme_name, rgb_color in theme_refs.items():
                        if theme_name not in theme_color_usage:
                            theme_color_usage[theme_name] = []
                        theme_color_usage[theme_name].append(rgb_color)
                        print(f"  Found theme color {theme_name}: RGB({rgb_color.r}, {rgb_color.g}, {rgb_color.b})")
            
            # Convert most common theme colors to our color scheme
            colors = self._map_theme_usage_to_scheme(theme_color_usage)
            
        except Exception as e:
            print(f"[THEME DEBUG] Theme usage analysis error: {e}")
        
        return colors
    
    def _get_theme_color_references(self, shape) -> Dict[str, RGBColor]:
        """Get theme color references from a shape"""
        theme_colors = {}
        
        try:
            # Check fill color
            if hasattr(shape, 'fill'):
                theme_color = self._get_theme_color_from_fill(shape.fill)
                if theme_color:
                    theme_colors.update(theme_color)
            
            # Check text colors  
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run.font, 'color'):
                            theme_color = self._get_theme_color_from_color_format(run.font.color)
                            if theme_color:
                                theme_colors.update(theme_color)
                                
        except Exception as e:
            print(f"[THEME DEBUG] Theme reference extraction error: {e}")
        
        return theme_colors
    
    def _get_theme_color_from_fill(self, fill) -> Optional[Dict[str, RGBColor]]:
        """Extract theme color from fill"""
        try:
            if hasattr(fill, 'fore_color'):
                return self._get_theme_color_from_color_format(fill.fore_color)
        except Exception:
            pass
        return None
    
    def _get_theme_color_from_color_format(self, color_format) -> Optional[Dict[str, RGBColor]]:
        """Extract theme color from color format object"""
        try:
            # Check if it's a theme color
            if hasattr(color_format, 'type') and color_format.type == MSO_COLOR_TYPE.THEME:
                if hasattr(color_format, 'theme_color') and hasattr(color_format, 'rgb'):
                    theme_name = str(color_format.theme_color)
                    rgb_color = color_format.rgb
                    return {theme_name: rgb_color}
            
            # Even if not explicitly theme, try to get RGB
            elif hasattr(color_format, 'rgb'):
                rgb_color = color_format.rgb
                # Guess theme role based on color properties
                if rgb_color.r > rgb_color.g and rgb_color.r > rgb_color.b:
                    return {'accent1': rgb_color}  # Likely primary color
                    
        except Exception as e:
            print(f"[THEME DEBUG] Color format analysis error: {e}")
        
        return None
    
    def _map_theme_usage_to_scheme(self, theme_usage: Dict[str, List[RGBColor]]) -> Dict[str, RGBColor]:
        """Map theme color usage to our color scheme"""
        colors = {}
        
        # Map common theme color names to our scheme
        theme_mapping = {
            'accent1': 'primary',
            'accent2': 'secondary',
            'accent3': 'accent',
            'ACCENT_1': 'primary',
            'ACCENT_2': 'secondary', 
            'ACCENT_3': 'accent',
            '4': 'primary',  # Sometimes theme colors are just numbers
            '5': 'secondary',
            '6': 'accent'
        }
        
        for theme_name, rgb_list in theme_usage.items():
            # Take the most common color for this theme
            if rgb_list:
                most_common_rgb = rgb_list[0]  # Simplified - could count frequencies
                
                # Map to our scheme
                mapped_name = theme_mapping.get(theme_name, 'accent')
                colors[mapped_name] = most_common_rgb
                print(f"[THEME DEBUG] Mapped theme {theme_name} -> {mapped_name}")
        
        return colors
    
    def _extract_direct_rgb_colors(self, prs: Presentation) -> Dict[str, RGBColor]:
        """Extract direct RGB colors (fallback method)"""
        return self._extract_colors_from_slides(prs)
    
    def _extract_colors_from_usage(self, prs: Presentation) -> Dict[str, RGBColor]:
        """Extract colors based on how they're used (size, position, etc.)"""
        colors = {}
        color_candidates = []
        
        try:
            for slide in prs.slides[:5]:
                for shape in slide.shapes:
                    shape_colors = self._extract_all_colors_from_shape(shape, "usage")
                    for role, color in shape_colors.items():
                        # Determine usage context
                        usage_context = self._determine_usage_context(shape, prs)
                        color_candidates.append((role, color, usage_context))
            
            # Analyze candidates and pick best colors
            colors = self._select_best_color_candidates(color_candidates)
            
        except Exception as e:
            print(f"[THEME DEBUG] Usage-based extraction error: {e}")
        
        return colors
    
    def _extract_all_colors_from_shape(self, shape, context="") -> Dict[str, RGBColor]:
        """Extract all possible colors from a shape"""
        colors = {}
        
        try:
            # Fill color
            if hasattr(shape, 'fill'):
                fill_color = self._extract_color_from_fill(shape.fill)
                if fill_color and not self._is_default_color(fill_color):
                    colors['fill'] = fill_color
                    print(f"[THEME DEBUG] {context} fill: RGB({fill_color.r}, {fill_color.g}, {fill_color.b})")
            
            # Text colors
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb'):
                            text_color = run.font.color.rgb
                            if not self._is_default_color(text_color):
                                colors['text'] = text_color
                                print(f"[THEME DEBUG] {context} text: RGB({text_color.r}, {text_color.g}, {text_color.b})")
            
            # Line color
            if hasattr(shape, 'line') and hasattr(shape.line, 'color'):
                if hasattr(shape.line.color, 'rgb'):
                    line_color = shape.line.color.rgb
                    if not self._is_default_color(line_color):
                        colors['line'] = line_color
                        print(f"[THEME DEBUG] {context} line: RGB({line_color.r}, {line_color.g}, {line_color.b})")
                        
        except Exception as e:
            print(f"[THEME DEBUG] Shape color extraction error: {e}")
        
        return colors
    
    def _extract_color_from_fill(self, fill) -> Optional[RGBColor]:
        """Extract color from fill object"""
        try:
            if hasattr(fill, 'type') and fill.type == 1:  # Solid fill
                if hasattr(fill, 'fore_color') and hasattr(fill.fore_color, 'rgb'):
                    return fill.fore_color.rgb
        except Exception:
            pass
        return None
    
    def _determine_usage_context(self, shape, prs) -> str:
        """Determine how a shape is being used (header, body, accent, etc.)"""
        try:
            # Large shapes are likely headers/primary
            if hasattr(shape, 'width') and hasattr(prs, 'slide_width'):
                if shape.width > prs.slide_width * 0.6:
                    return 'header'
            
            # Small shapes are likely accents
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                if (shape.width < prs.slide_width * 0.1 and 
                    shape.height < prs.slide_height * 0.1):
                    return 'accent'
            
            # Text-heavy shapes are body
            if hasattr(shape, 'text_frame') and shape.text_frame:
                if len(shape.text_frame.text) > 50:
                    return 'body'
                    
            return 'general'
            
        except Exception:
            return 'general'
    
    def _select_best_color_candidates(self, candidates) -> Dict[str, RGBColor]:
        """Select the best color candidates for our scheme"""
        colors = {}
        
        # Group by usage context
        header_colors = [c[1] for c in candidates if c[2] == 'header']
        accent_colors = [c[1] for c in candidates if c[2] == 'accent']
        body_colors = [c[1] for c in candidates if c[2] == 'body']
        
        # Pick representative colors
        if header_colors:
            colors['primary'] = header_colors[0]
        if accent_colors:
            colors['secondary'] = accent_colors[0]
        if body_colors:
            colors['text'] = body_colors[0]
            
        return colors
    
    def _combine_color_sources(self, theme_colors, rgb_colors, usage_colors) -> Dict[str, RGBColor]:
        """Combine colors from all sources with priority order"""
        # Start with defaults
        final_colors = {
            'primary': RGBColor(24, 58, 88),
            'secondary': RGBColor(181, 151, 91),
            'accent': RGBColor(64, 64, 64),
            'background': RGBColor(255, 255, 255),
            'text': RGBColor(64, 64, 64),
            'light_grey': RGBColor(240, 240, 240),
            'footer_grey': RGBColor(128, 128, 128)
        }
        
        # Apply in order of preference: usage -> rgb -> theme (theme has highest priority)
        for colors in [usage_colors, rgb_colors, theme_colors]:
            for name, color in colors.items():
                if name in final_colors and color:
                    final_colors[name] = color
                    print(f"[THEME DEBUG] Applied {name}: RGB({color.r}, {color.g}, {color.b})")
        
        return final_colors
    
    def _is_default_color(self, color: RGBColor) -> bool:
        """Check if color should be ignored"""
        if not color:
            return True
        return ((color.r == 255 and color.g == 255 and color.b == 255) or  # White
                (color.r == 0 and color.g == 0 and color.b == 0) or          # Black
                (color.r > 240 and color.g > 240 and color.b > 240))         # Very light
    
    def _extract_colors_from_slides(self, prs: Presentation) -> Dict[str, RGBColor]:
        """Basic slide color extraction (existing method)"""
        colors = {}
        
        try:
            for slide_idx, slide in enumerate(prs.slides[:5]):
                for shape in slide.shapes:
                    shape_colors = self._extract_all_colors_from_shape(shape, f"Slide {slide_idx + 1}")
                    
                    # Map first colors found to scheme
                    if 'fill' in shape_colors and 'primary' not in colors:
                        colors['primary'] = shape_colors['fill']
                    if 'text' in shape_colors and 'text' not in colors:
                        colors['text'] = shape_colors['text']
                    if 'line' in shape_colors and 'accent' not in colors:
                        colors['accent'] = shape_colors['line']
                        
        except Exception as e:
            print(f"[THEME DEBUG] Basic slide extraction error: {e}")
        
        return colors
    
    def _extract_fonts(self, prs: Presentation) -> Dict:
        """Extract typography (existing method)"""
        return {
            'primary_font': 'Arial',
            'title_size': 24,
            'header_size': 14,
            'body_size': 11,
            'small_size': 9
        }
    
    def _extract_header_style(self, prs: Presentation) -> Dict:
        """Extract header style (existing method)"""
        return {
            'type': 'line',
            'height': 0.05,
            'color': 'primary',
            'has_logo': False
        }
    
    def _get_default_brand_config(self) -> Dict:
        """Default config fallback"""
        return {
            'color_scheme': {
                'primary': RGBColor(24, 58, 88),
                'secondary': RGBColor(181, 151, 91),
                'accent': RGBColor(64, 64, 64),
                'background': RGBColor(255, 255, 255),
                'text': RGBColor(64, 64, 64),
                'light_grey': RGBColor(240, 240, 240),
                'footer_grey': RGBColor(128, 128, 128)
            },
            'typography': {
                'primary_font': 'Arial',
                'title_size': 24,
                'header_size': 14,
                'body_size': 11,
                'small_size': 9
            },
            'header_style': {
                'type': 'line',
                'height': 0.05,
                'color': 'primary',
                'has_logo': False
            },
            'layout_config': {
                'title_alignment': 'left',
                'header_type': 'line'
            }
        }
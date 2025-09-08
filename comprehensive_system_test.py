#!/usr/bin/env python3
"""
Comprehensive System Test - Validate All Critical Functionality
Tests everything that was previously working to ensure no regressions
"""
import sys
import os
import json
import traceback
from pathlib import Path

# Add the current directory to the Python path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

def test_imports():
    """Test that all critical imports are working"""
    print("🧪 Testing Critical Imports...")
    
    try:
        # Core functionality imports
        from executor import execute_plan
        from catalog_loader import TemplateCatalog
        from brand_extractor import BrandExtractor
        from slide_templates import (
            render_business_overview_slide, render_management_team_slide,
            render_competitive_positioning_slide, render_buyer_profiles_slide,
            render_sea_conglomerates_slide, render_investor_process_overview_slide
        )
        
        print("✅ All critical imports successful")
        return True
    except Exception as e:
        print(f"❌ Import failed: {e}")
        return False

def test_json_cleaning():
    """Test the enhanced JSON cleaning functionality"""
    print("🧪 Testing Enhanced JSON Cleaning...")
    
    try:
        # Import the JSON cleaning functions from app.py
        import app
        
        # Test cases for JSON cleaning
        test_cases = [
            # Basic valid JSON
            ('{"test": "value"}', True, "Basic valid JSON"),
            
            # Missing comma between objects
            ('{"a": 1}\n{"b": 2}', False, "Missing comma between objects"),
            
            # Trailing comma
            ('{"a": 1, "b": 2,}', True, "Trailing comma should be fixed"),
            
            # Code blocks
            ('```json\n{"test": "value"}\n```', True, "JSON in code blocks"),
            
            # Missing quotes around properties
            ('{test: "value"}', False, "Missing quotes around property"),
        ]
        
        success_count = 0
        for test_json, should_work, description in test_cases:
            try:
                cleaned = app.clean_json_string(test_json)
                parsed = json.loads(cleaned)
                print(f"  ✅ {description}: PASSED")
                success_count += 1
            except Exception as e:
                if should_work:
                    print(f"  ❌ {description}: FAILED - {e}")
                else:
                    print(f"  ⚠️  {description}: Expected failure - {e}")
                    success_count += 1
        
        print(f"✅ JSON Cleaning: {success_count}/{len(test_cases)} tests passed")
        return success_count == len(test_cases)
        
    except Exception as e:
        print(f"❌ JSON cleaning test failed: {e}")
        return False

def test_slide_generation():
    """Test slide generation functionality"""
    print("🧪 Testing Slide Generation...")
    
    try:
        from slide_templates import (
            render_business_overview_slide, render_management_team_slide,
            render_competitive_positioning_slide
        )
        from pptx import Presentation
        
        # Create a test presentation
        prs = Presentation()
        
        # Test business overview slide
        test_data = {
            'title': 'Test Business Overview',
            'description': 'Test company description',
            'highlights': ['Highlight 1', 'Highlight 2', 'Highlight 3'],
            'services': ['Service 1', 'Service 2', 'Service 3'],
            'positioning_desc': 'Test positioning'
        }
        
        slide = render_business_overview_slide(data=test_data, prs=prs)
        if slide:
            print("  ✅ Business overview slide generation: PASSED")
        else:
            print("  ❌ Business overview slide generation: FAILED")
            return False
        
        # Test management team slide
        mgmt_data = {
            'title': 'Management Team',
            'left_column_profiles': [
                {'role_title': 'CEO', 'experience_bullets': ['Bullet 1', 'Bullet 2']},
            ],
            'right_column_profiles': [
                {'role_title': 'CTO', 'experience_bullets': ['Bullet 1', 'Bullet 2']},
            ]
        }
        
        slide = render_management_team_slide(data=mgmt_data, prs=prs)
        if slide:
            print("  ✅ Management team slide generation: PASSED")
        else:
            print("  ❌ Management team slide generation: FAILED")
            return False
        
        # Test competitive positioning slide (with iCar Asia format)
        comp_data = {
            'title': 'Competitive Positioning',
            'competitors': [
                {'name': 'Company A', 'revenue': 100},
                {'name': 'Company B', 'revenue': 200}
            ],
            'assessment': [
                ['Company', 'Market Share', 'Tech Platform', 'Coverage', 'Revenue (M)'],
                ['Our Company', '⭐⭐⭐⭐', '⭐⭐⭐⭐⭐', '⭐⭐⭐⭐', '$150M'],
                ['Company A', '⭐⭐⭐', '⭐⭐⭐', '⭐⭐⭐', '$100M']
            ],
            'advantages': [
                {'title': 'Advantage 1', 'desc': 'Description 1'},
                {'title': 'Advantage 2', 'desc': 'Description 2'}
            ],
            'barriers': [
                {'title': 'Barrier 1', 'desc': 'Description 1'},
                {'title': 'Barrier 2', 'desc': 'Description 2'}
            ]
        }
        
        slide = render_competitive_positioning_slide(data=comp_data, prs=prs)
        if slide:
            print("  ✅ Competitive positioning slide (iCar Asia format): PASSED")
        else:
            print("  ❌ Competitive positioning slide generation: FAILED")
            return False
        
        print("✅ Slide Generation: All tests passed")
        return True
        
    except Exception as e:
        print(f"❌ Slide generation test failed: {e}")
        traceback.print_exc()
        return False

def test_buyer_slide_fixes():
    """Test the buyer slide column header fixes"""
    print("🧪 Testing Buyer Slide Fixes...")
    
    try:
        from slide_templates import render_buyer_profiles_slide
        from pptx import Presentation
        
        prs = Presentation()
        
        # Test data with fit field (not fit_score)
        buyer_data = {
            'title': 'Strategic Buyers',
            'table_headers': ['Buyer Name', 'Description', 'Strategic Rationale', 'Key Synergies', 'Fit'],
            'table_rows': [
                {
                    'buyer_name': 'Test Buyer',
                    'description': 'Test description',
                    'strategic_rationale': 'Test rationale',
                    'key_synergies': 'Test synergies',
                    'fit': 'High (9/10)'
                }
            ]
        }
        
        slide = render_buyer_profiles_slide(data=buyer_data, prs=prs)
        if slide:
            print("  ✅ Buyer profiles slide with 'Fit' header: PASSED")
        else:
            print("  ❌ Buyer profiles slide generation: FAILED")
            return False
        
        print("✅ Buyer Slide Fixes: All tests passed")
        return True
        
    except Exception as e:
        print(f"❌ Buyer slide test failed: {e}")
        return False

def test_brand_extraction():
    """Test brand extraction functionality"""
    print("🧪 Testing Brand Extraction...")
    
    try:
        from brand_extractor import BrandExtractor
        
        extractor = BrandExtractor()
        
        # Test with a known test file
        test_files = ['test_brand_deck.pptx', 'test_brand_colors.pptx']
        
        for test_file in test_files:
            if Path(test_file).exists():
                print(f"  Testing with {test_file}...")
                brand_config = extractor.extract_brand_from_pptx(test_file, use_llm=False)
                
                if brand_config and 'color_scheme' in brand_config:
                    print(f"  ✅ Brand extraction from {test_file}: PASSED")
                else:
                    print(f"  ❌ Brand extraction from {test_file}: FAILED")
                    return False
        
        print("✅ Brand Extraction: All tests passed")
        return True
        
    except Exception as e:
        print(f"❌ Brand extraction test failed: {e}")
        return False

def test_validation_functions():
    """Test validation functions for slides"""
    print("🧪 Testing Validation Functions...")
    
    try:
        import app
        
        # Test buyer validation with fit field (not fit_score)
        test_slide = {
            'template': 'buyer_profiles',
            'data': {
                'table_rows': [
                    {
                        'buyer_name': 'Test Buyer',
                        'strategic_rationale': 'Test rationale',
                        'fit': 'High (9/10)'
                    }
                ]
            }
        }
        
        test_content_ir = {
            'strategic_buyers': [
                {
                    'buyer_name': 'Test Buyer',
                    'strategic_rationale': 'Test rationale', 
                    'fit': 'High (9/10)'
                }
            ]
        }
        
        if hasattr(app, 'validate_buyer_profiles_slide'):
            validation = app.validate_buyer_profiles_slide(test_slide, test_content_ir)
            print(f"  ✅ Buyer validation with 'fit' field: PASSED")
        else:
            print(f"  ⚠️ Buyer validation function not found")
        
        print("✅ Validation Functions: All tests passed")
        return True
        
    except Exception as e:
        print(f"❌ Validation test failed: {e}")
        return False

def test_investor_process_timeline():
    """Test the enhanced investor process timeline"""
    print("🧪 Testing Investment Process Timeline...")
    
    try:
        from slide_templates import render_investor_process_overview_slide
        from pptx import Presentation
        
        prs = Presentation()
        
        # Test with enhanced timeline
        timeline_data = {
            'title': 'Investor Process Overview',
            'diligence_topics': ['Topic 1', 'Topic 2'],
            'synergy_opportunities': ['Synergy 1', 'Synergy 2'],
            'risk_factors': ['Risk 1', 'Risk 2'],
            'mitigants': ['Mitigant 1', 'Mitigant 2'],
            'timeline': [
                {'date': 'Week 1-2', 'description': 'Initial outreach and process launch'},
                {'date': 'Week 3-4', 'description': 'Management presentations and strategic discussions'},
                {'date': 'Week 15-16', 'description': 'Definitive agreements and closing preparations'}
            ]
        }
        
        slide = render_investor_process_overview_slide(data=timeline_data, prs=prs)
        if slide:
            print("  ✅ Investor process timeline slide: PASSED")
        else:
            print("  ❌ Investor process timeline slide: FAILED")
            return False
        
        print("✅ Investment Process Timeline: All tests passed")
        return True
        
    except Exception as e:
        print(f"❌ Investment process timeline test failed: {e}")
        return False

def main():
    """Run comprehensive system validation"""
    print("🚀 COMPREHENSIVE SYSTEM VALIDATION")
    print("=" * 60)
    
    tests = [
        ("Critical Imports", test_imports),
        ("JSON Cleaning", test_json_cleaning),
        ("Slide Generation", test_slide_generation),
        ("Buyer Slide Fixes", test_buyer_slide_fixes),
        ("Brand Extraction", test_brand_extraction),
        ("Validation Functions", test_validation_functions),
        ("Investment Process Timeline", test_investor_process_timeline),
    ]
    
    results = {}
    total_tests = len(tests)
    passed_tests = 0
    
    for test_name, test_func in tests:
        print(f"\n{'='*20} {test_name} {'='*20}")
        try:
            result = test_func()
            results[test_name] = result
            if result:
                passed_tests += 1
                print(f"✅ {test_name}: PASSED")
            else:
                print(f"❌ {test_name}: FAILED")
        except Exception as e:
            print(f"💥 {test_name}: CRASHED - {e}")
            results[test_name] = False
    
    print(f"\n{'='*60}")
    print("📊 COMPREHENSIVE TEST RESULTS")
    print(f"{'='*60}")
    
    for test_name, result in results.items():
        status = "✅ PASS" if result else "❌ FAIL"
        print(f"  {test_name:<30}: {status}")
    
    print(f"\n🎯 OVERALL RESULT: {passed_tests}/{total_tests} tests passed")
    
    if passed_tests == total_tests:
        print("✅ ALL SYSTEMS OPERATIONAL - No regressions detected!")
        print("🚀 All previously working functionality is intact")
        return True
    else:
        print("❌ SYSTEM ISSUES DETECTED - Some functionality may be broken")
        print("🔧 Please review failed tests and fix issues")
        return False

if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)
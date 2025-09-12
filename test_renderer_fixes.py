#!/usr/bin/env python3
"""
Test that slide renderers handle the fixed data structure properly
"""

from slide_templates import render_competitive_positioning_slide, render_precedent_transactions_slide
from pptx import Presentation
import json

def test_competitive_renderer():
    """Test competitive positioning slide with dynamic headers"""
    print("🔧 Testing competitive positioning renderer...")
    
    test_data = {
        "title": "Competitive Positioning",
        "competitors": [
            {"name": "NVIDIA", "revenue": 60900},
            {"name": "AMD", "revenue": 23500},
            {"name": "Intel", "revenue": 76000}
        ],
        "assessment": [
            ["Company", "AI Computing Leadership", "Software Ecosystem", "Data Center Penetration", "Innovation Velocity"],
            ["NVIDIA", "⭐⭐⭐⭐⭐", "⭐⭐⭐⭐⭐", "⭐⭐⭐⭐⭐", "⭐⭐⭐⭐⭐"],
            ["AMD", "⭐⭐⭐", "⭐⭐", "⭐⭐⭐", "⭐⭐⭐⭐"],
            ["Intel", "⭐⭐", "⭐⭐", "⭐⭐⭐⭐", "⭐⭐⭐"]
        ],
        "barriers": [
            {"title": "CUDA Ecosystem Lock-in", "desc": "4M+ developers create high switching costs"}
        ],
        "advantages": [
            {"title": "AI Computing Dominance", "desc": "90%+ market share in AI training chips"}
        ]
    }
    
    try:
        prs = render_competitive_positioning_slide(
            data=test_data,
            company_name="NVIDIA"
        )
        
        if prs and len(prs.slides) > 0:
            print("✅ Competitive positioning slide rendered successfully")
            print(f"   - Assessment matrix: {len(test_data['assessment'])} rows x {len(test_data['assessment'][0])} cols")
            print(f"   - Dynamic headers: {test_data['assessment'][0]}")
            return True
        else:
            print("❌ Competitive positioning slide rendering failed")
            return False
            
    except Exception as e:
        print(f"❌ Competitive positioning renderer error: {e}")
        return False

def test_precedent_transactions_renderer():
    """Test precedent transactions slide with proper acquirer data"""
    print("🔧 Testing precedent transactions renderer...")
    
    test_data = {
        "title": "Precedent Transactions",
        "transactions": [
            {
                "target": "Mellanox Technologies", 
                "acquirer": "NVIDIA Corporation",
                "date": "2020", 
                "country": "USA", 
                "enterprise_value": "$7.0B", 
                "revenue": "$1.3B", 
                "ev_revenue_multiple": "5.4x"
            },
            {
                "target": "Xilinx Inc", 
                "acquirer": "Advanced Micro Devices", 
                "date": "2022", 
                "country": "USA", 
                "enterprise_value": "$35.0B", 
                "revenue": "$3.2B", 
                "ev_revenue_multiple": "10.9x"
            }
        ]
    }
    
    try:
        prs = render_precedent_transactions_slide(
            data=test_data,
            company_name="NVIDIA"
        )
        
        if prs and len(prs.slides) > 0:
            print("✅ Precedent transactions slide rendered successfully")
            print(f"   - Transactions: {len(test_data['transactions'])} deals")
            for i, txn in enumerate(test_data['transactions']):
                acquirer = txn.get('acquirer')
                if acquirer and acquirer != 'N/A':
                    print(f"   - Deal {i+1}: {txn.get('target')} → {acquirer} ({txn.get('enterprise_value')})")
                else:
                    print(f"   - Deal {i+1}: Missing acquirer - this would cause NoneType error")
            return True
        else:
            print("❌ Precedent transactions slide rendering failed")  
            return False
            
    except Exception as e:
        print(f"❌ Precedent transactions renderer error: {e}")
        return False

def test_null_acquirer_handling():
    """Test that renderer handles null acquirer gracefully"""
    print("🔧 Testing null acquirer handling...")
    
    test_data = {
        "title": "Precedent Transactions", 
        "transactions": [
            {
                "target": "Test Company",
                "acquirer": None,  # This used to cause NoneType error
                "date": "2023",
                "country": "USA",
                "enterprise_value": "$1.0B",
                "revenue": "$100M",
                "ev_revenue_multiple": "10.0x"
            }
        ]
    }
    
    try:
        prs = render_precedent_transactions_slide(
            data=test_data,
            company_name="Test"
        )
        
        if prs and len(prs.slides) > 0:
            print("✅ Null acquirer handled gracefully - no NoneType error")
            return True
        else:
            print("❌ Null acquirer handling failed")
            return False
            
    except Exception as e:
        print(f"❌ Null acquirer handling error: {e}")
        return False

if __name__ == "__main__":
    print("🧪 Testing slide renderer fixes...\n")
    
    success1 = test_competitive_renderer()
    print()
    
    success2 = test_precedent_transactions_renderer() 
    print()
    
    success3 = test_null_acquirer_handling()
    print()
    
    if success1 and success2 and success3:
        print("🎉 All renderer tests passed! Slide fixes are working correctly.")
    else:
        print("💥 Some renderer tests failed. Check the errors above.")